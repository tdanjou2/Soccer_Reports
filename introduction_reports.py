# Imports
import requests
import json
import pandas as pd  
import numpy as np
import matplotlib as mpl
import matplotlib.pyplot as plt
from matplotlib import rcParams
from matplotlib.ticker import StrMethodFormatter
from pathlib import Path
import time
from matplotlib.lines import Line2D
import matplotlib.patches as mpatches
import plotly.graph_objects as go
from pptx.util import Inches
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
import shutil  
from shutil import copyfile
#from pptx.enum.shapes import MSO_SHAPE
#from pptx.enum.dml import MSO_FILL
import math

from pptx.util import Inches, Pt,Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import ColorFormat, RGBColor
import codecs

from parameters_package import *
from core_functions import * # For nan_to_zero function
from output_v2 import *
from summarize import *

import seaborn as sns
import warnings; warnings.filterwarnings(action='once')

'''
Package parameters
------------------------------------------------------------------------------------
'''

# The parameters for the font
mpl.rcParams["font.family"] = font_1

'''
Transfer functions
------------------------------------------------------------------------------------
'''

# OK - 
def introduction_master(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_introduction,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3):
    
    print('start of introduction_master')
    pres_path_0 = link_reports + central_folder + '/' + 'report' +'/'+'report' + '.pptx' 
    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    pres_path_old = link_reports + game_fixture_id + '/' + 'report' +'/'+'report_old' + '.pptx' 
    
    copyfile(pres_path, pres_path_old)
    copyfile(pres_path_0, pres_path)

    # Save the current version of the report
    prs = Presentation(pres_path)

    # Charts
    values_home, values_away = introduction_reports_chart_section1(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_introduction,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3)
    values_home_home, values_away_away = introduction_reports_chart_section3(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_introduction,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3)
    delta_home_mean, delta_away_mean, var_home_mean, var_away_mean, delta_max_home, BP_max_home,BC_max_home, delta_max_away, BP_max_away, BC_max_away, delta_min_home, BP_min_home, BC_min_home, delta_min_away, BP_min_away, BC_min_away = introduction_reports_chart_section2(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_introduction,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3)
    delta_home_mean_home, delta_max_home_home, BP_max_home_home, BC_max_home_home, delta_min_home_home, BP_min_home_home, BC_min_home_home = introduction_reports_chart_section4(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_introduction,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3)
    delta_away_mean_away, delta_max_away_away, BP_max_away_away, BC_max_away_away, delta_min_away_away, BP_min_away_away, BC_min_away_away = introduction_reports_chart_section5(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_introduction,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3)
    
    # Other photos and media
    transfer_image_to_ppt_introduction(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_introduction,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3)

    # Texts
    introduction_adv = transfer_text_to_ppt_introduction(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_introduction,values_home, values_away, values_home_home, values_away_away, delta_home_mean, delta_away_mean,var_home_mean, var_away_mean, delta_max_home, BP_max_home,BC_max_home, delta_max_away, BP_max_away, BC_max_away, delta_min_home, BP_min_home, BC_min_home, delta_min_away, BP_min_away, BC_min_away, delta_home_mean_home, delta_max_home_home, BP_max_home_home, BC_max_home_home, delta_min_home_home, BP_min_home_home, BC_min_home_home,delta_away_mean_away, delta_max_away_away, BP_max_away_away, BC_max_away_away, delta_min_away_away, BP_min_away_away, BC_min_away_away,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3)
    prs = Presentation(pres_path)

    # Advantage
    move_box_advantage(pres_path,"La forme et performances récentes donnent l'avantage à",str(introduction_adv),matrix_positions_introduction['introduction_advantage']['pos_left']+2*buffer_adv,matrix_positions_introduction['introduction_advantage']['pos_top']+2*buffer_adv)

    print('end of introduction_master')

# OK - 
def transfer_image_to_ppt_introduction(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_introduction,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3):
    
    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    
    # Charge ppt presentation
    prs = Presentation(pres_path)

    #x = pd.read_pickle(link_data+'{}'.format(team_id)+'/'+str(game_fixture_id) + '/' +'fixtures_package'+'/'+'fixtures_team_id_next_s'+ '_'+ str(team_id) +'_'+ str(nex) +'.pkl')
    #homeTeam_id = int(x.iloc[0,x.columns.get_loc('homeTeam')]['team_id'])
    
    # Logo home
    page_number, pos_left, pos_top = matrix_positions_introduction['logo_team_home']['page_number'], matrix_positions_introduction['logo_team_home']['pos_left'],matrix_positions_introduction['logo_team_home']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(homeTeam_id) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_2,prs))
    
    # Logo away
    page_number, pos_left, pos_top = matrix_positions_introduction['logo_team_away']['page_number'], matrix_positions_introduction['logo_team_away']['pos_left'],matrix_positions_introduction['logo_team_away']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(awayTeam_id) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_2,prs))
    
    # Graph 1
    page_number, pos_left, pos_top = matrix_positions_introduction['introduction_reports_chart_section1']['page_number'], matrix_positions_introduction['introduction_reports_chart_section1']['pos_left'],matrix_positions_introduction['introduction_reports_chart_section1']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports + game_fixture_id + '/' + name_folder_introduction + '/' +'introduction_reports_chart_section1' + graph_format , left=cm_to_ppx(pos_left-buffer,prs), top=cm_to_ppx(pos_top-buffer,prs), height= cm_to_ppx(height_1+2.2*buffer,prs))

    # Graph 2
    page_number, pos_left, pos_top = matrix_positions_introduction['introduction_reports_chart_section2']['page_number'], matrix_positions_introduction['introduction_reports_chart_section2']['pos_left'],matrix_positions_introduction['introduction_reports_chart_section2']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports + game_fixture_id + '/' + name_folder_introduction + '/' +'introduction_reports_chart_section2' + graph_format , left=cm_to_ppx(pos_left-buffer,prs), top=cm_to_ppx(pos_top-buffer,prs), height= cm_to_ppx(height_1+2.2*buffer,prs))
    
    # Graph 3
    page_number, pos_left, pos_top = matrix_positions_introduction['introduction_reports_chart_section3']['page_number'], matrix_positions_introduction['introduction_reports_chart_section3']['pos_left'],matrix_positions_introduction['introduction_reports_chart_section3']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports + game_fixture_id + '/' + name_folder_introduction + '/' +'introduction_reports_chart_section3' + graph_format , left=cm_to_ppx(pos_left-buffer,prs), top=cm_to_ppx(pos_top-buffer,prs), height= cm_to_ppx(height_1+2.2*buffer,prs))
    
    # Graph 4
    page_number, pos_left, pos_top = matrix_positions_introduction['introduction_reports_chart_section4']['page_number'], matrix_positions_introduction['introduction_reports_chart_section4']['pos_left'],matrix_positions_introduction['introduction_reports_chart_section4']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports + game_fixture_id + '/' + name_folder_introduction + '/' +'introduction_reports_chart_section4' + graph_format , left=cm_to_ppx(pos_left-buffer,prs), top=cm_to_ppx(pos_top-buffer,prs), height= cm_to_ppx(height_6+2.2*buffer,prs)) # the height to fit the page
   
   # Graph 5
    page_number, pos_left, pos_top = matrix_positions_introduction['introduction_reports_chart_section5']['page_number'], matrix_positions_introduction['introduction_reports_chart_section5']['pos_left'],matrix_positions_introduction['introduction_reports_chart_section5']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports + game_fixture_id + '/' + name_folder_introduction + '/' +'introduction_reports_chart_section5' + graph_format , left=cm_to_ppx(pos_left-buffer,prs), top=cm_to_ppx(pos_top-buffer,prs), height= cm_to_ppx(height_6+2.2*buffer,prs)) # the height to fit the page
   
    # Save ppt presentation 
    prs.save(pres_path)

# OK - 
def transfer_text_to_ppt_introduction(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_introduction,values_home, values_away, values_home_home, values_away_away, delta_home_mean, delta_away_mean,var_home_mean, var_away_mean, delta_max_home, BP_max_home,BC_max_home, delta_max_away, BP_max_away, BC_max_away, delta_min_home, BP_min_home, BC_min_home, delta_min_away, BP_min_away, BC_min_away, delta_home_mean_home, delta_max_home_home, BP_max_home_home, BC_max_home_home, delta_min_home_home, BP_min_home_home, BC_min_home_home,delta_away_mean_away, delta_max_away_away, BP_max_away_away, BC_max_away_away, delta_min_away_away, BP_min_away_away, BC_min_away_away,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3): 

    # Advantage parameters
    advantage_introduction_home = 0
    advantage_introduction_away = 0

    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 

    # Useful data part 1
    fixtures_home_away = pd.read_pickle(link_data+'{}'.format(team_id)+'/'+str(game_fixture_id) + '/' +'fixtures_package'+'/'+'fixtures_team_id_next_s'+ '_'+ str(team_id) +'_'+ str(nex) +'.pkl')
    rounds = fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('round')]
    league_data = pd.read_pickle(link_data+ central_folder +'/'+link_data_package_leagues_central+'/'+'leagues_league_id_master_s'+ '_'+ str(fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')]) +'.pkl')
    teams_league_data = pd.read_pickle(link_data+ central_folder +'/'+link_data_package_teams_central+'/'+'teams_league_id_s'+ '_'+ str(fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')]) +'.pkl')
    homeTeam_data = teams_league_data.loc[teams_league_data['team_id'] == int(homeTeam_id)]
    
    
    # Header


    # Titles
    
    
    # Beginning of populating of the section


    # Useful data part 2
    global_ranking = output_ranking(rounds,0,fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')],homeTeam_id)
    global_ranking_extract_home = output_ranking_extract(rounds,0,fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')],homeTeam_id,homeTeam_id)
    global_ranking_extract_home_previous = output_ranking_extract(rounds[:-1]+str(int(rounds[-1])-1),1,fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')],homeTeam_id,homeTeam_id)
    global_ranking_extract_away = output_ranking_extract(rounds,0,fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')],awayTeam_id,awayTeam_id)
    global_ranking_extract_away_previous = output_ranking_extract(rounds[:-1]+str(int(rounds[-1])-1),1,fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')],awayTeam_id,awayTeam_id)
    
    home_ranking = output_ranking_home(rounds,0,fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')],homeTeam_id)
    home_ranking_extract = output_ranking_extract_home(rounds,0,fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')],homeTeam_id,homeTeam_id)
    home_ranking_extract_previous = output_ranking_extract_home(rounds[:-1]+str(int(rounds[-1])-1),1,fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')],homeTeam_id,homeTeam_id)

    away_ranking = output_ranking_away(rounds,0,fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')],homeTeam_id)
    away_ranking_extract = output_ranking_extract_home(rounds,0,fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')],awayTeam_id,awayTeam_id)
    away_ranking_extract_previous = output_ranking_extract_home(rounds[:-1]+str(int(rounds[-1])-1),1,fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')],awayTeam_id,awayTeam_id)

    rank_global_home, rank_global_away, rank_home, rank_away = 0,0,0,0

    # Global Home
    global_ranking_ppt_home = {1:{'#':0,'var':0,'team_id':0,'team':0,'Pts':0,'W':0,'D':0,'L':0,'del':0},2:{'#':0,'var':0,'team_id':0,'team':0,'Pts':0,'W':0,'D':0,'L':0,'del':0},3:{'#':0,'var':0,'team_id':0,'team':0,'Pts':0,'W':0,'D':0,'L':0,'del':0},4:{'#':0,'var':0,'team_id':0,'team':0,'Pts':0,'W':0,'D':0,'L':0,'del':0},5:{'#':0,'var':0,'team_id':0,'team':0,'Pts':0,'W':0,'D':0,'L':0,'del':0},6:{'#':0,'var':0,'team_id':0,'team':0,'Pts':0,'W':0,'D':0,'L':0,'del':0},7:{'#':0,'var':0,'team_id':0,'team':0,'Pts':0,'W':0,'D':0,'L':0,'del':0}}
    
    nb_inf, nb_sup =  0, 0
    if len(global_ranking)-global_ranking_extract_home[1]>2:
        nb_inf = 3
    else:
        nb_inf = len(global_ranking)-global_ranking_extract_home[1]

    if global_ranking_extract_home[1]>3:
        nb_sup = 3
    else:
        nb_sup = global_ranking_extract_home[1]-1
    if nb_inf != 3:
        nb_sup += 3-nb_inf
    elif nb_sup != 3:
        nb_inf += 3-nb_sup
    
    cpt_3 = 7
    for cpt_1 in range (nb_inf,0,-1):
        x = output_ranking_extract(rounds,0,fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')],homeTeam_id,int(global_ranking.iloc[cpt_1+global_ranking_extract_home[1]-1].name))
        y = output_ranking_extract(rounds[:-1]+str(int(rounds[-1])-1),1,fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')],homeTeam_id,int(global_ranking.iloc[cpt_1+global_ranking_extract_home[1]-1].name))
        global_ranking_ppt_home[cpt_3]['team_id'] = int(global_ranking.iloc[cpt_1+global_ranking_extract_home[1]-1].name)
        global_ranking_ppt_home[cpt_3]['#'] = x[1]
        global_ranking_ppt_home[cpt_3]['Pts'] = x[2]
        global_ranking_ppt_home[cpt_3]['team'] = x[0]
        global_ranking_ppt_home[cpt_3]['W'] = x[4]
        global_ranking_ppt_home[cpt_3]['D'] = x[5]
        global_ranking_ppt_home[cpt_3]['L'] = x[6]
        global_ranking_ppt_home[cpt_3]['del'] = x[7]-x[8]
        if int(x[1]-y[1])>0:
            global_ranking_ppt_home[cpt_3]['var'] = str(-int(x[1]-y[1]))
        elif int(x[1]-y[1])<0:
            global_ranking_ppt_home[cpt_3]['var'] = "+" + str(-int(x[1]-y[1]))
        else:
            global_ranking_ppt_home[cpt_3]['var'] = "="
        cpt_3 -= 1
    
    rank_global_home = cpt_3
    global_ranking_ppt_home[cpt_3]['team_id'] = homeTeam_id
    global_ranking_ppt_home[cpt_3]['#'] = global_ranking_extract_home[1]
    global_ranking_ppt_home[cpt_3]['Pts'] = global_ranking_extract_home[2]
    global_ranking_ppt_home[cpt_3]['team'] = global_ranking_extract_home[0]
    global_ranking_ppt_home[cpt_3]['W'] = global_ranking_extract_home[4]
    global_ranking_ppt_home[cpt_3]['D'] = global_ranking_extract_home[5]
    global_ranking_ppt_home[cpt_3]['L'] = global_ranking_extract_home[6]
    global_ranking_ppt_home[cpt_3]['del'] = global_ranking_extract_home[7]-global_ranking_extract_home[8]
    if (global_ranking_extract_home[1]-global_ranking_extract_home_previous[1])>0:
        global_ranking_ppt_home[cpt_3]['var'] = str(-int(global_ranking_extract_home[1]-global_ranking_extract_home_previous[1]))
    elif (global_ranking_extract_home[1]-global_ranking_extract_home_previous[1])<0:
        global_ranking_ppt_home[cpt_3]['var'] = "+"+ str(-int(global_ranking_extract_home[1]-global_ranking_extract_home_previous[1]))
    else:
        global_ranking_ppt_home[cpt_3]['var'] = "="
    cpt_3 -= 1
    
    for cpt_1 in range (1,nb_sup+1):
        x = output_ranking_extract(rounds,0,fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')],homeTeam_id,int(global_ranking.iloc[-cpt_1+global_ranking_extract_home[1]-1].name))
        y = output_ranking_extract(rounds[:-1]+str(int(rounds[-1])-1),1,fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')],homeTeam_id,int(global_ranking.iloc[-cpt_1+global_ranking_extract_home[1]-1].name))
        global_ranking_ppt_home[cpt_3]['team_id'] = int(global_ranking.iloc[-cpt_1+global_ranking_extract_home[1]-1].name)
        global_ranking_ppt_home[cpt_3]['#'] = x[1]
        global_ranking_ppt_home[cpt_3]['Pts'] = x[2]
        global_ranking_ppt_home[cpt_3]['team'] = x[0]
        global_ranking_ppt_home[cpt_3]['W'] = x[4]
        global_ranking_ppt_home[cpt_3]['D'] = x[5]
        global_ranking_ppt_home[cpt_3]['L'] = x[6]
        global_ranking_ppt_home[cpt_3]['del'] = x[7]-x[8]
        if int(x[1]-y[1])>0:
            global_ranking_ppt_home[cpt_3]['var'] = str(-int(x[1]-y[1]))
        elif int(x[1]-y[1])<0:
            global_ranking_ppt_home[cpt_3]['var'] = "+" + str(-int(x[1]-y[1]))
        else:
            global_ranking_ppt_home[cpt_3]['var'] = "="
        cpt_3 -= 1

    # Global Away
    global_ranking_ppt_away = {1:{'#':0,'var':0,'team_id':0,'team':0,'Pts':0,'W':0,'D':0,'L':0,'del':0},2:{'#':0,'var':0,'team_id':0,'team':0,'Pts':0,'W':0,'D':0,'L':0,'del':0},3:{'#':0,'var':0,'team_id':0,'team':0,'Pts':0,'W':0,'D':0,'L':0,'del':0},4:{'#':0,'var':0,'team_id':0,'team':0,'Pts':0,'W':0,'D':0,'L':0,'del':0},5:{'#':0,'var':0,'team_id':0,'team':0,'Pts':0,'W':0,'D':0,'L':0,'del':0},6:{'#':0,'var':0,'team_id':0,'team':0,'Pts':0,'W':0,'D':0,'L':0,'del':0},7:{'#':0,'var':0,'team_id':0,'team':0,'Pts':0,'W':0,'D':0,'L':0,'del':0}}
    
    nb_inf, nb_sup =  0, 0
    if len(global_ranking)-global_ranking_extract_away[1]>2:
        nb_inf = 3
    else:
        nb_inf = len(global_ranking)-global_ranking_extract_away[1]

    if global_ranking_extract_away[1]>3:
        nb_sup = 3
    else:
        nb_sup = global_ranking_extract_away[1]-1
    if nb_inf != 3:
        nb_sup += 3-nb_inf
    elif nb_sup != 3:
        nb_inf += 3-nb_sup

    cpt_3 = 7
    for cpt_1 in range (nb_inf,0,-1):
        x = output_ranking_extract(rounds,0,fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')],awayTeam_id,int(global_ranking.iloc[cpt_1+global_ranking_extract_away[1]-1].name))
        y = output_ranking_extract(rounds[:-1]+str(int(rounds[-1])-1),1,fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')],awayTeam_id,int(global_ranking.iloc[cpt_1+global_ranking_extract_away[1]-1].name))
        global_ranking_ppt_away[cpt_3]['#'] = x[1]
        global_ranking_ppt_away[cpt_3]['team_id'] = int(global_ranking.iloc[cpt_1+global_ranking_extract_away[1]-1].name)
        global_ranking_ppt_away[cpt_3]['Pts'] = x[2]
        global_ranking_ppt_away[cpt_3]['team'] = x[0]
        global_ranking_ppt_away[cpt_3]['W'] = x[4]
        global_ranking_ppt_away[cpt_3]['D'] = x[5]
        global_ranking_ppt_away[cpt_3]['L'] = x[6]
        global_ranking_ppt_away[cpt_3]['del'] = x[7]-x[8]
        if int(x[1]-y[1])>0:
            global_ranking_ppt_away[cpt_3]['var'] = str(-int(x[1]-y[1]))
        elif int(x[1]-y[1])<0:
            global_ranking_ppt_away[cpt_3]['var'] = "+"+ str(-int(x[1]-y[1]))
        else:
            global_ranking_ppt_away[cpt_3]['var'] = "="
        cpt_3 -= 1
    
    rank_global_away = cpt_3
    global_ranking_ppt_away[cpt_3]['team_id'] = awayTeam_id
    global_ranking_ppt_away[cpt_3]['#'] = global_ranking_extract_away[1]
    global_ranking_ppt_away[cpt_3]['Pts'] = global_ranking_extract_away[2]
    global_ranking_ppt_away[cpt_3]['team'] = global_ranking_extract_away[0]
    global_ranking_ppt_away[cpt_3]['W'] = global_ranking_extract_away[4]
    global_ranking_ppt_away[cpt_3]['D'] = global_ranking_extract_away[5]
    global_ranking_ppt_away[cpt_3]['L'] = global_ranking_extract_away[6]
    global_ranking_ppt_away[cpt_3]['del'] = global_ranking_extract_away[7]-global_ranking_extract_away[8]
    if (global_ranking_extract_away[1]-global_ranking_extract_away_previous[1])>0:
        global_ranking_ppt_away[cpt_3]['var'] = str(-int(global_ranking_extract_away[1]-global_ranking_extract_away_previous[1]))
    elif (global_ranking_extract_away[1]-global_ranking_extract_away_previous[1])<0:
        global_ranking_ppt_away[cpt_3]['var'] = "+"+ str(-int(global_ranking_extract_away[1]-global_ranking_extract_away_previous[1]))
    else:
        global_ranking_ppt_away[cpt_3]['var'] = "="
    cpt_3 -= 1
    
    for cpt_1 in range (1,nb_sup+1):
        x = output_ranking_extract(rounds,0,fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')],awayTeam_id,int(global_ranking.iloc[-cpt_1+global_ranking_extract_away[1]-1].name))
        y = output_ranking_extract(rounds[:-1]+str(int(rounds[-1])-1),1,fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')],awayTeam_id,int(global_ranking.iloc[-cpt_1+global_ranking_extract_away[1]-1].name))
        global_ranking_ppt_away[cpt_3]['#'] = x[1]
        global_ranking_ppt_away[cpt_3]['team_id'] = int(global_ranking.iloc[-cpt_1+global_ranking_extract_away[1]-1].name)
        global_ranking_ppt_away[cpt_3]['Pts'] = x[2]
        global_ranking_ppt_away[cpt_3]['team'] = x[0]
        global_ranking_ppt_away[cpt_3]['W'] = x[4]
        global_ranking_ppt_away[cpt_3]['D'] = x[5]
        global_ranking_ppt_away[cpt_3]['L'] = x[6]
        global_ranking_ppt_away[cpt_3]['del'] = x[7]-x[8]
        if int(x[1]-y[1])>0:
            global_ranking_ppt_away[cpt_3]['var'] = str(-int(x[1]-y[1]))
        elif int(x[1]-y[1])<0:
            global_ranking_ppt_away[cpt_3]['var'] = "+"+ str(-int(x[1]-y[1]))
        else:
            global_ranking_ppt_away[cpt_3]['var'] = "="
        cpt_3 -= 1

    # Home
    home_ranking_ppt = {1:{'#':0,'var':0,'team_id':0,'team':0,'Pts':0,'W':0,'D':0,'L':0,'del':0},2:{'#':0,'var':0,'team_id':0,'team':0,'Pts':0,'W':0,'D':0,'L':0,'del':0},3:{'#':0,'var':0,'team_id':0,'team':0,'Pts':0,'W':0,'D':0,'L':0,'del':0},4:{'#':0,'var':0,'team_id':0,'team':0,'Pts':0,'W':0,'D':0,'L':0,'del':0},5:{'#':0,'var':0,'team_id':0,'team':0,'Pts':0,'W':0,'D':0,'L':0,'del':0},6:{'#':0,'var':0,'team_id':0,'team':0,'Pts':0,'W':0,'D':0,'L':0,'del':0},7:{'#':0,'var':0,'team_id':0,'team':0,'Pts':0,'W':0,'D':0,'L':0,'del':0}}
    
    nb_inf, nb_sup =  0, 0
    if len(home_ranking)-home_ranking_extract[1]>2:
        nb_inf = 3
    else:
        nb_inf = len(home_ranking)-home_ranking_extract[1]

    if home_ranking_extract[1]>3:
        nb_sup = 3
    else:
        nb_sup = home_ranking_extract[1]-1
    if nb_inf != 3:
        nb_sup += 3-nb_inf
    elif nb_sup != 3:
        nb_inf += 3-nb_sup

    cpt_3 = 7
    for cpt_1 in range (nb_inf,0,-1):
        x = output_ranking_extract_home(rounds,0,fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')],homeTeam_id,int(home_ranking.iloc[cpt_1+home_ranking_extract[1]-1].name))
        y = output_ranking_extract_home(rounds[:-1]+str(int(rounds[-1])-1),1,fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')],homeTeam_id,int(home_ranking.iloc[cpt_1+home_ranking_extract[1]-1].name))
        home_ranking_ppt[cpt_3]['#'] = x[1]
        home_ranking_ppt[cpt_3]['team_id'] = int(home_ranking.iloc[cpt_1+home_ranking_extract[1]-1].name)
        home_ranking_ppt[cpt_3]['team'] = x[0]
        home_ranking_ppt[cpt_3]['Pts'] = x[2]
        home_ranking_ppt[cpt_3]['W'] = x[4]
        home_ranking_ppt[cpt_3]['D'] = x[5]
        home_ranking_ppt[cpt_3]['L'] = x[6]
        home_ranking_ppt[cpt_3]['del'] = x[7]-x[8]
        if int(x[1]-y[1])>0:
            home_ranking_ppt[cpt_3]['var'] = str(-int(x[1]-y[1]))
        elif int(x[1]-y[1])<0:
            home_ranking_ppt[cpt_3]['var'] = "+"+ str(-int(x[1]-y[1]))
        else:
            home_ranking_ppt[cpt_3]['var'] = "="
        cpt_3 -= 1
    
    rank_home = cpt_3
    home_ranking_ppt[cpt_3]['#'] = home_ranking_extract[1]
    home_ranking_ppt[cpt_3]['Pts'] = home_ranking_extract[2]
    home_ranking_ppt[cpt_3]['team_id'] = homeTeam_id
    home_ranking_ppt[cpt_3]['team'] = home_ranking_extract[0]
    home_ranking_ppt[cpt_3]['W'] = home_ranking_extract[4]
    home_ranking_ppt[cpt_3]['D'] = home_ranking_extract[5]
    home_ranking_ppt[cpt_3]['L'] = home_ranking_extract[6]
    home_ranking_ppt[cpt_3]['del'] = home_ranking_extract[7]-home_ranking_extract[8]
    if (home_ranking_extract[1]-home_ranking_extract_previous[1])>0:
        home_ranking_ppt[cpt_3]['var'] = str(-int(home_ranking_extract[1]-home_ranking_extract_previous[1]))
    elif (home_ranking_extract[1]-home_ranking_extract_previous[1])<0:
        home_ranking_ppt[cpt_3]['var'] = "+"+ str(-int(home_ranking_extract[1]-home_ranking_extract_previous[1]))
    else:
        home_ranking_ppt[cpt_3]['var'] = "="
    cpt_3 -= 1
    
    for cpt_1 in range (1,nb_sup+1):
        x = output_ranking_extract_home(rounds,0,fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')],homeTeam_id,int(home_ranking.iloc[-cpt_1+home_ranking_extract[1]-1].name))
        y = output_ranking_extract_home(rounds[:-1]+str(int(rounds[-1])-1),1,fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')],homeTeam_id,int(home_ranking.iloc[-cpt_1+home_ranking_extract[1]-1].name))
        home_ranking_ppt[cpt_3]['#'] = x[1]
        home_ranking_ppt[cpt_3]['team_id'] = int(home_ranking.iloc[-cpt_1+home_ranking_extract[1]-1].name)
        home_ranking_ppt[cpt_3]['Pts'] = x[2]
        home_ranking_ppt[cpt_3]['team'] = x[0]
        home_ranking_ppt[cpt_3]['W'] = x[4]
        home_ranking_ppt[cpt_3]['D'] = x[5]
        home_ranking_ppt[cpt_3]['L'] = x[6]
        home_ranking_ppt[cpt_3]['del'] = x[7]-x[8]
        if int(x[1]-y[1])>0:
            home_ranking_ppt[cpt_3]['var'] = str(-int(x[1]-y[1]))
        elif int(x[1]-y[1])<0:
            home_ranking_ppt[cpt_3]['var'] = "+"+ str(-int(x[1]-y[1]))
        else:
            home_ranking_ppt[cpt_3]['var'] = "="
        cpt_3 -= 1

    # Away
    away_ranking_ppt = {1:{'#':0,'var':0,'team_id':0,'team':0,'Pts':0,'W':0,'D':0,'L':0,'del':0},2:{'#':0,'var':0,'team_id':0,'team':0,'Pts':0,'W':0,'D':0,'L':0,'del':0},3:{'#':0,'var':0,'team_id':0,'team':0,'Pts':0,'W':0,'D':0,'L':0,'del':0},4:{'#':0,'var':0,'team_id':0,'team':0,'Pts':0,'W':0,'D':0,'L':0,'del':0},5:{'#':0,'var':0,'team_id':0,'team':0,'Pts':0,'W':0,'D':0,'L':0,'del':0},6:{'#':0,'var':0,'team_id':0,'team':0,'Pts':0,'W':0,'D':0,'L':0,'del':0},7:{'#':0,'var':0,'team_id':0,'team':0,'Pts':0,'W':0,'D':0,'L':0,'del':0}}
    
    nb_inf, nb_sup =  0, 0
    if len(away_ranking)-away_ranking_extract[1]>2:
        nb_inf = 3
    else:
        nb_inf = len(away_ranking)-away_ranking_extract[1]

    if away_ranking_extract[1]>3:
        nb_sup = 3
    else:
        nb_sup = away_ranking_extract[1]-1
    if nb_inf != 3:
        nb_sup += 3-nb_inf
    elif nb_sup != 3:
        nb_inf += 3-nb_sup

    cpt_3 = 7
    for cpt_1 in range (nb_inf,0,-1):
        x = output_ranking_extract_away(rounds,0,fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')],awayTeam_id,int(away_ranking.iloc[cpt_1+away_ranking_extract[1]-1].name))
        y = output_ranking_extract_away(rounds[:-1]+str(int(rounds[-1])-1),1,fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')],awayTeam_id,int(away_ranking.iloc[cpt_1+away_ranking_extract[1]-1].name))
        away_ranking_ppt[cpt_3]['#'] = x[1]
        away_ranking_ppt[cpt_3]['team_id'] = int(away_ranking.iloc[cpt_1+away_ranking_extract[1]-1].name)
        away_ranking_ppt[cpt_3]['Pts'] = x[2]
        away_ranking_ppt[cpt_3]['team'] = x[0]
        away_ranking_ppt[cpt_3]['W'] = x[4]
        away_ranking_ppt[cpt_3]['D'] = x[5]
        away_ranking_ppt[cpt_3]['L'] = x[6]
        away_ranking_ppt[cpt_3]['del'] = x[7]-x[8]
        if int(x[1]-y[1])>0:
            away_ranking_ppt[cpt_3]['var'] = str(-int(x[1]-y[1]))
        elif int(x[1]-y[1])<0:
            away_ranking_ppt[cpt_3]['var'] = "+"+  str(-int(x[1]-y[1]))
        else:
            away_ranking_ppt[cpt_3]['var'] = "="
        cpt_3 -= 1
    
    rank_away = cpt_3
    away_ranking_ppt[cpt_3]['#'] = away_ranking_extract[1]
    away_ranking_ppt[cpt_3]['Pts'] = away_ranking_extract[2]
    away_ranking_ppt[cpt_3]['team_id'] = awayTeam_id
    away_ranking_ppt[cpt_3]['team'] = away_ranking_extract[0]
    away_ranking_ppt[cpt_3]['W'] = away_ranking_extract[4]
    away_ranking_ppt[cpt_3]['D'] = away_ranking_extract[5]
    away_ranking_ppt[cpt_3]['L'] = away_ranking_extract[6]
    away_ranking_ppt[cpt_3]['del'] = away_ranking_extract[7]-away_ranking_extract[8]
    if (away_ranking_extract[1]-away_ranking_extract_previous[1])>0:
        away_ranking_ppt[cpt_3]['var'] = str(-int(away_ranking_extract[1]-away_ranking_extract_previous[1]))
    elif (away_ranking_extract[1]-away_ranking_extract_previous[1])<0:
        away_ranking_ppt[cpt_3]['var'] = "+"+ str(-int(away_ranking_extract[1]-away_ranking_extract_previous[1]))
    else:
        away_ranking_ppt[cpt_3]['var'] = "="
    cpt_3 -= 1
    
    for cpt_1 in range (1,nb_sup+1):
        x = output_ranking_extract_away(rounds,0,fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')],awayTeam_id,int(away_ranking.iloc[-cpt_1+away_ranking_extract[1]-1].name))
        y = output_ranking_extract_away(rounds[:-1]+str(int(rounds[-1])-1),1,fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')],awayTeam_id,int(away_ranking.iloc[-cpt_1+away_ranking_extract[1]-1].name))
        away_ranking_ppt[cpt_3]['#'] = x[1]
        away_ranking_ppt[cpt_3]['team_id'] = int(away_ranking.iloc[-cpt_1+away_ranking_extract[1]-1].name)
        away_ranking_ppt[cpt_3]['Pts'] = x[2]
        away_ranking_ppt[cpt_3]['team'] = x[0]
        away_ranking_ppt[cpt_3]['W'] = x[4]
        away_ranking_ppt[cpt_3]['D'] = x[5]
        away_ranking_ppt[cpt_3]['L'] = x[6]
        away_ranking_ppt[cpt_3]['del'] = x[7]-x[8]
        if int(x[1]-y[1])>0:
            away_ranking_ppt[cpt_3]['var'] = str(-int(x[1]-y[1]))
        elif int(x[1]-y[1])<0:
            away_ranking_ppt[cpt_3]['var'] = "+"+ str(-int(x[1]-y[1]))
        else:
            away_ranking_ppt[cpt_3]['var'] = "="
        cpt_3 -= 1

    home_away_para = "H"
    for cpt_4 in range (1,8):
        
        # Global ranking of home team
        if global_ranking_ppt_home[cpt_4]['team_id'] != homeTeam_id :
            text_modification_paragraph_center(pres_path,"{}{}R".format(cpt_4,home_away_para),global_ranking_ppt_home[cpt_4]['#'],font_size_ranking,font_1,white_color_report,True,None)
            text_modification_paragraph_center(pres_path,"{}{}T".format(cpt_4,home_away_para),global_ranking_ppt_home[cpt_4]['team'],font_size_ranking,font_1,white_color_report,True,None)
            text_modification_paragraph_center(pres_path,"{}{}Pts".format(cpt_4,home_away_para),global_ranking_ppt_home[cpt_4]['Pts'],font_size_ranking,font_1,white_color_report,True,None)
            text_modification_paragraph_center(pres_path,"{}{}W".format(cpt_4,home_away_para),global_ranking_ppt_home[cpt_4]['W'],font_size_ranking,font_1,white_color_report,True,None)
            text_modification_paragraph_center(pres_path,"{}{}D".format(cpt_4,home_away_para),global_ranking_ppt_home[cpt_4]['D'],font_size_ranking,font_1,white_color_report,True,None)
            text_modification_paragraph_center(pres_path,"{}{}L".format(cpt_4,home_away_para),global_ranking_ppt_home[cpt_4]['L'],font_size_ranking,font_1,white_color_report,True,None)
            text_modification_paragraph_center(pres_path,"{}{}d".format(cpt_4,home_away_para),global_ranking_ppt_home[cpt_4]['del'],font_size_ranking,font_1,white_color_report,True,None)
        
        else:
            text_modification_paragraph_center_line(pres_path,"{}{}R".format(cpt_4,home_away_para),global_ranking_ppt_home[cpt_4]['#'],font_size_ranking,font_1,white_color_report,True,None,int(25000)) # The width of the border
            text_modification_paragraph_center_line(pres_path,"{}{}T".format(cpt_4,home_away_para),global_ranking_ppt_home[cpt_4]['team'],font_size_ranking,font_1,white_color_report,True,None,int(25000)) # The width of the border
            text_modification_paragraph_center_line(pres_path,"{}{}Pts".format(cpt_4,home_away_para),global_ranking_ppt_home[cpt_4]['Pts'],font_size_ranking,font_1,white_color_report,True,None,int(25000)) # The width of the border
            text_modification_paragraph_center_line(pres_path,"{}{}W".format(cpt_4,home_away_para),global_ranking_ppt_home[cpt_4]['W'],font_size_ranking,font_1,white_color_report,True,None,int(25000)) # The width of the border
            text_modification_paragraph_center_line(pres_path,"{}{}D".format(cpt_4,home_away_para),global_ranking_ppt_home[cpt_4]['D'],font_size_ranking,font_1,white_color_report,True,None,int(25000)) # The width of the border
            text_modification_paragraph_center_line(pres_path,"{}{}L".format(cpt_4,home_away_para),global_ranking_ppt_home[cpt_4]['L'],font_size_ranking,font_1,white_color_report,True,None,int(25000)) # The width of the border
            text_modification_paragraph_center_line(pres_path,"{}{}d".format(cpt_4,home_away_para),global_ranking_ppt_home[cpt_4]['del'],font_size_ranking,font_1,white_color_report,True,None,int(25000)) # The width of the border

        if global_ranking_ppt_home[cpt_4]['var'][-1] == "=":
            text_modification_paragraph_center(pres_path,"{}{}+".format(cpt_4,home_away_para),global_ranking_ppt_home[cpt_4]['var'],font_size_ranking,font_1,black_color_report,True,None)
        else:
            if global_ranking_ppt_home[cpt_4]['var'][0] == "+":
                text_modification_paragraph_center(pres_path,"{}{}+".format(cpt_4,home_away_para),global_ranking_ppt_home[cpt_4]['var'],font_size_ranking,font_1,green_color_report,True,None)
            else:
                text_modification_paragraph_center(pres_path,"{}{}+".format(cpt_4,home_away_para),global_ranking_ppt_home[cpt_4]['var'],font_size_ranking,font_1,red_color_report,True,None)
        
        # Home ranking of home team
        if home_ranking_ppt[cpt_4]['team_id'] != homeTeam_id :
            text_modification_paragraph_center(pres_path,"{}{}{}R".format(cpt_4,home_away_para,home_away_para),home_ranking_ppt[cpt_4]['#'],font_size_ranking,font_1,white_color_report,True,None)        
            text_modification_paragraph_center(pres_path,"{}{}{}T".format(cpt_4,home_away_para,home_away_para),home_ranking_ppt[cpt_4]['team'],font_size_ranking,font_1,white_color_report,True,None)
            text_modification_paragraph_center(pres_path,"{}{}{}Pts".format(cpt_4,home_away_para,home_away_para),home_ranking_ppt[cpt_4]['Pts'],font_size_ranking,font_1,white_color_report,True,None)
            text_modification_paragraph_center(pres_path,"{}{}{}W".format(cpt_4,home_away_para,home_away_para),home_ranking_ppt[cpt_4]['W'],font_size_ranking,font_1,white_color_report,True,None)
            text_modification_paragraph_center(pres_path,"{}{}{}D".format(cpt_4,home_away_para,home_away_para),home_ranking_ppt[cpt_4]['D'],font_size_ranking,font_1,white_color_report,True,None)
            text_modification_paragraph_center(pres_path,"{}{}{}L".format(cpt_4,home_away_para,home_away_para),home_ranking_ppt[cpt_4]['L'],font_size_ranking,font_1,white_color_report,True,None)
            text_modification_paragraph_center(pres_path,"{}{}{}d".format(cpt_4,home_away_para,home_away_para),home_ranking_ppt[cpt_4]['del'],font_size_ranking,font_1,white_color_report,True,None)
        
        else:
            text_modification_paragraph_center_line(pres_path,"{}{}{}R".format(cpt_4,home_away_para,home_away_para),home_ranking_ppt[cpt_4]['#'],font_size_ranking,font_1,white_color_report,True,None,int(25000)) # The width of the border    
            text_modification_paragraph_center_line(pres_path,"{}{}{}T".format(cpt_4,home_away_para,home_away_para),home_ranking_ppt[cpt_4]['team'],font_size_ranking,font_1,white_color_report,True,None,int(25000)) # The width of the border
            text_modification_paragraph_center_line(pres_path,"{}{}{}Pts".format(cpt_4,home_away_para,home_away_para),home_ranking_ppt[cpt_4]['Pts'],font_size_ranking,font_1,white_color_report,True,None,int(25000)) # The width of the border
            text_modification_paragraph_center_line(pres_path,"{}{}{}W".format(cpt_4,home_away_para,home_away_para),home_ranking_ppt[cpt_4]['W'],font_size_ranking,font_1,white_color_report,True,None,int(25000)) # The width of the border
            text_modification_paragraph_center_line(pres_path,"{}{}{}D".format(cpt_4,home_away_para,home_away_para),home_ranking_ppt[cpt_4]['D'],font_size_ranking,font_1,white_color_report,True,None,int(25000)) # The width of the border
            text_modification_paragraph_center_line(pres_path,"{}{}{}L".format(cpt_4,home_away_para,home_away_para),home_ranking_ppt[cpt_4]['L'],font_size_ranking,font_1,white_color_report,True,None,int(25000)) # The width of the border
            text_modification_paragraph_center_line(pres_path,"{}{}{}d".format(cpt_4,home_away_para,home_away_para),home_ranking_ppt[cpt_4]['del'],font_size_ranking,font_1,white_color_report,True,None,int(25000)) # The width of the border

        if home_ranking_ppt[cpt_4]['var'][-1] == "=":
            text_modification_paragraph_center(pres_path,"{}{}{}+".format(cpt_4,home_away_para,home_away_para),home_ranking_ppt[cpt_4]['var'],font_size_ranking,font_1,black_color_report,True,None)
        else:
            if home_ranking_ppt[cpt_4]['var'][0] == "+":
                text_modification_paragraph_center(pres_path,"{}{}{}+".format(cpt_4,home_away_para,home_away_para),home_ranking_ppt[cpt_4]['var'],font_size_ranking,font_1,green_color_report,True,None)
            else:
                text_modification_paragraph_center(pres_path,"{}{}{}+".format(cpt_4,home_away_para,home_away_para),home_ranking_ppt[cpt_4]['var'],font_size_ranking,font_1,red_color_report,True,None)
        print(cpt_4,home_away_para)

    home_away_para = "A"
    for cpt_4 in range (1,8):
        # Global ranking of away team
        if global_ranking_ppt_away[cpt_4]['team_id'] != awayTeam_id :
            text_modification_paragraph_center(pres_path,"{}{}R".format(cpt_4,home_away_para),global_ranking_ppt_away[cpt_4]['#'],font_size_ranking,font_1,white_color_report,True,None)
            text_modification_paragraph_center(pres_path,"{}{}T".format(cpt_4,home_away_para),global_ranking_ppt_away[cpt_4]['team'],font_size_ranking,font_1,white_color_report,True,None)
            text_modification_paragraph_center(pres_path,"{}{}Pts".format(cpt_4,home_away_para),global_ranking_ppt_away[cpt_4]['Pts'],font_size_ranking,font_1,white_color_report,True,None)
            text_modification_paragraph_center(pres_path,"{}{}W".format(cpt_4,home_away_para),global_ranking_ppt_away[cpt_4]['W'],font_size_ranking,font_1,white_color_report,True,None)
            text_modification_paragraph_center(pres_path,"{}{}D".format(cpt_4,home_away_para),global_ranking_ppt_away[cpt_4]['D'],font_size_ranking,font_1,white_color_report,True,None)
            text_modification_paragraph_center(pres_path,"{}{}L".format(cpt_4,home_away_para),global_ranking_ppt_away[cpt_4]['L'],font_size_ranking,font_1,white_color_report,True,None)
            text_modification_paragraph_center(pres_path,"{}{}d".format(cpt_4,home_away_para),global_ranking_ppt_away[cpt_4]['del'],font_size_ranking,font_1,white_color_report,True,None)
        
        else:
            text_modification_paragraph_center_line(pres_path,"{}{}R".format(cpt_4,home_away_para),global_ranking_ppt_away[cpt_4]['#'],font_size_ranking,font_1,white_color_report,True,None,int(25000)) # The width of the border
            text_modification_paragraph_center_line(pres_path,"{}{}T".format(cpt_4,home_away_para),global_ranking_ppt_away[cpt_4]['team'],font_size_ranking,font_1,white_color_report,True,None,int(25000)) # The width of the border
            text_modification_paragraph_center_line(pres_path,"{}{}Pts".format(cpt_4,home_away_para),global_ranking_ppt_away[cpt_4]['Pts'],font_size_ranking,font_1,white_color_report,True,None,int(25000)) # The width of the border
            text_modification_paragraph_center_line(pres_path,"{}{}W".format(cpt_4,home_away_para),global_ranking_ppt_away[cpt_4]['W'],font_size_ranking,font_1,white_color_report,True,None,int(25000)) # The width of the border
            text_modification_paragraph_center_line(pres_path,"{}{}D".format(cpt_4,home_away_para),global_ranking_ppt_away[cpt_4]['D'],font_size_ranking,font_1,white_color_report,True,None,int(25000)) # The width of the border
            text_modification_paragraph_center_line(pres_path,"{}{}L".format(cpt_4,home_away_para),global_ranking_ppt_away[cpt_4]['L'],font_size_ranking,font_1,white_color_report,True,None,int(25000)) # The width of the border
            text_modification_paragraph_center_line(pres_path,"{}{}d".format(cpt_4,home_away_para),global_ranking_ppt_away[cpt_4]['del'],font_size_ranking,font_1,white_color_report,True,None,int(25000)) # The width of the border

        if global_ranking_ppt_away[cpt_4]['var'][-1] == "=":
            text_modification_paragraph_center(pres_path,"{}{}+".format(cpt_4,home_away_para),global_ranking_ppt_away[cpt_4]['var'],font_size_ranking,font_1,black_color_report,True,None)
        else:
            if global_ranking_ppt_away[cpt_4]['var'][0] == "+":
                text_modification_paragraph_center(pres_path,"{}{}+".format(cpt_4,home_away_para),global_ranking_ppt_away[cpt_4]['var'],font_size_ranking,font_1,green_color_report,True,None)
            else:
                text_modification_paragraph_center(pres_path,"{}{}+".format(cpt_4,home_away_para),global_ranking_ppt_away[cpt_4]['var'],font_size_ranking,font_1,red_color_report,True,None)

        # Away ranking of away team
        if away_ranking_ppt[cpt_4]['team_id'] != awayTeam_id :
            text_modification_paragraph_center(pres_path,"{}{}{}R".format(cpt_4,home_away_para,home_away_para),away_ranking_ppt[cpt_4]['#'],font_size_ranking,font_1,white_color_report,True,None)
            text_modification_paragraph_center(pres_path,"{}{}{}T".format(cpt_4,home_away_para,home_away_para),away_ranking_ppt[cpt_4]['team'],font_size_ranking,font_1,white_color_report,True,None)
            text_modification_paragraph_center(pres_path,"{}{}{}Pts".format(cpt_4,home_away_para,home_away_para),away_ranking_ppt[cpt_4]['Pts'],font_size_ranking,font_1,white_color_report,True,None)
            text_modification_paragraph_center(pres_path,"{}{}{}W".format(cpt_4,home_away_para,home_away_para),away_ranking_ppt[cpt_4]['W'],font_size_ranking,font_1,white_color_report,True,None)
            text_modification_paragraph_center(pres_path,"{}{}{}D".format(cpt_4,home_away_para,home_away_para),away_ranking_ppt[cpt_4]['D'],font_size_ranking,font_1,white_color_report,True,None)
            text_modification_paragraph_center(pres_path,"{}{}{}L".format(cpt_4,home_away_para,home_away_para),away_ranking_ppt[cpt_4]['L'],font_size_ranking,font_1,white_color_report,True,None)
            text_modification_paragraph_center(pres_path,"{}{}{}d".format(cpt_4,home_away_para,home_away_para),away_ranking_ppt[cpt_4]['del'],font_size_ranking,font_1,white_color_report,True,None)
        
        else:    
            text_modification_paragraph_center_line(pres_path,"{}{}{}R".format(cpt_4,home_away_para,home_away_para),away_ranking_ppt[cpt_4]['#'],font_size_ranking,font_1,white_color_report,True,None,int(25000)) # The width of the border
            text_modification_paragraph_center_line(pres_path,"{}{}{}T".format(cpt_4,home_away_para,home_away_para),away_ranking_ppt[cpt_4]['team'],font_size_ranking,font_1,white_color_report,True,None,int(25000)) # The width of the border
            text_modification_paragraph_center_line(pres_path,"{}{}{}Pts".format(cpt_4,home_away_para,home_away_para),away_ranking_ppt[cpt_4]['Pts'],font_size_ranking,font_1,white_color_report,True,None,int(25000)) # The width of the border
            text_modification_paragraph_center_line(pres_path,"{}{}{}W".format(cpt_4,home_away_para,home_away_para),away_ranking_ppt[cpt_4]['W'],font_size_ranking,font_1,white_color_report,True,None,int(25000)) # The width of the border
            text_modification_paragraph_center_line(pres_path,"{}{}{}D".format(cpt_4,home_away_para,home_away_para),away_ranking_ppt[cpt_4]['D'],font_size_ranking,font_1,white_color_report,True,None,int(25000)) # The width of the border
            text_modification_paragraph_center_line(pres_path,"{}{}{}L".format(cpt_4,home_away_para,home_away_para),away_ranking_ppt[cpt_4]['L'],font_size_ranking,font_1,white_color_report,True,None,int(25000)) # The width of the border
            text_modification_paragraph_center_line(pres_path,"{}{}{}d".format(cpt_4,home_away_para,home_away_para),away_ranking_ppt[cpt_4]['del'],font_size_ranking,font_1,white_color_report,True,None,int(25000)) # The width of the border
        
        if away_ranking_ppt[cpt_4]['var'][-1] == "=":
            text_modification_paragraph_center(pres_path,"{}{}{}+".format(cpt_4,home_away_para,home_away_para),away_ranking_ppt[cpt_4]['var'],font_size_ranking,font_1,black_color_report,True,None)
        else:
            if away_ranking_ppt[cpt_4]['var'][0] == "+":
                text_modification_paragraph_center(pres_path,"{}{}{}+".format(cpt_4,home_away_para,home_away_para),away_ranking_ppt[cpt_4]['var'],font_size_ranking,font_1,green_color_report,True,None)
            else:
                text_modification_paragraph_center(pres_path,"{}{}{}+".format(cpt_4,home_away_para,home_away_para),away_ranking_ppt[cpt_4]['var'],font_size_ranking,font_1,red_color_report,True,None)
        print(cpt_4,home_away_para) 

    # Charge ppt presentation
    prs = Presentation(pres_path)
    
    text_modification_paragraph_center(pres_path,"HOME TEAM NAME 1",homeTeam_name,14,font_1,white_color_report,True,None)
    text_modification_paragraph_center(pres_path,"AWAY TEAM NAME 1",awayTeam_name,14,font_1,white_color_report,True,None)
    text_modification_paragraph_center(pres_path,"LEAGUE SEASON 1","{} - {}".format(fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league')]['name'],league_data.iloc[0,league_data.columns.get_loc('season')]),16,font_1,white_color_report,True,None)
    text_modification_paragraph_center(pres_path,"TIME DATE 1",fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('event_date')],16,font_1,white_color_report,True,None)
    text_modification_paragraph_center(pres_path,"L2","{}".format(last_fixtures_lineups_duration),16,font_1,white_color_report,True,None)

    # The biggest defeats and wins in the last 20 games 
    a, b, c, d = 0, 0, 0, 0
    if df_raw_home_1.iloc[delta_max_home,df_raw_home_1.columns.get_loc('away_home')] == "home":
        text_modification_paragraph_center(pres_path,"home team win home",homeTeam_name,7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"home team win away",df_raw_home_1.iloc[delta_max_home,df_raw_home_1.columns.get_loc('opp_team_name')],7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"z1",str(int(BP_max_home)),14,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"z2",str(int(BC_max_home)),14,font_1,white_color_report,True,None) 
        a = BP_max_home

    else:
        text_modification_paragraph_center(pres_path,"home team win home",df_raw_home_1.iloc[delta_max_home,df_raw_home_1.columns.get_loc('opp_team_name')],7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"home team win away",homeTeam_name,7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"z1",str(int(BC_max_home)),14,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"z2",str(int(BP_max_home)),14,font_1,white_color_report,True,None) 
        a = BC_max_home

    if df_raw_home_1.iloc[delta_min_home,df_raw_home_1.columns.get_loc('away_home')] == "home":
        text_modification_paragraph_center(pres_path,"home team lose home",homeTeam_name,7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"home team lose away",df_raw_home_1.iloc[delta_min_home,df_raw_home_1.columns.get_loc('opp_team_name')],7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"z5",str(int(BP_min_home)),14,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"z6",str(int(BC_min_home)),14,font_1,white_color_report,True,None) 
        b = BP_min_home

    else:
        text_modification_paragraph_center(pres_path,"home team lose home",df_raw_home_1.iloc[delta_min_home,df_raw_home_1.columns.get_loc('opp_team_name')],7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"home team lose away",homeTeam_name,7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"z5",str(int(BC_min_home)),14,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"z6",str(int(BP_min_home)),14,font_1,white_color_report,True,None) 
        b = BC_min_home

    if df_raw_away_1.iloc[delta_max_away,df_raw_away_1.columns.get_loc('away_home')] == "home":
        text_modification_paragraph_center(pres_path,"away team win home",awayTeam_name,7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"away team win away",df_raw_away_1.iloc[delta_max_away,df_raw_away_1.columns.get_loc('opp_team_name')],7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"z3",str(int(BP_max_away)),14,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"z4",str(int(BC_max_away)),14,font_1,white_color_report,True,None) 
        c = BP_max_away

    else:
        text_modification_paragraph_center(pres_path,"away team win home",df_raw_away_1.iloc[delta_max_away,df_raw_away_1.columns.get_loc('opp_team_name')],7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"away team win away",awayTeam_name,7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"z3",str(int(BC_max_away)),14,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"z4",str(int(BP_max_away)),14,font_1,white_color_report,True,None) 
        c = BC_max_away

    if df_raw_away_1.iloc[delta_min_away,df_raw_away_1.columns.get_loc('away_home')] == "home":
        text_modification_paragraph_center(pres_path,"away team lose home",awayTeam_name,7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"away team lose away",df_raw_away_1.iloc[delta_min_away,df_raw_away_1.columns.get_loc('opp_team_name')],7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"z7",str(int(BP_min_away)),14,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"z8",str(int(BC_min_away)),14,font_1,white_color_report,True,None) 
        d = BP_min_away

    else:
        text_modification_paragraph_center(pres_path,"away team lose home",df_raw_away_1.iloc[delta_min_away,df_raw_away_1.columns.get_loc('opp_team_name')],7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"away team lose away",awayTeam_name,7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"z7",str(int(BC_min_away)),14,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"z8",str(int(BP_min_away)),14,font_1,white_color_report,True,None)  
        d = BC_min_away
    
    # The biggest defeats and wins in the last years
    
    if df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest win')]["homeaway"] == "home":
        text_modification_paragraph_center(pres_path,"home team win homey",homeTeam_name,7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"home team win awayy",df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest win')]["team_name"],7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"y1",df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest win')]["BP"],14,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"y2",df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest win')]["BC"],14,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"Date win home",df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest win')]["Date"][:10],12,font_1,draw_color_report,True,True)    
    else:
        text_modification_paragraph_center(pres_path,"home team win homey",df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest win')]["team_name"],7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"home team win awayy",homeTeam_name,7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"y1",df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest win')]["BC"],14,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"y2",df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest win')]["BP"],14,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"Date win home",df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest win')]["Date"][:10],12,font_1,draw_color_report,True,True)   

    if df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest defeat')]["homeaway"] == "home":
        text_modification_paragraph_center(pres_path,"home team lose homey",homeTeam_name,7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"home team lose awayy",df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest defeat')]["team_name"],7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"y5",df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest defeat')]["BP"],14,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"y6",df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest defeat')]["BC"],14,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"Date lose home",df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest defeat')]["Date"][:10],12,font_1,draw_color_report,True,True)    
    else:
        text_modification_paragraph_center(pres_path,"home team lose homey",df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest defeat')]["team_name"],7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"home team lose awayy",homeTeam_name,7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"y5",df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest defeat')]["BC"],14,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"y6",df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest defeat')]["BP"],14,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"Date lose home",df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest defeat')]["Date"][:10],12,font_1,draw_color_report,True,True)   

    if df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest win')]["homeaway"] == "home":
        text_modification_paragraph_center(pres_path,"away team win homey",awayTeam_name,7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"away team win awayy",df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest win')]["team_name"],7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"y3",df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest win')]["BP"],14,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"y4",df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest win')]["BC"],14,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"Date win away",df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest win')]["Date"][:10],12,font_1,draw_color_report,True,True)    
    else:
        text_modification_paragraph_center(pres_path,"away team win homey",df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest win')]["team_name"],7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"away team win awayy",awayTeam_name,7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"y3",df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest win')]["BC"],14,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"y4",df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest win')]["BP"],14,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"Date win away",df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest win')]["Date"][:10],12,font_1,draw_color_report,True,True) 

    if df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest defeat')]["homeaway"] == "home":
        text_modification_paragraph_center(pres_path,"away team lose homey",awayTeam_name,7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"away team lose awayy",df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest defeat')]["team_name"],7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"y7",df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest defeat')]["BP"],14,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"y8",df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest defeat')]["BC"],14,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"Date lose away",df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest defeat')]["Date"][:10],12,font_1,draw_color_report,True,True)    
    else:
        text_modification_paragraph_center(pres_path,"away team lose homey",df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest defeat')]["team_name"],7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"away team lose awayy",awayTeam_name,7,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"y7",df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest defeat')]["BC"],14,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"y8",df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest defeat')]["BP"],14,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"Date lose away",df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest defeat')]["Date"][:10],12,font_1,draw_color_report,True,True) 

    # 1 criteria for the advantage
    if a>c :
        advantage_introduction_home += 1
    else :
        advantage_introduction_away += 1
    # 2 criteria for the advantage
    if b>d :
        advantage_introduction_away += 1
    else :
        advantage_introduction_home += 1

    ######################################## FRENCH BELOW ########################################
    
    if english == 0:
        text_modification_paragraph_center(pres_path,"ROUNDS 1","{}".format(fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('round')]),16,font_1,white_color_report,True,None)
        text_modification_paragraph_center(pres_path,"STADIUM CAPACITY 1","{} | capacité: {} places".format(homeTeam_data.iloc[0,homeTeam_data.columns.get_loc('venue_name')],homeTeam_data.iloc[0,homeTeam_data.columns.get_loc('venue_capacity')]),16,font_1,white_color_report,True,None) 
        if type(fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('referee')]) == type(None):
            text_modification_paragraph_center(pres_path,"REFEREE 1","-----",16,font_1,white_color_report,True,None) 
        else:
            text_modification_paragraph_center(pres_path,"REFEREE 1","Arbitre: {}".format(fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('referee')]),16,font_1,white_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"Focus on HOME TEAM","Zoom sur {}".format(homeTeam_name),11,font_1,black_color_report,True,None) 
        text_modification_paragraph_center(pres_path,"Focus on AWAY TEAM","Zoom sur {}".format(awayTeam_name),11,font_1,black_color_report,True,None) 
        text_modification_paragraph_left(pres_path,"HOME RANKING for HOME TEAM","Class. domicile pour {}".format(homeTeam_name),11,font_1,black_color_report,True,None) 
        text_modification_paragraph_left(pres_path,"AWAY RANKING for AWAY TEAM","Class. extérieur pour {}".format(awayTeam_name),11,font_1,black_color_report,True,None)
        text_modification_paragraph_center(pres_path,"Last 20 games biggest win","Plus grosse victoire des {} derniers match".format(last_fixtures_lineups_duration),14,font_1,black_color_report,True,None)
        text_modification_paragraph_center(pres_path,"Last 20 games biggest lose","Plus grosse défaite des {} derniers match".format(last_fixtures_lineups_duration),14,font_1,black_color_report,True,None)
        text_modification_paragraph_center(pres_path,"What about recent shocks","Qu'en est-il des derniers chocs...?",18,font_1,black_color_report,True,None)
        text_modification_paragraph_center(pres_path,"& before that","Et avant...?",18,font_1,draw_color_report,True,True)
        text_modification_paragraph_center(pres_path,"In the last1","Dans les",14,font_1,black_color_report,True,None)
        text_modification_paragraph_center(pres_path,"Games1","derniers matchs",14,font_1,black_color_report,True,None)
        
        text_modification_paragraph_center(pres_path,"Observing the past 20 games parameters including:","Analyse durant les {} derniers matchs des paramètres suivants:".format(last_fixtures_lineups_duration),font_size_titleh1,font_1,black_color_report,True,None)
        text_modification_paragraph_center(pres_path,"Rankings in the league","Classements des équipes",14,font_1,black_color_report,None,True)
        text_modification_paragraph_center(pres_path,"Game results","Résultats des matchs récents",14,font_1,black_color_report,None,True)
        text_modification_paragraph_center(pres_path,"Goals against & for","Buts pour et buts contre",14,font_1,black_color_report,None,True)

        text_modification_paragraph_center(pres_path,"RECENT FORM & PERFORMANCES","FORME & PERFORMANCE RECENTE",font_size_titleh1,font_1,black_color_report,True,None)
        text_modification_paragraph_center(pres_path,"GENERAL RANKING","CLASSEMENT GENERAL",font_size_titleh2,font_1,black_color_report,True,None)
        text_modification_paragraph_center(pres_path,"HOME AWAY RANKING","CLASSEMENT DOMICILE-EXTERIEUR",font_size_titleh2,font_1,black_color_report,True,None)
        text_modification_paragraph_center(pres_path,"GENERAL RESULTS","RESULTATS GENERAUX",font_size_titleh2,font_1,black_color_report,True,None)
        text_modification_paragraph_center(pres_path,"HOME AWAY RESULTS","RESULTATS DOMICILE-EXTERIEUR",font_size_titleh2,font_1,black_color_report,True,None)
        text_modification_paragraph_center(pres_path,"GENERAL GOALS AGAINST & FOR","BP & BC GENERAUX",font_size_titleh2,font_1,black_color_report,True,None)
        text_modification_paragraph_center(pres_path,"HOME AWAY GOALS AGAINST & FOR","BP & BC DOMICILE-EXTERIEUR",font_size_titleh2,font_1,black_color_report,True,None)
                
        # Comments on general rankings
        
        # Ranks
        if int(global_ranking_ppt_home[rank_global_home]['#']-global_ranking_ppt_away[rank_global_away]['#'])>0:
            text_modification_paragraph_center(pres_path,"+xr1","+ "+ "{}".format(int(global_ranking_ppt_home[rank_global_home]['#']-global_ranking_ppt_away[rank_global_away]['#'])),font_size_bubble_number,font_1,black_color_report,True,None) 
            if int(global_ranking_ppt_home[rank_global_home]['#']-global_ranking_ppt_away[rank_global_away]['#'])>1:
                text_modification_paragraph_center(pres_path,"ranks1","rangs",font_size_bubble,font_1,black_color_report,True,None) 
            else:
                text_modification_paragraph_center(pres_path,"ranks1","rang",font_size_bubble,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"for team1","{}".format(homeTeam_name),font_size_bubble,font_1,home_color_report,True,None) 
            advantage_introduction_home += 1 # 3 criteria for the advantage
        elif int(global_ranking_ppt_home[rank_global_home]['#']-global_ranking_ppt_away[rank_global_away]['#'])<0:
            text_modification_paragraph_center(pres_path,"+xr1","+ "+ "{}".format(-int(global_ranking_ppt_home[rank_global_home]['#']-global_ranking_ppt_away[rank_global_away]['#'])),font_size_bubble_number,font_1,black_color_report,True,None) 
            if -int(global_ranking_ppt_home[rank_global_home]['#']-global_ranking_ppt_away[rank_global_away]['#'])>1:
                text_modification_paragraph_center(pres_path,"ranks1","rangs",font_size_bubble,font_1,black_color_report,True,None) 
            else:
                text_modification_paragraph_center(pres_path,"ranks1","rang",font_size_bubble,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"for team1","{}".format(awayTeam_name),font_size_bubble,font_1,away_color_report,True,None) 
            advantage_introduction_away += 1 # 3 criteria for the advantage
        else:
            text_modification_paragraph_center(pres_path,"+xr1","=",font_size_bubble_number,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"ranks1","Rangs identiques",font_size_bubble,font_1,black_color_report,True,None)
            text_modification_paragraph_center(pres_path,"for team1","",font_size_bubble,font_1,black_color_report,True,None)

        # Wins
        if int(global_ranking_ppt_home[rank_global_home]['W']-global_ranking_ppt_away[rank_global_away]['W'])>0:
            text_modification_paragraph_center(pres_path,"+xw1","+ "+ "{}".format(int(global_ranking_ppt_home[rank_global_home]['W']-global_ranking_ppt_away[rank_global_away]['W'])),font_size_bubble_number,font_1,black_color_report,True,None) 
            if int(global_ranking_ppt_home[rank_global_home]['W']-global_ranking_ppt_away[rank_global_away]['W'])>1:
                text_modification_paragraph_center(pres_path,"wins1","victoires",font_size_bubble,font_1,black_color_report,True,None) 
            else:
                text_modification_paragraph_center(pres_path,"wins1","victoire",font_size_bubble,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"for team2","{}".format(homeTeam_name),font_size_bubble,font_1,home_color_report,True,None) 
            advantage_introduction_home += 1 # 4 criteria for the advantage
        elif int(global_ranking_ppt_home[rank_global_home]['W']-global_ranking_ppt_away[rank_global_away]['W'])<0:
            text_modification_paragraph_center(pres_path,"+xw1","+ "+ "{}".format(-int(global_ranking_ppt_home[rank_global_home]['W']-global_ranking_ppt_away[rank_global_away]['W'])),font_size_bubble_number,font_1,black_color_report,True,None) 
            if -int(global_ranking_ppt_home[rank_global_home]['W']-global_ranking_ppt_away[rank_global_away]['W'])>1:
                text_modification_paragraph_center(pres_path,"wins1","victoires",font_size_bubble,font_1,black_color_report,True,None) 
            else:
                text_modification_paragraph_center(pres_path,"wins1","victoire",font_size_bubble,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"for team2","{}".format(awayTeam_name),font_size_bubble,font_1,away_color_report,True,None) 
            advantage_introduction_away += 1 # 4 criteria for the advantage
        else:
            text_modification_paragraph_center(pres_path,"+xw1","=",font_size_bubble_number,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"wins1","#Victoires identiques",font_size_bubble,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"for team2","",font_size_bubble,font_1,black_color_report,True,None) 

        # Goal Average
        if int(global_ranking_ppt_home[rank_global_home]['del']-global_ranking_ppt_away[rank_global_away]['del'])>0:
            text_modification_paragraph_center(pres_path,"+xa1","+ "+ "{}".format(int(global_ranking_ppt_home[rank_global_home]['del']-global_ranking_ppt_away[rank_global_away]['del'])),font_size_bubble_number,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"goal average1","goal average",font_size_bubble,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"for team3","{}".format(homeTeam_name),font_size_bubble,font_1,home_color_report,True,None) 
            advantage_introduction_home += 1 # 5 criteria for the advantage
        elif int(global_ranking_ppt_home[rank_global_home]['del']-global_ranking_ppt_away[rank_global_away]['del'])<0:
            text_modification_paragraph_center(pres_path,"+xa1","+ "+ "{}".format(-int(global_ranking_ppt_home[rank_global_home]['del']-global_ranking_ppt_away[rank_global_away]['del'])),font_size_bubble_number,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"goal average1","goal average",font_size_bubble,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"for team3","{}".format(awayTeam_name),font_size_bubble,font_1,away_color_report,True,None) 
            advantage_introduction_away += 1 # 5 criteria for the advantage
        else:
            text_modification_paragraph_center(pres_path,"+xa1","=",font_size_bubble_number,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"goal average1","Goal average identique",font_size_bubble,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"for team3","",font_size_bubble,font_1,black_color_report,True,None) 

        # Comments on Home / Away rank
        # Ranks
        if int(home_ranking_ppt[rank_home]['#']-away_ranking_ppt[rank_away]['#'])>0:
            text_modification_paragraph_center(pres_path,"+xr2","+ "+ "{}".format(int(home_ranking_ppt[rank_home]['#']-away_ranking_ppt[rank_away]['#'])),font_size_bubble_number,font_1,black_color_report,True,None) 
            if int(home_ranking_ppt[rank_home]['#']-away_ranking_ppt[rank_away]['#'])>1:
                text_modification_paragraph_center(pres_path,"ranks2","rangs",font_size_bubble,font_1,black_color_report,True,None) 
            else:
                text_modification_paragraph_center(pres_path,"ranks2","rang",font_size_bubble,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"for team4","{}".format(homeTeam_name),font_size_bubble,font_1,home_color_report,True,None) 
            advantage_introduction_home += 1 # 6 criteria for the advantage
        elif int(home_ranking_ppt[rank_home]['#']-away_ranking_ppt[rank_away]['#'])<0:
            text_modification_paragraph_center(pres_path,"+xr2","+ "+ "{}".format(-int(home_ranking_ppt[rank_home]['#']-away_ranking_ppt[rank_away]['#'])),font_size_bubble_number,font_1,black_color_report,True,None) 
            if -int(home_ranking_ppt[rank_home]['#']-away_ranking_ppt[rank_away]['#'])>1:
                text_modification_paragraph_center(pres_path,"ranks2","rangs",font_size_bubble,font_1,black_color_report,True,None) 
            else:
                text_modification_paragraph_center(pres_path,"ranks2","rang",font_size_bubble,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"for team4","{}".format(awayTeam_name),font_size_bubble,font_1,away_color_report,True,None) 
            advantage_introduction_away += 1 # 6 criteria for the advantage
        else:
            text_modification_paragraph_center(pres_path,"+xr2","=",font_size_bubble_number,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"ranks2","Rangs identiques",font_size_bubble,font_1,black_color_report,True,None)
            text_modification_paragraph_center(pres_path,"for team4","",font_size_bubble,font_1,black_color_report,True,None)

        # Wins
        if int(home_ranking_ppt[rank_home]['W']-away_ranking_ppt[rank_away]['W'])>0:
            text_modification_paragraph_center(pres_path,"+xw2","+ "+ "{}".format(int(home_ranking_ppt[rank_home]['W']-away_ranking_ppt[rank_away]['W'])),font_size_bubble_number,font_1,black_color_report,True,None) 
            if int(home_ranking_ppt[rank_home]['W']-away_ranking_ppt[rank_away]['W']):
                text_modification_paragraph_center(pres_path,"wins2","victoires",font_size_bubble,font_1,black_color_report,True,None) 
            else:
                text_modification_paragraph_center(pres_path,"wins2","victoire",font_size_bubble,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"for team5","{}".format(homeTeam_name),font_size_bubble,font_1,home_color_report,True,None) 
            advantage_introduction_home += 1 # 7 criteria for the advantage
        elif int(home_ranking_ppt[rank_home]['W']-away_ranking_ppt[rank_away]['W'])<0:
            text_modification_paragraph_center(pres_path,"+xw2","+ "+ "{}".format(-int(home_ranking_ppt[rank_home]['W']-away_ranking_ppt[rank_away]['W'])),font_size_bubble_number,font_1,black_color_report,True,None) 
            if -int(home_ranking_ppt[rank_home]['W']-away_ranking_ppt[rank_away]['W'])>1:
                text_modification_paragraph_center(pres_path,"wins2","victoires",font_size_bubble,font_1,black_color_report,True,None) 
            else:
                text_modification_paragraph_center(pres_path,"wins2","victoire",font_size_bubble,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"for team5","{}".format(awayTeam_name),font_size_bubble,font_1,away_color_report,True,None) 
            advantage_introduction_away += 1 # 7 criteria for the advantage
        else:
            text_modification_paragraph_center(pres_path,"+xw2","=",font_size_bubble_number,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"wins2","#Victoires identiques",font_size_bubble,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"for team5","",font_size_bubble,font_1,black_color_report,True,None) 

        # Goal Average
        if int(home_ranking_ppt[rank_home]['del']-away_ranking_ppt[rank_away]['del'])>0:
            text_modification_paragraph_center(pres_path,"+xa2","+ "+ "{}".format(int(home_ranking_ppt[rank_home]['del']-away_ranking_ppt[rank_away]['del'])),font_size_bubble_number,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"goal average2","goal average",font_size_bubble,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"for team6","{}".format(homeTeam_name),font_size_bubble,font_1,home_color_report,True,None) 
            advantage_introduction_home += 1 # 8 criteria for the advantage
        elif int(home_ranking_ppt[rank_home]['del']-away_ranking_ppt[rank_away]['del'])<0:
            text_modification_paragraph_center(pres_path,"+xa2","+ "+ "{}".format(-int(home_ranking_ppt[rank_home]['del']-away_ranking_ppt[rank_away]['del'])),font_size_bubble_number,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"goal average2","goal average",font_size_bubble,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"for team6","{}".format(awayTeam_name),font_size_bubble,font_1,away_color_report,True,None) 
            advantage_introduction_away += 1 # 8 criteria for the advantage
        else:
            text_modification_paragraph_center(pres_path,"+xa2","=",font_size_bubble_number,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"goal average2","Goal average identique",font_size_bubble,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"for team6","",font_size_bubble,font_1,black_color_report,True,None) 
        
        # Comments on figure 1

        # Assignement of variable
        p = 0
        
        if values_home[2]/sum(values_home) != values_away[2]/sum(values_away): 
            if values_home[2]/sum(values_home)>values_away[2]/sum(values_away):
                p = int(100*round(values_home[2]/sum(values_home)-values_away[2]/sum(values_away),2))
                text_modification_paragraph_center(pres_path,"+XA%","+ "+ "{}%".format(p),20,font_1,black_color_report,True,None) 
                text_modification_paragraph_center(pres_path,"wins3","victoires",font_size_bubble,font_1,black_color_report,True,None) 
                text_modification_paragraph_center(pres_path,"for team7","{}".format(homeTeam_name),font_size_bubble,font_1,home_color_report,True,None) 
                advantage_introduction_home += 1 # 9 criteria for the advantage
            else:
                p = int(100*round(-values_home[2]/sum(values_home)+values_away[2]/sum(values_away),2))
                text_modification_paragraph_center(pres_path,"+XA%","+ "+ "{}%".format(-p),20,font_1,black_color_report,True,None) 
                text_modification_paragraph_center(pres_path,"wins3","victoires",font_size_bubble,font_1,black_color_report,True,None) 
                text_modification_paragraph_center(pres_path,"for team7","{}".format(awayTeam_name),font_size_bubble,font_1,away_color_report,True,None) 
                advantage_introduction_away += 1 # 9 criteria for the advantage
        
        elif values_home[1]/sum(values_home) != values_away[1]/sum(values_away): 
            if values_home[1]/sum(values_home)>values_away[1]/sum(values_away):
                p = int(100*round(values_home[1]/sum(values_home)-values_away[1]/sum(values_away),2))
                text_modification_paragraph_center(pres_path,"+XA%","+ "+ str(p) +  "%",20,font_1,black_color_report,True,None) 
                text_modification_paragraph_center(pres_path,"wins3","défaites",font_size_bubble,font_1,black_color_report,True,None) 
                text_modification_paragraph_center(pres_path,"for team7","{}".format(homeTeam_name),font_size_bubble,font_1,home_color_report,True,None) 
                advantage_introduction_away += 1 # 9 criteria for the advantage
            else:
                p = int(100*round(-values_home[1]/sum(values_home)+values_away[1]/sum(values_away),2))
                text_modification_paragraph_center(pres_path,"+XA%","+ "+ str(-p) +  "%",20,font_1,black_color_report,True,None) 
                text_modification_paragraph_center(pres_path,"wins3","défaites",font_size_bubble,font_1,black_color_report,True,None) 
                text_modification_paragraph_center(pres_path,"for team7","{}".format(awayTeam_name),font_size_bubble,font_1,away_color_report,True,None) 
                advantage_introduction_home += 1 # 9 criteria for the advantage
        else:
            text_modification_paragraph_center(pres_path,"+XA%","=",24,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"wins3","Même forme",font_size_bubble,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"for team7","",font_size_bubble,font_1,black_color_report,True,None)         

        if values_home_home[2]/sum(values_home_home) != values_away_away[2]/sum(values_away_away): 
            if values_home_home[2]/sum(values_home_home)>values_away_away[2]/sum(values_away_away):
                p = int(100*round(values_home_home[2]/sum(values_home_home)-values_away_away[2]/sum(values_away_away),2))
                text_modification_paragraph_center(pres_path,"+XA2%","+ "+ "{}%".format(p),20,font_1,black_color_report,True,None) 
                text_modification_paragraph_center(pres_path,"wins4","victoires",font_size_bubble,font_1,black_color_report,True,None) 
                text_modification_paragraph_center(pres_path,"for team9","{}".format(homeTeam_name),font_size_bubble,font_1,home_color_report,True,None) 
                advantage_introduction_home += 1 # 10 criteria for the advantage
            else:
                p = int(100*round(-values_home_home[2]/sum(values_home_home)+values_away_away[2]/sum(values_away_away),2))
                text_modification_paragraph_center(pres_path,"+XA2%","+ "+ "{}%".format(-p),20,font_1,black_color_report,True,None) 
                text_modification_paragraph_center(pres_path,"wins4","victoires",font_size_bubble,font_1,black_color_report,True,None) 
                text_modification_paragraph_center(pres_path,"for team9","{}".format(awayTeam_name),font_size_bubble,font_1,away_color_report,True,None) 
                advantage_introduction_away += 1 # 10 criteria for the advantage
        
        elif values_home_home[1]/sum(values_home_home) != values_away_away[1]/sum(values_away_away): 
            if values_home_home[1]/sum(values_home_home)>values_away_away[1]/sum(values_away_away):
                p = int(100*round(values_home_home[1]/sum(values_home_home)-values_away_away[1]/sum(values_away_away),2))
                text_modification_paragraph_center(pres_path,"+XA2%","+ "+ str(p) +  "%",20,font_1,black_color_report,True,None) 
                text_modification_paragraph_center(pres_path,"wins4","défaites",font_size_bubble,font_1,black_color_report,True,None) 
                text_modification_paragraph_center(pres_path,"for team9","{}".format(homeTeam_name),font_size_bubble,font_1,home_color_report,True,None) 
                advantage_introduction_away += 1 # 10 criteria for the advantage
            else:
                p = int(100*round(-values_home_home[1]/sum(values_home_home)+values_away_away[1]/sum(values_away_away),2))
                text_modification_paragraph_center(pres_path,"+XA2%","+ "+ str(-p) +  "%",20,font_1,black_color_report,True,None) 
                text_modification_paragraph_center(pres_path,"wins4","défaites",font_size_bubble,font_1,black_color_report,True,None) 
                text_modification_paragraph_center(pres_path,"for team9","{}".format(awayTeam_name),font_size_bubble,font_1,away_color_report,True,None) 
                advantage_introduction_home += 1 # 10 criteria for the advantage
        else:
            text_modification_paragraph_center(pres_path,"+XA2%","=",24,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"wins4","Même forme",font_size_bubble,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"for team9","",font_size_bubble,font_1,black_color_report,True,None)     

        # Comments on general performances
        if int(values_home[1]) != int(values_away[1]):
            if int(values_home[1])>int(values_away[1]):
                text_modification_paragraph_left(pres_path,"More defeats for1","Plus de défaites pour {}".format(homeTeam_name),16,font_1,lose_color_report,True,None) 
            else:
                text_modification_paragraph_left(pres_path,"More defeats for1","Plus de défaites pour {}".format(awayTeam_name),16,font_1,lose_color_report,True,None) 
        else:
                text_modification_paragraph_left(pres_path,"More defeats for1","Même nombre de défaites",16,font_1,lose_color_report,True,None) 

        if int(values_home[2]) != int(values_away[2]):
            if int(values_home[2])>int(values_away[2]):
                text_modification_paragraph_left(pres_path,"More wins for1","Plus de victoires pour {}".format(homeTeam_name),16,font_1,win_color_report,True,None) 
            else:
                text_modification_paragraph_left(pres_path,"More wins for1","Plus de victoires pour {}".format(awayTeam_name),16,font_1,win_color_report,True,None) 
        else:
                text_modification_paragraph_left(pres_path,"More wins for1","Même nombre de victoires",16,font_1,win_color_report,True,None) 

        if int(values_home[0]) != int(values_away[0]):
            if int(values_home[0])>int(values_away[0]):
                text_modification_paragraph_left(pres_path,"More draws for1","Plus de nuls pour {}".format(homeTeam_name),16,font_1,black_color_report,True,None) 
            else:
                text_modification_paragraph_left(pres_path,"More draws for1","Plus de nuls pour {}".format(awayTeam_name),16,font_1,black_color_report,True,None) 
        else:
                text_modification_paragraph_left(pres_path,"More draws for1","Même nombre de nuls",16,font_1,black_color_report,True,None)   
        
        # Comments on home-away performances
        if int(values_home_home[1]) != int(values_away_away[1]):
            if int(values_home_home[1])>int(values_away_away[1]):
                text_modification_paragraph_left(pres_path,"More defeats for2","Plus de défaites pour {}".format(homeTeam_name),16,font_1,lose_color_report,True,None) 
            else:
                text_modification_paragraph_left(pres_path,"More defeats for2","Plus de défaites pour {}".format(awayTeam_name),16,font_1,lose_color_report,True,None) 
        else:
                text_modification_paragraph_left(pres_path,"More defeats for2","Même nombre de défaites",16,font_1,lose_color_report,True,None) 

        if int(values_home_home[2]) != int(values_away_away[2]):
            if int(values_home_home[2])>int(values_away_away[2]):
                text_modification_paragraph_left(pres_path,"More wins for2","Plus de victoires pour {}".format(homeTeam_name),16,font_1,win_color_report,True,None) 
            else:
                text_modification_paragraph_left(pres_path,"More wins for2","Plus de victoires pour {}".format(awayTeam_name),16,font_1,win_color_report,True,None) 
        else:
                text_modification_paragraph_left(pres_path,"More wins for2","Même nombre de victoires",16,font_1,win_color_report,True,None) 

        if int(values_home_home[0]) != int(values_away_away[0]):
            if int(values_home_home[0])>int(values_away_away[0]):
                text_modification_paragraph_left(pres_path,"More draws for2","Plus de nuls pour {}".format(homeTeam_name),16,font_1,black_color_report,True,None) 
            else:
                text_modification_paragraph_left(pres_path,"More draws for2","Plus de nuls pour {}".format(awayTeam_name),16,font_1,black_color_report,True,None) 
        else:
                text_modification_paragraph_left(pres_path,"More draws for2","Même nombre de nuls",16,font_1,black_color_report,True,None)   

        # Comments on figure 2
        # Assignement of variable
        p = 0
        
        if delta_home_mean>delta_away_mean:
            text_modification_paragraph_center(pres_path,"+XA3","+ "+ "{}".format(round(delta_home_mean-delta_away_mean,1)),28,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"goal average3","goal average",font_size_bubble,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"for team8","{}".format(homeTeam_name),font_size_bubble,font_1,home_color_report,True,None) 
            advantage_introduction_home += 1 # 11 criteria for the advantage
        elif delta_home_mean<delta_away_mean:
            text_modification_paragraph_center(pres_path,"+XA3","+ "+ "{}".format(round(-delta_home_mean+delta_away_mean,1)),28,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"goal average3","goal average",font_size_bubble,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"for team8","{}".format(awayTeam_name),font_size_bubble,font_1,away_color_report,True,None) 
            advantage_introduction_away += 1 # 11 criteria for the advantage
        else:
            text_modification_paragraph_center(pres_path,"+XA3","=",font_size_bubble_number,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"goal average3","Même goal average",font_size_bubble,font_1,black_color_report,True,None) 
            text_modification_paragraph_center(pres_path,"for team8","",font_size_bubble,font_1,black_color_report,True,None) 
        
        if var_home_mean>var_away_mean:
            text_modification_paragraph_center(pres_path,"More variance for team","Resultats plus hétérogènes pour {} (variance: {})".format(homeTeam_name,round(var_home_mean,1)),16,font_1,black_color_report,True,None) 
        elif var_home_mean<var_away_mean:
            text_modification_paragraph_center(pres_path,"More variance for team","Resultats plus hétérogènes pour {} (variance: {})".format(awayTeam_name,round(var_away_mean,1)),16,font_1,black_color_report,True,None) 
        else:
            text_modification_paragraph_center(pres_path,"More variance for team","Resultats autant dispersés pour les 2 équipes (variance: {})".format(round(var_away_mean,1)),16,font_1,black_color_report,True,None) 

        if delta_home_mean_home>delta_away_mean_away:
            text_modification_paragraph_center(pres_path,"Difference of x in adavantage of x","Difference de {} à l'avantage de {}".format(round(delta_home_mean_home-delta_away_mean_away,1),homeTeam_name),18,font_1,black_color_report,True,None)   
        else:
            text_modification_paragraph_center(pres_path,"Difference of x in adavantage of x","Difference de {} à l'avantage de {}".format(round(-delta_home_mean_home+delta_away_mean_away,1),awayTeam_name),18,font_1,black_color_report,True,None)   


    ######################################## ENGLISH BELOW ########################################
    ######################################## END OF ENGLISH ########################################
    
    introduction_adv = 0
    
    if advantage_introduction_home>advantage_introduction_away:
        text_modification_paragraph_center(pres_path,"+ ADVANTAGE INTRODUCTION","+ " + str(homeTeam_name),24,font_1,white_color_report,True,None)  
        introduction_adv = math.ceil(10*advantage_introduction_home/(advantage_introduction_home+advantage_introduction_away))
        
    elif advantage_introduction_home<advantage_introduction_away:
        text_modification_paragraph_center(pres_path,"+ ADVANTAGE INTRODUCTION","+ " + str(awayTeam_name),24,font_1,white_color_report,True,None)  
        introduction_adv = math.ceil(10*advantage_introduction_away/(advantage_introduction_home+advantage_introduction_away))
        
    else:
        text_modification_paragraph_center(pres_path,"+ ADVANTAGE INTRODUCTION","Perfect equality",24,font_1,white_color_report,True,None)  

    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)
    
    # Photos 
    print("ok")
    for cpt_5 in range (1,8):
        print(int(global_ranking_ppt_home[cpt_5]['team_id']))

        page_number, pos_left, pos_top = matrix_positions_introduction['introduction_global_photo_home_{}'.format(cpt_5)]['page_number'], matrix_positions_introduction['introduction_global_photo_home_{}'.format(cpt_5)]['pos_left'],matrix_positions_introduction['introduction_global_photo_home_{}'.format(cpt_5)]['pos_top'] 
        pic = prs.slides[page_number].shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(global_ranking_ppt_home[cpt_5]['team_id'])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_3,prs))
        
        page_number, pos_left, pos_top = matrix_positions_introduction['introduction_global_photo_away_{}'.format(cpt_5)]['page_number'], matrix_positions_introduction['introduction_global_photo_away_{}'.format(cpt_5)]['pos_left'],matrix_positions_introduction['introduction_global_photo_away_{}'.format(cpt_5)]['pos_top'] 
        pic = prs.slides[page_number].shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(global_ranking_ppt_away[cpt_5]['team_id'])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_3,prs))

        page_number, pos_left, pos_top = matrix_positions_introduction['introduction_home_photo_{}'.format(cpt_5)]['page_number'], matrix_positions_introduction['introduction_home_photo_{}'.format(cpt_5)]['pos_left'],matrix_positions_introduction['introduction_home_photo_{}'.format(cpt_5)]['pos_top'] 
        pic = prs.slides[page_number].shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(home_ranking_ppt[cpt_5]['team_id'])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_3,prs))
    
        page_number, pos_left, pos_top = matrix_positions_introduction['introduction_away_photo_{}'.format(cpt_5)]['page_number'], matrix_positions_introduction['introduction_away_photo_{}'.format(cpt_5)]['pos_left'],matrix_positions_introduction['introduction_away_photo_{}'.format(cpt_5)]['pos_top'] 
        pic = prs.slides[page_number].shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(away_ranking_ppt[cpt_5]['team_id'])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_3,prs)) 
        
    
    # The logos of the biggest wins and defeats
    if df_raw_home_1.iloc[delta_max_home,df_raw_home_1.columns.get_loc('away_home')] == "home":
            page_number, pos_left, pos_top = matrix_positions_introduction['logo_1']['page_number'], matrix_positions_introduction['logo_1']['pos_left'],matrix_positions_introduction['logo_1']['pos_top'] 
            slide = prs.slides[page_number]
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(homeTeam_id) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))

            page_number, pos_left, pos_top = matrix_positions_introduction['logo_2']['page_number'], matrix_positions_introduction['logo_2']['pos_left'],matrix_positions_introduction['logo_2']['pos_top'] 
            slide = prs.slides[page_number]
            try:
                pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_home_1.iloc[delta_max_home,df_raw_home_1.columns.get_loc('opp_team_id')])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))
            except:
                logo_team_id_s(int(df_raw_home_1.iloc[delta_max_home,df_raw_home_1.columns.get_loc('opp_team_id')]))
                pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_home_1.iloc[delta_max_home,df_raw_home_1.columns.get_loc('opp_team_id')])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))
    else:
        page_number, pos_left, pos_top = matrix_positions_introduction['logo_2']['page_number'], matrix_positions_introduction['logo_2']['pos_left'],matrix_positions_introduction['logo_2']['pos_top'] 
        slide = prs.slides[page_number]
        pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(homeTeam_id) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))

        page_number, pos_left, pos_top = matrix_positions_introduction['logo_1']['page_number'], matrix_positions_introduction['logo_1']['pos_left'],matrix_positions_introduction['logo_1']['pos_top'] 
        slide = prs.slides[page_number]
        try:
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_home_1.iloc[delta_max_home,df_raw_home_1.columns.get_loc('opp_team_id')])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))
        except:
            logo_team_id_s(int(df_raw_home_1.iloc[delta_max_home,df_raw_home_1.columns.get_loc('opp_team_id')]))
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_home_1.iloc[delta_max_home,df_raw_home_1.columns.get_loc('opp_team_id')])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))

    if df_raw_home_1.iloc[delta_min_home,df_raw_home_1.columns.get_loc('away_home')] == "home":
        page_number, pos_left, pos_top = matrix_positions_introduction['logo_5']['page_number'], matrix_positions_introduction['logo_5']['pos_left'],matrix_positions_introduction['logo_5']['pos_top'] 
        slide = prs.slides[page_number]
        pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(homeTeam_id) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))

        page_number, pos_left, pos_top = matrix_positions_introduction['logo_6']['page_number'], matrix_positions_introduction['logo_6']['pos_left'],matrix_positions_introduction['logo_6']['pos_top'] 
        slide = prs.slides[page_number]
        try:
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_home_1.iloc[delta_min_home,df_raw_home_1.columns.get_loc('opp_team_id')])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))
        except:
            logo_team_id_s(int(df_raw_home_1.iloc[delta_min_home,df_raw_home_1.columns.get_loc('opp_team_id')]))
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_home_1.iloc[delta_min_home,df_raw_home_1.columns.get_loc('opp_team_id')])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))
    else:
        page_number, pos_left, pos_top = matrix_positions_introduction['logo_6']['page_number'], matrix_positions_introduction['logo_6']['pos_left'],matrix_positions_introduction['logo_6']['pos_top'] 
        slide = prs.slides[page_number]
        pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(homeTeam_id) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))

        page_number, pos_left, pos_top = matrix_positions_introduction['logo_5']['page_number'], matrix_positions_introduction['logo_5']['pos_left'],matrix_positions_introduction['logo_5']['pos_top'] 
        slide = prs.slides[page_number]
        try:
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_home_1.iloc[delta_min_home,df_raw_home_1.columns.get_loc('opp_team_id')])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))
        except:
            logo_team_id_s(int(df_raw_home_1.iloc[delta_min_home,df_raw_home_1.columns.get_loc('opp_team_id')]))
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_home_1.iloc[delta_min_home,df_raw_home_1.columns.get_loc('opp_team_id')])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))
    
    if df_raw_away_1.iloc[delta_max_away,df_raw_away_1.columns.get_loc('away_home')] == "home":
        page_number, pos_left, pos_top = matrix_positions_introduction['logo_3']['page_number'], matrix_positions_introduction['logo_3']['pos_left'],matrix_positions_introduction['logo_3']['pos_top'] 
        slide = prs.slides[page_number]
        pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(awayTeam_id) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))

        page_number, pos_left, pos_top = matrix_positions_introduction['logo_4']['page_number'], matrix_positions_introduction['logo_4']['pos_left'],matrix_positions_introduction['logo_4']['pos_top'] 
        slide = prs.slides[page_number]
        try:
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_away_1.iloc[delta_max_away,df_raw_away_1.columns.get_loc('opp_team_id')])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))
        except:
            logo_team_id_s(int(df_raw_away_1.iloc[delta_max_away,df_raw_away_1.columns.get_loc('opp_team_id')]))
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_away_1.iloc[delta_max_away,df_raw_away_1.columns.get_loc('opp_team_id')])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))
    else:
        page_number, pos_left, pos_top = matrix_positions_introduction['logo_4']['page_number'], matrix_positions_introduction['logo_4']['pos_left'],matrix_positions_introduction['logo_4']['pos_top'] 
        slide = prs.slides[page_number]
        pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(awayTeam_id) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))

        page_number, pos_left, pos_top = matrix_positions_introduction['logo_3']['page_number'], matrix_positions_introduction['logo_3']['pos_left'],matrix_positions_introduction['logo_3']['pos_top'] 
        slide = prs.slides[page_number]
        try:
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_away_1.iloc[delta_max_away,df_raw_away_1.columns.get_loc('opp_team_id')])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))
        except:
            logo_team_id_s(int(df_raw_away_1.iloc[delta_max_away,df_raw_away_1.columns.get_loc('opp_team_id')]))
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_away_1.iloc[delta_max_away,df_raw_away_1.columns.get_loc('opp_team_id')])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))

    if df_raw_away_1.iloc[delta_min_away,df_raw_away_1.columns.get_loc('away_home')] == "home":
        page_number, pos_left, pos_top = matrix_positions_introduction['logo_7']['page_number'], matrix_positions_introduction['logo_7']['pos_left'],matrix_positions_introduction['logo_7']['pos_top'] 
        slide = prs.slides[page_number]
        pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(awayTeam_id) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))

        page_number, pos_left, pos_top = matrix_positions_introduction['logo_8']['page_number'], matrix_positions_introduction['logo_8']['pos_left'],matrix_positions_introduction['logo_8']['pos_top'] 
        slide = prs.slides[page_number]
        try:
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_away_1.iloc[delta_min_away,df_raw_away_1.columns.get_loc('opp_team_id')])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))
        except:
            logo_team_id_s(int(df_raw_away_1.iloc[delta_min_away,df_raw_away_1.columns.get_loc('opp_team_id')]))
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_away_1.iloc[delta_min_away,df_raw_away_1.columns.get_loc('opp_team_id')])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))
    else:
        page_number, pos_left, pos_top = matrix_positions_introduction['logo_8']['page_number'], matrix_positions_introduction['logo_8']['pos_left'],matrix_positions_introduction['logo_8']['pos_top'] 
        slide = prs.slides[page_number]
        pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(awayTeam_id) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))

        page_number, pos_left, pos_top = matrix_positions_introduction['logo_7']['page_number'], matrix_positions_introduction['logo_7']['pos_left'],matrix_positions_introduction['logo_7']['pos_top'] 
        slide = prs.slides[page_number]
        try:
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_away_1.iloc[delta_min_away,df_raw_away_1.columns.get_loc('opp_team_id')])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))
        except:
            logo_team_id_s(int(df_raw_away_1.iloc[delta_min_away,df_raw_away_1.columns.get_loc('opp_team_id')]))
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_away_1.iloc[delta_min_away,df_raw_away_1.columns.get_loc('opp_team_id')])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))

    if df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest win')]["homeaway"] == "home":
        page_number, pos_left, pos_top = matrix_positions_introduction['logo_9']['page_number'], matrix_positions_introduction['logo_9']['pos_left'],matrix_positions_introduction['logo_9']['pos_top'] 
        slide = prs.slides[page_number]
        pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(homeTeam_id) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))

        page_number, pos_left, pos_top = matrix_positions_introduction['logo_10']['page_number'], matrix_positions_introduction['logo_10']['pos_left'],matrix_positions_introduction['logo_10']['pos_top'] 
        slide = prs.slides[page_number]
        try:
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest win')]["team_id"])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))
        except:
            logo_team_id_s(int(df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest win')]["team_id"]))
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest win')]["team_id"])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))
    else:
        page_number, pos_left, pos_top = matrix_positions_introduction['logo_10']['page_number'], matrix_positions_introduction['logo_10']['pos_left'],matrix_positions_introduction['logo_10']['pos_top'] 
        slide = prs.slides[page_number]
        pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(homeTeam_id) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))

        page_number, pos_left, pos_top = matrix_positions_introduction['logo_9']['page_number'], matrix_positions_introduction['logo_9']['pos_left'],matrix_positions_introduction['logo_9']['pos_top'] 
        slide = prs.slides[page_number]
        try:
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest win')]["team_id"])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))
        except:
            logo_team_id_s(int(df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest win')]["team_id"]))
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest win')]["team_id"])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))


    if df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest defeat')]["homeaway"] == "home":
        page_number, pos_left, pos_top = matrix_positions_introduction['logo_13']['page_number'], matrix_positions_introduction['logo_13']['pos_left'],matrix_positions_introduction['logo_13']['pos_top'] 
        slide = prs.slides[page_number]
        pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(homeTeam_id) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))

        page_number, pos_left, pos_top = matrix_positions_introduction['logo_14']['page_number'], matrix_positions_introduction['logo_14']['pos_left'],matrix_positions_introduction['logo_14']['pos_top'] 
        slide = prs.slides[page_number]
        try:
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest defeat')]["team_id"])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))
        except:
            logo_team_id_s(int(df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest defeat')]["team_id"]))
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest defeat')]["team_id"])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))
    else:
        page_number, pos_left, pos_top = matrix_positions_introduction['logo_14']['page_number'], matrix_positions_introduction['logo_14']['pos_left'],matrix_positions_introduction['logo_14']['pos_top'] 
        slide = prs.slides[page_number]
        pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(homeTeam_id) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))

        page_number, pos_left, pos_top = matrix_positions_introduction['logo_13']['page_number'], matrix_positions_introduction['logo_13']['pos_left'],matrix_positions_introduction['logo_13']['pos_top'] 
        slide = prs.slides[page_number]
        try:
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest defeat')]["team_id"])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))
        except:
            logo_team_id_s(int(df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest defeat')]["team_id"]))
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('biggest defeat')]["team_id"])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))

    if df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest win')]["homeaway"] == "home":
        page_number, pos_left, pos_top = matrix_positions_introduction['logo_11']['page_number'], matrix_positions_introduction['logo_11']['pos_left'],matrix_positions_introduction['logo_11']['pos_top'] 
        slide = prs.slides[page_number]
        pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(awayTeam_id) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))

        page_number, pos_left, pos_top = matrix_positions_introduction['logo_12']['page_number'], matrix_positions_introduction['logo_12']['pos_left'],matrix_positions_introduction['logo_12']['pos_top'] 
        slide = prs.slides[page_number]
        try:
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest win')]["team_id"])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))
        except:
            logo_team_id_s(int(df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest win')]["team_id"]))
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest win')]["team_id"])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))
    else:
        page_number, pos_left, pos_top = matrix_positions_introduction['logo_12']['page_number'], matrix_positions_introduction['logo_12']['pos_left'],matrix_positions_introduction['logo_12']['pos_top'] 
        slide = prs.slides[page_number]
        pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(awayTeam_id) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))

        page_number, pos_left, pos_top = matrix_positions_introduction['logo_11']['page_number'], matrix_positions_introduction['logo_11']['pos_left'],matrix_positions_introduction['logo_11']['pos_top'] 
        slide = prs.slides[page_number]
        try:
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest win')]["team_id"])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))
        except:
            logo_team_id_s(int(df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest win')]["team_id"]))
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest win')]["team_id"])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))

    if df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest defeat')]["homeaway"] == "home":
        page_number, pos_left, pos_top = matrix_positions_introduction['logo_15']['page_number'], matrix_positions_introduction['logo_15']['pos_left'],matrix_positions_introduction['logo_15']['pos_top'] 
        slide = prs.slides[page_number]
        pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(awayTeam_id) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))

        page_number, pos_left, pos_top = matrix_positions_introduction['logo_16']['page_number'], matrix_positions_introduction['logo_16']['pos_left'],matrix_positions_introduction['logo_16']['pos_top'] 
        slide = prs.slides[page_number]
        try:
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest defeat')]["team_id"])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))
        except:
            logo_team_id_s(int(df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest defeat')]["team_id"]))
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest defeat')]["team_id"])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))
    else:
        page_number, pos_left, pos_top = matrix_positions_introduction['logo_16']['page_number'], matrix_positions_introduction['logo_16']['pos_left'],matrix_positions_introduction['logo_16']['pos_top'] 
        slide = prs.slides[page_number]
        pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(awayTeam_id) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))

        page_number, pos_left, pos_top = matrix_positions_introduction['logo_15']['page_number'], matrix_positions_introduction['logo_15']['pos_left'],matrix_positions_introduction['logo_15']['pos_top'] 
        slide = prs.slides[page_number]
        try:
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest defeat')]["team_id"])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))
        except:
            logo_team_id_s(int(df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest defeat')]["team_id"]))
            pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(int(df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('biggest defeat')]["team_id"])) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_4,prs))
    
    prs.save(pres_path)

    return introduction_adv


''' 
Creation function 
------------------------------------------------------------------------------------
'''

# OK -
def introduction_reports_chart_section1(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_introduction,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3): #NameOfReports_Reports_SectionXX
    
    ticks_size = 13
    label_size = ticks_size*1.15
    legend_size = ticks_size*1.15
    title_size = ticks_size*1.4
    
    fixtures_home_away = pd.read_pickle(link_data+'{}'.format(team_id)+'/'+str(game_fixture_id) + '/' +'fixtures_package'+'/'+'fixtures_team_id_next_s'+ '_'+ str(team_id) +'_'+ str(nex) +'.pkl')
    
    # A. Language parameters
    if english == 1 :
        labels = ["Draw","Lose","Win"] # Keep the same order
    elif english == 0 :
        labels = ["Nul","Défaite","Victoire"] # Keep the same order
    
    fig, ax = plt.subplots(1,2)

    # B. Plot Graph 1
    df_home = df_raw_home_1.groupby('W-D-L').size() 
    df_home = df_home.sort_index(ascending=True)

    df_away = df_raw_away_1.groupby('W-D-L').size() 
    df_away = df_away.sort_index(ascending=True)
    
    # Assignements of variables
    values_home, values_away = list(df_home), list(df_away)

    patches, texts, autotexts = ax[0].pie(values_home, autopct=make_autopct(values_home),shadow=False, labels = labels, colors = [draw_color_plot,lose_color_plot,win_color_plot], wedgeprops = {'linewidth': 6,'edgecolor' : 'white'}, rotatelabels=False,pctdistance=0.6,labeldistance=None, textprops = {'fontsize': label_size,'color':"black",'weight':"bold"})
    #plt.setp(autotexts, size=8, weight="bold") # This is for the labels in the wedges 
    plt.setp(texts, color='grey',weight='bold',size=label_size) # This is for the main labels
    
    center, r = patches[0].center, patches[0].r
    centre_circle = plt.Circle(center,0.25,fc='white')
    edge_circle = plt.Circle(center,r,fill=False,ec=home_color_plot)
    ax[0].add_artist(centre_circle)
    ax[0].add_artist(edge_circle)

    #fig = plt.gcf()
    #ax = fig.gca()

    ax[0].axis('equal')
    ax[0].set_title(homeTeam_name,weight='bold',fontsize=title_size,color= home_color_plot)
    
    # C. Plot Graph 2
    patches, texts, autotexts = ax[1].pie(values_away, autopct=make_autopct(values_away),shadow=False, labels=labels, colors=[draw_color_plot,lose_color_plot,win_color_plot], wedgeprops = {'linewidth': 6,'edgecolor' : 'white'}, rotatelabels=False,pctdistance=0.6,labeldistance=None, textprops = {'fontsize': label_size,'color':"black",'weight':"bold"})
    plt.setp(texts, color='grey',weight='bold',size=label_size) # This is for the main labels
    
    center, r = patches[0].center, patches[0].r
    centre_circle = plt.Circle(center,0.25,fc='white')
    edge_circle = plt.Circle(center,r,fill=False,ec=away_color_plot)
    ax[1].add_artist(centre_circle)
    ax[1].add_artist(edge_circle)

    ax[1].axis('equal')
    ax[1].set_title(awayTeam_name,weight='bold',fontsize=title_size,color = away_color_plot)
    
    # D. Plot Figure
    w_patch, d_patch, l_patch = mpatches.Patch(color=win_color_plot, label=labels[2]), mpatches.Patch(color=draw_color_plot, label=labels[0]),mpatches.Patch(color=lose_color_plot, label=labels[1])
    ax[1].legend(handles=[w_patch,d_patch,l_patch] ,ncol= len(labels),borderpad=1,labelspacing=1,frameon=False, prop={'size': legend_size},loc = 'lower center', bbox_to_anchor = (0,0,1,1),bbox_transform = plt.gcf().transFigure)
    
    #plt.tight_layout()
    plt.savefig(fname=link_reports + str(game_fixture_id) +'/'+name_folder_introduction+'/'+'introduction_reports_chart_section1'+graph_format, dpi = dpi_1)
    plt.show()

    return values_home,values_away


# OK -
def introduction_reports_chart_section3(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_introduction,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3): #NameOfReports_Reports_SectionXX
    
    ticks_size = 13
    label_size = ticks_size*1.15
    legend_size = ticks_size*1.15
    title_size = ticks_size*1.4
    
    fixtures_home_away = pd.read_pickle(link_data+'{}'.format(team_id)+'/'+str(game_fixture_id) + '/' +'fixtures_package'+'/'+'fixtures_team_id_next_s'+ '_'+ str(team_id) +'_'+ str(nex) +'.pkl')
    
    # A. Language parameters
    if english == 1 :
        labels = ["Draw","Lose","Win"] # Keep the same order
    elif english == 0 :
        labels = ["Nul","Défaite","Victoire"] # Keep the same order
    
    fig, ax = plt.subplots(1,2)

    values_home_home, values_away_away = [0,0,0],[0,0,0]

    # D L W is the way it has to be ranked

    for i in range (len(df_raw_home_1)):
        if df_raw_home_1.iloc[i,df_raw_home_1.columns.get_loc('away_home')] == "home":
            if df_raw_home_1.iloc[i,df_raw_home_1.columns.get_loc('W-D-L')] == "W":
                values_home_home[2] += 1
            elif df_raw_home_1.iloc[i,df_raw_home_1.columns.get_loc('W-D-L')] == "D":
                values_home_home[0] += 1
            else:
                values_home_home[1] += 1
        if df_raw_away_1.iloc[i,df_raw_away_1.columns.get_loc('away_home')] == "away":
            if df_raw_away_1.iloc[i,df_raw_away_1.columns.get_loc('W-D-L')] == "W":
                values_away_away[2] += 1
            elif df_raw_away_1.iloc[i,df_raw_away_1.columns.get_loc('W-D-L')] == "D":
                values_away_away[0] += 1
            else:
                values_away_away[1] += 1  

    patches, texts, autotexts = ax[0].pie(values_home_home, autopct=make_autopct(values_home_home),shadow=False, labels = labels, colors = [draw_color_plot,lose_color_plot,win_color_plot], wedgeprops = {'linewidth': 6,'edgecolor' : 'white'}, rotatelabels=False,pctdistance=0.6,labeldistance=None, textprops = {'fontsize': label_size,'color':"black",'weight':"bold"})
    #plt.setp(autotexts, size=8, weight="bold") # This is for the labels in the wedges 
    plt.setp(texts, color='grey',weight='bold',size=label_size) # This is for the main labels
    
    center, r = patches[0].center, patches[0].r
    centre_circle = plt.Circle(center,0.25,fc='white')
    edge_circle = plt.Circle(center,r,fill=False,ec=home_color_plot)
    ax[0].add_artist(centre_circle)
    ax[0].add_artist(edge_circle)

    #fig = plt.gcf()
    #ax = fig.gca()

    ax[0].axis('equal')
    ax[0].set_title(homeTeam_name,weight='bold',fontsize=title_size,color= home_color_plot)
    
    # C. Plot Graph 2
    patches, texts, autotexts = ax[1].pie(values_away_away, autopct=make_autopct(values_away_away),shadow=False, labels=labels, colors=[draw_color_plot,lose_color_plot,win_color_plot], wedgeprops = {'linewidth': 6,'edgecolor' : 'white'}, rotatelabels=False,pctdistance=0.6,labeldistance=None, textprops = {'fontsize': label_size,'color':"black",'weight':"bold"})
    plt.setp(texts, color='grey',weight='bold',size=label_size) # This is for the main labels
    
    center, r = patches[0].center, patches[0].r
    centre_circle = plt.Circle(center,0.25,fc='white')
    edge_circle = plt.Circle(center,r,fill=False,ec=away_color_plot)
    ax[1].add_artist(centre_circle)
    ax[1].add_artist(edge_circle)

    ax[1].axis('equal')
    ax[1].set_title(awayTeam_name,weight='bold',fontsize=title_size,color = away_color_plot)
    
    # D. Plot Figure
    w_patch, d_patch, l_patch = mpatches.Patch(color=win_color_plot, label=labels[2]), mpatches.Patch(color=draw_color_plot, label=labels[0]),mpatches.Patch(color=lose_color_plot, label=labels[1])
    ax[1].legend(handles=[w_patch,d_patch,l_patch] ,ncol= len(labels),borderpad=1,labelspacing=1,frameon=False, prop={'size': legend_size},loc = 'lower center', bbox_to_anchor = (0,0,1,1),bbox_transform = plt.gcf().transFigure)
    
    #plt.tight_layout()
    plt.savefig(fname=link_reports + str(game_fixture_id) +'/'+name_folder_introduction+'/'+'introduction_reports_chart_section3'+graph_format, dpi = dpi_1)
    plt.show()

    return values_home_home, values_away_away

# OK -
def introduction_reports_chart_section2(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_introduction,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3): #NameOfReports_Reports_SectionXX

    ticks_size = 13
    label_size = ticks_size*1.15
    legend_size = ticks_size*1.15
    title_size = ticks_size*1.4
    
    fixtures_home_away = pd.read_pickle(link_data+'{}'.format(team_id)+'/'+str(game_fixture_id) + '/' +'fixtures_package'+'/'+'fixtures_team_id_next_s'+ '_'+ str(team_id) +'_'+ str(nex) +'.pkl')
    homeTeam_id = int(fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('homeTeam')]['team_id'])
    awayTeam_id = int(fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('awayTeam')]['team_id'])
    homeTeam_name = fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('homeTeam')]['team_name']
    awayTeam_name = fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('awayTeam')]['team_name']

    # A. Language parameters
    if english == 1 :
        title = str("Goal difference previous games")
        average = "Average"
    elif english == 0 :
        title = str("Différence de buts")
        average = "Moyenne"
    
    fig, ax = plt.subplots(1,1)
    
    # B. Plot Graph 
    BP_home, BC_home, BP_away, BC_away = 0, 0, 0, 0

    BP_home = list(df_raw_home_1['BP'])
    BC_home = list(df_raw_home_1['BC'])
    BP_away = list(df_raw_away_1['BP'])
    BC_away = list(df_raw_away_1['BC'])

    delta_home = list(np.array(BP_home) - np.array(BC_home))
    delta_away = list(np.array(BP_away) - np.array(BC_away))

    
    for i in range (len(delta_home)):
        try:
            delta_home[i] = int(delta_home[i])
        except:
            delta_home[i] = 0
        try:
            delta_away[i] = int(delta_away[i])
        except:
            delta_away[i] = 0
            
    delta_max_home = delta_home.index(int(max(np.array(delta_home))))
    delta_max_away = delta_away.index(int(max(np.array(delta_away))))
    delta_min_home = delta_home.index(int(min(np.array(delta_home))))
    delta_min_away = delta_away.index(int(min(np.array(delta_away))))

    BP_max_home, BC_max_home, BP_min_home, BC_min_home, BP_max_away, BC_max_away ,BP_min_away ,BC_min_away = BP_home[delta_max_home], BC_home[delta_max_home], BP_home[delta_min_home], BC_home[delta_min_home], BP_away[delta_max_away], BC_away[delta_max_away], BP_away[delta_min_away], BC_away[delta_min_away]

    N = len(delta_home)
    delta_home_mean = [np.mean(delta_home)]*N
    delta_away_mean = [np.mean(delta_away)]*N
    delta_home.reverse()
    delta_away.reverse()
    ind = np.arange(N) 
    width = 0.30       
    ax.bar(ind, delta_home, width, label=homeTeam_name, color=N*[home_color_plot],edgecolor='black')
    ax.bar(ind + width, delta_away, width,label=awayTeam_name, color=N*[away_color_plot],edgecolor='black')
    #ax.set_yticks([], minor=True)
    #ax.set_yticks([])
    ax.plot(ind,delta_home_mean, color=home_color_plot, linestyle='--')
    ax.plot(ind,delta_away_mean, color=away_color_plot, linestyle='--')

    ax.set_title(title,fontsize=title_size,weight='bold')
    plt.xticks(ind + width / 2, ["N-{}".format(i) for i in range (N,0,-1)],fontsize=ticks_size-2)
    plt.yticks(fontsize=ticks_size)
    patch = [mpatches.Patch(color=home_color_plot, label=homeTeam_name), mpatches.Patch(color=away_color_plot, label=awayTeam_name) , matplotlib.lines.Line2D(xdata=[1],ydata=[1],linestyle='--', color='black',label=average)]
    ax.legend(handles=patch,loc='lower center',frameon=False, ncol=len(patch), bbox_to_anchor = (0,0,1,1),bbox_transform = plt.gcf().transFigure,fontsize=legend_size)
    
    every_nth = 3
    for n, label in enumerate(ax.xaxis.get_ticklabels()):
        if n % every_nth != 0:
            label.set_visible(False)

    plt.tight_layout()
    fig.autofmt_xdate()
    #ax.set_edgecolor('white')
    #plt.figure(edgecolor='white')
    #plt.axes(frameon=False)
    plt.savefig(fname=link_reports + str(game_fixture_id) +'/'+name_folder_introduction + '/'+'introduction_reports_chart_section2'+graph_format, dpi = dpi_1)
    plt.show()

    return np.mean(delta_home), np.mean(delta_away), np.var(delta_home), np.var(delta_away), delta_max_home, BP_max_home,BC_max_home, delta_max_away, BP_max_away, BC_max_away, delta_min_home, BP_min_home, BC_min_home, delta_min_away, BP_min_away, BC_min_away

# OK -
def introduction_reports_chart_section4(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_introduction,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3): #NameOfReports_Reports_SectionXX

    ticks_size = 13
    label_size = ticks_size*1.15
    legend_size = ticks_size*1.15
    title_size = ticks_size*1.4
    
    fixtures_home_away = pd.read_pickle(link_data+'{}'.format(team_id)+'/'+str(game_fixture_id) + '/' +'fixtures_package'+'/'+'fixtures_team_id_next_s'+ '_'+ str(team_id) +'_'+ str(nex) +'.pkl')

    # A. Language parameters
    if english == 1 :
        title = str("Goal difference previous games at home")
        average = "Average"
    elif english == 0 :
        title = str("Différence de buts à domicile")
        average = "Moyenne"
    
    fig, ax = plt.subplots(1,1)
    
    # B. Plot Graph     
    BP_home, BC_home = 0, 0
    
    BP_home = list(df_raw_home_1.loc[df_raw_home_1['away_home'] == "home"]['BP'])
    BC_home = list(df_raw_home_1.loc[df_raw_home_1['away_home'] == "home"]['BC'])

    delta_home = list(np.array(BP_home) - np.array(BC_home))

    delta_max_home = delta_home.index(int(max(np.array(delta_home))))
    delta_min_home = delta_home.index(int(min(np.array(delta_home))))

    BP_max_home, BC_max_home, BP_min_home, BC_min_home = BP_home[delta_max_home], BC_home[delta_max_home], BP_home[delta_min_home], BC_home[delta_min_home]
    N = len(delta_home)

    delta_home.reverse()

    ax.stackplot(["N-{}".format(i) for i in range (N,0,-1)], delta_home, alpha=0.8, colors=home_color_plot,labels=homeTeam_name)
    ax.set_title(title,fontsize=title_size,weight='bold')

    delta_home_mean = [np.mean(delta_home)]*N
    ax.plot(["N-{}".format(i) for i in range (N,0,-1)],delta_home_mean, color=black_color_plot, linestyle='--')

    plt.tight_layout()
    fig.autofmt_xdate()
    #ax.set_edgecolor('white')
    #plt.figure(edgecolor='white')
    #plt.axes(frameon=False)
    plt.savefig(fname=link_reports + str(game_fixture_id) +'/'+name_folder_introduction + '/'+'introduction_reports_chart_section4'+graph_format, dpi = dpi_1)
    plt.show()

    return np.mean(delta_home), delta_max_home, BP_max_home, BC_max_home, delta_min_home, BP_min_home, BC_min_home

# OK -
def introduction_reports_chart_section5(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_introduction,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3): #NameOfReports_Reports_SectionXX

    ticks_size = 13
    label_size = ticks_size*1.15
    legend_size = ticks_size*1.15
    title_size = ticks_size*1.4
    
    fixtures_home_away = pd.read_pickle(link_data+'{}'.format(team_id)+'/'+str(game_fixture_id) + '/' +'fixtures_package'+'/'+'fixtures_team_id_next_s'+ '_'+ str(team_id) +'_'+ str(nex) +'.pkl')

    # A. Language parameters
    if english == 1 :
        title = str("Goal difference previous games outside")
        average = "Average"
    elif english == 0 :
        title = str("Différence de buts à l'extérieur")
        average = "Moyenne"
    
    fig, ax = plt.subplots(1,1)
    
    # B. Plot Graph 
    BP_away, BC_away = 0, 0

    BP_away = list(df_raw_away_1.loc[df_raw_away_1['away_home'] == "away"]['BP'])
    BC_away = list(df_raw_away_1.loc[df_raw_away_1['away_home'] == "away"]['BC'])

    delta_away = list(np.array(BP_away) - np.array(BC_away))

    delta_max_away = delta_away.index(int(max(np.array(delta_away))))
    delta_min_away = delta_away.index(int(min(np.array(delta_away))))

    BP_max_away, BC_max_away ,BP_min_away ,BC_min_away = BP_away[delta_max_away], BC_away[delta_max_away], BP_away[delta_min_away], BC_away[delta_min_away]
    N = len(delta_away)

    delta_away.reverse()

    ax.stackplot(["N-{}".format(i) for i in range (N,0,-1)], delta_away, alpha=0.8, colors=away_color_plot,labels=awayTeam_name)
    ax.set_title(title,fontsize=title_size,weight='bold')

    delta_away_mean = [np.mean(delta_away)]*N
    ax.plot(["N-{}".format(i) for i in range (N,0,-1)],delta_away_mean, color=black_color_plot, linestyle='--')

    plt.tight_layout()
    fig.autofmt_xdate()
    #ax.set_edgecolor('white')
    #plt.figure(edgecolor='white')
    #plt.axes(frameon=False)
    plt.savefig(fname=link_reports + str(game_fixture_id) +'/'+name_folder_introduction + '/'+'introduction_reports_chart_section5'+graph_format, dpi = dpi_1)
    plt.show()

    return np.mean(delta_away), delta_max_away, BP_max_away, BC_max_away, delta_min_away, BP_min_away, BC_min_away

'''
 Functions 
------------------------------------------------------------------------------------
'''

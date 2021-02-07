# Imports
import requests
import json
import pandas as pd  
import numpy as np
import matplotlib as mpl
import matplotlib.pyplot as plt
from matplotlib import rcParams
from matplotlib.ticker import StrMethodFormatter
from pathlib import Path
import time
from matplotlib.lines import Line2D
import matplotlib.patches as mpatches
import plotly.graph_objects as go
from pptx.util import Inches
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
import shutil  
from shutil import copyfile
#from pptx.enum.shapes import MSO_SHAPE
#from pptx.enum.dml import MSO_FILL
from colour import Color
import webcolors
import math

from PIL import ImageColor

from pptx.util import Inches, Pt,Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import ColorFormat, RGBColor
import codecs

from parameters_package import *
from core_functions import * # For nan_to_zero function
from output_v2 import *
from summarize import *

import seaborn as sns
import warnings; warnings.filterwarnings(action='once')


'''
Package parameters
------------------------------------------------------------------------------------
'''

# The parameters for the font
mpl.rcParams["font.family"] = font_1

'''
Transfer functions
------------------------------------------------------------------------------------
'''

global team_id 
global opp_id  
global team_name 
global opp_name  
global game_fixture_id  

# OK - 
def pressure_master(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_presuure,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3):

    print("start of pressure_master")
    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    
    prs = Presentation(pres_path)

    # Charts
    
    # Other photos and docs to be added

    # Texts
    pressure_adv = transfer_text_to_ppt_pressure(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_pressure,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3)
    prs = Presentation(pres_path)
    
    # Advantage
    move_box_advantage(pres_path,"L'analyse de la réaction des équipes à la pression et l'enjeu donne l'avantage à",str(pressure_adv),matrix_positions_pressure['pressure_advantage']['pos_left']+2*buffer_adv,matrix_positions_pressure['pressure_advantage']['pos_top']+2*buffer_adv)
    
    print("end of pressure_master")

def transfer_text_to_ppt_pressure(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_pressure,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3):
    
    # Advantage parameters
    advantage_pressure_home = 0
    advantage_pressure_away = 0

    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    

    # Useful data part 1
    fixtures_home_away = pd.read_pickle(link_data+'{}'.format(team_id)+'/'+str(game_fixture_id) + '/' +'fixtures_package'+'/'+'fixtures_team_id_next_s'+ '_'+ str(team_id) +'_'+ str(nex) +'.pkl')
    rounds = fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('round')]
    
    league_data = pd.read_pickle(link_data+ central_folder +'/'+link_data_package_leagues_central+'/'+'leagues_league_id_master_s'+ '_'+ str(fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')]) +'.pkl')
    teams_league_data = pd.read_pickle(link_data+ central_folder +'/'+link_data_package_teams_central+'/'+'teams_league_id_s'+ '_'+ str(fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')]) +'.pkl')
    homeTeam_data = teams_league_data.loc[teams_league_data['team_id'] == int(homeTeam_id)]
    
    # Header
    text_modification_paragraph_center(pres_path,"STAKES OF THE GAME","ENJEU ET PRESSION DU MATCH",font_size_titleh1,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"Observing teams behaviour under similar level of pressure and stakes","Analyse des équipes face à la pression et l'enjeu du match",14,font_1,black_color_report,None,True)
    text_modification_paragraph_center(pres_path,"Future games","Prochains match",14,font_1,black_color_report,None,True)
    text_modification_paragraph_center(pres_path,"Observing direct competitors in the league","Analyse des concurrents direct dans la ligue",14,font_1,black_color_report,None,True)
    text_modification_paragraph_center(pres_path,"XXX","XXX",14,font_1,black_color_report,None,True)

    # Titles
    text_modification_paragraph_center(pres_path,"FUTURE GAMES OF BOTH TEAMS","PROCHAINS MATCH DES EQUIPES",font_size_titleh2,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"DIRECT COMPETITORS IN THE LEAGUE","ADVERSAIRES DIRECTS DANS LA LIGUE",font_size_titleh2,font_1,black_color_report,True,None)
    
    # Beginning of populating of the section
    text_modification_paragraph_center(pres_path,"Home team1",homeTeam_name,14,font_1,home_color_report,True,None)
    text_modification_paragraph_center(pres_path,"Away team1",awayTeam_name,14,font_1,home_color_report,True,None)

    away_home_list = ["h","a"]
    upre1, upre2, upre3, upre4, upre5, upre6, nupre2, nupre5 = 0,0,0,0,0,0,0,0
    for away_home in away_home_list:
        for i in range (1,5): # The number of games we are looking at
            if away_home == "h":
                matrix = df_raw_home_1
            else:
                matrix = df_raw_away_1
            
            n1 = matrix.iloc[0,matrix.columns.get_loc("N+{}_team".format(i+1))]
            
            if type(n1[0]) != type({}):
                text_modification_table_left(pres_path,i,1,"{}namen{}".format(away_home,i),str(n1[0]),font_size_table,font_1,black_color_report,None,None)
                text_modification_table_left(pres_path,i,2,"{}leaguen{}".format(away_home,i),str(n1[-2]),font_size_table,font_1,black_color_report,None,None)
                text_modification_table_left(pres_path,i,3,"{}rn{}".format(away_home,i),str(n1[1]),font_size_table,font_1,black_color_report,None,None)
                text_modification_table_left(pres_path,i,4,"{}pn{}".format(away_home,i),str(n1[2]),font_size_table,font_1,black_color_report,None,None)
                text_modification_table_left(pres_path,i,5,"{}wn{}".format(away_home,i),str(n1[4]),font_size_table,font_1,black_color_report,None,None)
                text_modification_table_left(pres_path,i,6,"{}dn{}".format(away_home,i),str(n1[5]),font_size_table,font_1,black_color_report,None,None)
                text_modification_table_left(pres_path,i,7,"{}ln{}".format(away_home,i),str(n1[6]),font_size_table,font_1,black_color_report,None,None)
                
                if away_home == "h":
                    upre2 += int(n1[1])
                    nupre2 += 1
                else:
                    upre5 += int(n1[1])
                    nupre5 += 1
                if int(n1[1])>int(matrix.iloc[0,matrix.columns.get_loc("rank_team")][2]):
                    if away_home == "h":
                        upre3 += 1
                    else:
                        upre6 += 1
                if abs(int(n1[2])-int(matrix.iloc[0,matrix.columns.get_loc("rank_team")][2]))<6 : # The points delta to deem a game as at high stake
                    text_modification_table_left(pres_path,i,8,"{}staken{}".format(away_home,i),"Fort enjeu",11,font_1,red_color_report,True,None)
                    if away_home == "h":
                        upre1 += 1
                    else:
                        upre4 += 1
                else:
                    text_modification_table_left(pres_path,i,8,"{}staken{}".format(away_home,i),"Faible enjeu",11,font_1,black_color_report,None,None)


                prs = Presentation(pres_path)
                page_number, pos_left, pos_top = matrix_positions_pressure['logo_team_{}_{}'.format(away_home,i)]['page_number'], matrix_positions_pressure['logo_team_{}_{}'.format(away_home,i)]['pos_left'],matrix_positions_pressure['logo_team_{}_{}'.format(away_home,i)]['pos_top'] 
                try:
                    pic = prs.slides[page_number].shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(n1[-1]) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_3,prs))                    
                except:
                    logo_team_id_s(int(n1[-1]))
                    pic = prs.slides[page_number].shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(n1[-1]) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_3,prs))
                prs.save(pres_path)

            else:
                text_modification_table_left(pres_path,i,1,"{}namen{}".format(away_home,i),str(n1[0]['team_name']),11,font_1,black_color_report,None,None)
                text_modification_table_left(pres_path,i,2,"{}leaguen{}".format(away_home,i),str(n1[0]['league_name']),11,font_1,black_color_report,None,None)
                text_modification_table_left(pres_path,i,3,"{}rn{}".format(away_home,i),"-",11,font_1,black_color_report,None,None)
                text_modification_table_left(pres_path,i,4,"{}pn{}".format(away_home,i),"-",11,font_1,black_color_report,None,None)
                text_modification_table_left(pres_path,i,5,"{}wn{}".format(away_home,i),"-",11,font_1,black_color_report,None,None)
                text_modification_table_left(pres_path,i,6,"{}dn{}".format(away_home,i),"-",11,font_1,black_color_report,None,None)
                text_modification_table_left(pres_path,i,7,"{}ln{}".format(away_home,i),"-",11,font_1,black_color_report,None,None)
                text_modification_table_left(pres_path,i,8,"{}staken{}".format(away_home,i),"Autre ligue",11,font_1,black_color_report,True,None)
    
                prs = Presentation(pres_path)
                
                page_number, pos_left, pos_top = matrix_positions_pressure['logo_team_{}_{}'.format(away_home,i)]['page_number'], matrix_positions_pressure['logo_team_{}_{}'.format(away_home,i)]['pos_left'],matrix_positions_pressure['logo_team_{}_{}'.format(away_home,i)]['pos_top'] 
                try:
                    pic = prs.slides[page_number].shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(n1[0]['team_id']) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_3,prs))
                except:
                    logo_team_id_s(int(n1[0]['team_id']))
                    pic = prs.slides[page_number].shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(n1[0]['team_id']) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_3,prs))
                prs.save(pres_path)

            print(away_home,i)

    draw_home_a,lose_home_a,win_home_a,draw_away_a,lose_away_a,win_away_a,draw_home_b,lose_home_b,win_home_b,draw_away_b,lose_away_b,win_away_b,draw_home_c,lose_home_c,win_home_c,draw_away_c,lose_away_c,win_away_c = pressure_reports_section1(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_pressure,upre1,upre2/nupre2,upre3,upre4,upre5/nupre5,upre6,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3)

    number_id = 1
    text_modification_paragraph_center(pres_path,"upre{}".format(number_id),upre1,20,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"mpre{}".format(number_id),"matchs à fort enjeu",font_size_bubble,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"lpre{}".format(number_id),"pour "+ homeTeam_name,font_size_bubble,font_1,home_color_report,True,None)

    number_id = 2
    text_modification_paragraph_center(pres_path,"upre{}".format(number_id),str(round(upre2/nupre2,1)),20,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"mpre{}".format(number_id),"classement moyen",font_size_bubble,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"lpre{}".format(number_id),"adversaires",font_size_bubble,font_1,black_color_report,True,None)

    number_id = 3
    text_modification_paragraph_center(pres_path,"upre{}".format(number_id),upre3,20,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"mpre{}".format(number_id),"matchs vs.",font_size_bubble,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"lpre{}".format(number_id),"meilleurs adversaires",font_size_bubble,font_1,black_color_report,True,None)

    number_id = 4
    text_modification_paragraph_center(pres_path,"upre{}".format(number_id),upre4,20,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"mpre{}".format(number_id),"matchs à fort enjeu",font_size_bubble,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"lpre{}".format(number_id),"pour "+ awayTeam_name,font_size_bubble,font_1,away_color_report,True,None)

    number_id = 5
    text_modification_paragraph_center(pres_path,"upre{}".format(number_id),str(round(upre5/nupre5,1)),20,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"mpre{}".format(number_id),"classement moyen",font_size_bubble,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"lpre{}".format(number_id),"adversaires",font_size_bubble,font_1,black_color_report,True,None)

    number_id = 6
    text_modification_paragraph_center(pres_path,"upre{}".format(number_id),upre6,20,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"mpre{}".format(number_id),"matchs vs.",font_size_bubble,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"lpre{}".format(number_id),"meilleurs adversaires",font_size_bubble,font_1,black_color_report,True,None)
    
    number_id = 7
    text_modification_paragraph_center(pres_path,"upre{}".format(number_id),win_home_a,20,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"mpre{}".format(number_id),"victoires pour matchs",font_size_bubble,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"lpre{}".format(number_id),"avec même enjeu",font_size_bubble,font_1,black_color_report,True,None)

    number_id = 8
    text_modification_paragraph_center(pres_path,"upre{}".format(number_id),lose_away_a,20,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"mpre{}".format(number_id),"défaites pour matchs",font_size_bubble,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"lpre{}".format(number_id),"avec même enjeu",font_size_bubble,font_1,black_color_report,True,None)

    text_modification_paragraph_left(pres_path,"Reaction to pressure and high stake gives advantage to","L'analyse de la réaction des équipes à la pression et l'enjeu donne l'avantage à",14,font_1,white_color_report,True,None)

    if win_home_a != win_away_a:
        if win_home_a>win_away_a:
            text_modification_paragraph_left(pres_path,"With high stake games",homeTeam_name+" a de meilleurs résultats avec le même niveau d'enjeu de la rencontre",16,font_1,black_color_report,True,None)
        else:
            text_modification_paragraph_left(pres_path,"With high stake games",awayTeam_name+" a de meilleurs résultats avec le même niveau d'enjeu de la rencontre",16,font_1,black_color_report,True,None)
    else:
        text_modification_paragraph_left(pres_path,"With high stake games","Mêmes résultats lorsque calendrier de matchs avec enjeu similaire",16,font_1,black_color_report,True,None)


    if win_home_c != win_away_c:
        if win_home_c>win_away_c:
            text_modification_paragraph_left(pres_path,"Against better teams",homeTeam_name+" a de meilleurs résultats lorsqu'elle se prépare à afronter des adversaires de niveau similaire",16,font_1,black_color_report,True,None)
        else:
            text_modification_paragraph_left(pres_path,"Against better teams",awayTeam_name+" a de meilleurs résultats lorsqu'elle se prépare à afronter des adversaires de niveau similaire",16,font_1,black_color_report,True,None)
    else:
        text_modification_paragraph_left(pres_path,"Against better teams","Mêmes résultats futurs adversaires de niveau similaire",16,font_1,black_color_report,True,None)


    # Advantage
    text_modification_paragraph_left(pres_path,"Reaction to pressure and high stake gives advantage to","L'analyse de la réaction des équipes à la pression et l'enjeu donne l'avantage à",14,font_1,white_color_report,True,None)

    prs = Presentation(pres_path)
    prs.save(pres_path)

    pressure_adv = 0
    
    if advantage_pressure_home>advantage_pressure_away:
        text_modification_paragraph_center(pres_path,"+ ADVANTAGE PRESSURE","+ " + str(homeTeam_name),24,font_1,white_color_report,True,None)  
        pressure_adv = math.ceil(10*advantage_pressure_home/(advantage_pressure_home+advantage_pressure_away))
        
    elif advantage_pressure_home<advantage_pressure_away:
        text_modification_paragraph_center(pres_path,"+ ADVANTAGE PRESSURE","+ " + str(awayTeam_name),24,font_1,white_color_report,True,None)  
        pressure_adv = math.ceil(10*advantage_pressure_away/(advantage_pressure_home+advantage_pressure_away))
        
    else:
        text_modification_paragraph_center(pres_path,"+ ADVANTAGE PRESSURE","Perfect equality",24,font_1,white_color_report,True,None)  

    return pressure_adv


def pressure_reports_section1(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_pressure,upre1,upre2,upre3,upre4,upre5,upre6,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3):
    
    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)
    
    ticks_size = 13
    label_size = ticks_size*1.15
    legend_size = ticks_size*1.15
    title_size = ticks_size*1.4
    
    draw_home_a,lose_home_a,win_home_a,draw_away_a,lose_away_a,win_away_a = 0,0,0,0,0,0 # This is for the criteria : fort enjeu 
    draw_home_b,lose_home_b,win_home_b,draw_away_b,lose_away_b,win_away_b = 0,0,0,0,0,0 # This is for the criteria : classement moyen
    draw_home_c,lose_home_c,win_home_c,draw_away_c,lose_away_c,win_away_c = 0,0,0,0,0,0 # This is for the criteria : better oppponent 
    
    for i in range (len(df_raw_home_1)-1):
        c_data_home, c_data_away = 0,0
        a_data_home, b_data_home, b_data_home_n, a_data_away, b_data_away, b_data_away_n = 0,0,0,0,0,0
        n1_home,n2_home,n3_home,n4_home,n1_away,n2_away,n3_away,n4_away = df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("N+1_team")],df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("N+2_team")],df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("N+3_team")],df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("N+4_team")],\
            df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("N+1_team")],df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("N+2_team")],df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("N+3_team")],df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("N+4_team")]

        try:
            if type(n1_home[0]) != type({}):
                if abs(df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("N+1_team")][2]-df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("rank_team")][2])<6:
                    a_data_home += 1
                b_data_home += df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("N+1_team")][1]
                b_data_home_n += 1
                if df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("N+1_team")][1]>df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("rank_team")][1]:
                    c_data_home += 1

            if type(n2_home[0]) != type({}):
                if abs(df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("N+2_team")][2]-df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("rank_team")][2])<6:
                    a_data_home += 1
                b_data_home += df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("N+2_team")][1]
                b_data_home_n += 1
                if df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("N+2_team")][1]>df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("rank_team")][1]:
                    c_data_home += 1

            if type(n3_home[0]) != type({}):
                if abs(df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("N+3_team")][2]-df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("rank_team")][2])<6:
                    a_data_home += 1
                b_data_home += df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("N+3_team")][1]
                b_data_home_n += 1
                if df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("N+3_team")][1]>df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("rank_team")][1]:
                    c_data_home += 1

            if type(n4_home[0]) != type({}):
                if abs(df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("N+4_team")][2]-df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("rank_team")][2])<6:
                    a_data_home += 1
                b_data_home += df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("N+4_team")][1]
                b_data_home_n += 1
                if df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("N+4_team")][1]>df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("rank_team")][1]:
                    c_data_home += 1

            if type(n1_away[0]) != type({}):
                if abs(df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("N+1_team")][2]-df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("rank_team")][2])<6:
                    a_data_away += 1
                b_data_away += df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("N+1_team")][1]
                b_data_away_n += 1
                if df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("N+1_team")][1]>df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("rank_team")][1]:
                    c_data_away += 1
            
            if type(n2_away[0]) != type({}):
                if abs(df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("N+2_team")][2]-df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("rank_team")][2])<6:
                    a_data_away += 1
                b_data_away += df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("N+2_team")][1]
                b_data_away_n += 1
                if df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("N+2_team")][1]>df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("rank_team")][1]:
                    c_data_away += 1
            
            if type(n3_away[0]) != type({}):
                if abs(df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("N+3_team")][2]-df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("rank_team")][2])<6:
                    a_data_away += 1
                b_data_away += df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("N+3_team")][1]
                b_data_away_n += 1
                if df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("N+3_team")][1]>df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("rank_team")][1]:
                    c_data_away += 1
            
            if type(n4_away[0]) != type({}):
                if abs(df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("N+4_team")][2]-df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("rank_team")][2])<6:
                    a_data_away += 1
                b_data_away += df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("N+4_team")][1]
                b_data_away_n += 1
                if df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("N+4_team")][1]>df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("rank_team")][1]:
                    c_data_away += 1

            if a_data_home == upre1 :
                if df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("W-D-L")] == "W":
                    win_home_a += 1
                elif df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("W-D-L")] == "D":
                    draw_home_a += 1
                else:
                    lose_home_a += 1
            if abs((b_data_home/b_data_home_n) - upre2)<3: # 3 is the limit for considering close enough to be similar
                if df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("W-D-L")] == "W":
                    win_home_b += 1
                elif df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("W-D-L")] == "D":
                    draw_home_b += 1
                else:
                    lose_home_b += 1
            if a_data_away == upre4 :
                if df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("W-D-L")] == "W":
                    win_away_a += 1
                elif df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("W-D-L")] == "D":
                    draw_away_a += 1
                else:
                    lose_away_a += 1
            if abs((b_data_away/b_data_away_n) - upre5)<3: # 3 is the limit for considering close enough to be similar
                if df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("W-D-L")] == "W":
                    win_away_b += 1
                elif df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("W-D-L")] == "D":
                    draw_away_b += 1
                else:
                    lose_away_b += 1

            if c_data_home == upre3:
                if df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("W-D-L")] == "W":
                    win_home_a += 1
                elif df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("W-D-L")] == "D":
                    draw_home_a += 1
                else:
                    lose_home_a += 1
            if c_data_away == upre6:
                if df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("W-D-L")] == "W":
                    win_away_b += 1
                elif df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("W-D-L")] == "D":
                    draw_away_b += 1
                else:
                    lose_away_b += 1

            print("home")
            print(a_data_home,b_data_home,c_data_home)
            print(upre1,upre2,upre3)

            print("away")
            print(a_data_away,b_data_away,c_data_away)
            print(upre4,upre5,upre6)

        except Exception as e:
            print(e)

    # Graph number 1
    if sum([draw_home_a,lose_home_a,win_home_a]) != 0:
        if english == 1 :
            title = str("Last 20 games results with similar stake")
            labels = ["Draw","Lose","Win"] # Keep the same order
        elif english == 0 :
            title = str("Résultats 20 derniers matchs lors du même enjeu")
            labels = ["Nul","Défaite","Victoire"] # Keep the same order

        fig, ax = plt.subplots(1,1)
        ind = np.arange(len(labels))
        width = 0.30
        ax.bar(ind, [draw_home_a,lose_home_a,win_home_a], width, label=labels, color=[draw_color_plot,lose_color_plot,win_color_plot],edgecolor='black')

        ax.set_title(homeTeam_name,weight='bold',fontsize=title_size,color=home_color_plot)

        fig.suptitle(title, fontsize=16)
        fig.subplots_adjust(top=0.8)

        for p in ax.patches:
            width = p.get_width()
            height = p.get_height()
            x, y = p.get_xy() 
            ax.annotate(f'{height/(draw_home_a+lose_home_a+win_home_a):.0%}', (x + width/2, y + height*1.15), ha='center',fontsize=13)

        ax.set_ylim([0,max([draw_home_a,lose_home_a,win_home_a])+3])
        ax.set_yticks([], minor=True)
        ax.set_yticks([])
        
        plt.xticks(ind , labels,fontsize=ticks_size-2)
        plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_pressure+'/'+'pressure_reports_chart_section1'+graph_format, dpi = dpi_1)
        plt.show()
        
        page_number, pos_left, pos_top = matrix_positions_pressure['pressure_reports_chart_section1']['page_number'], matrix_positions_pressure['pressure_reports_chart_section1']['pos_left'],matrix_positions_pressure['pressure_reports_chart_section1']['pos_top'] 
        slide = prs.slides[page_number]
        pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_pressure+'/'+'pressure_reports_chart_section1'+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_6,prs))
    
    else:
        if english == 1 :
            title = str("Last 20 games results when future games against similar opponents")
            labels = ["Draw","Lose","Win"] # Keep the same order
        elif english == 0 :
            title = str("Résultats 20 derniers matchs quand futurs matchs contre meme adversaires")
            labels = ["Nul","Défaite","Victoire"] # Keep the same order

        fig, ax = plt.subplots(1,1)
        ind = np.arange(len(labels))
        width = 0.30
        ax.bar(ind, [draw_home_b,lose_home_b,win_home_b], width, label=labels, color=[draw_color_plot,lose_color_plot,win_color_plot],edgecolor='black')

        ax.set_title(homeTeam_name,weight='bold',fontsize=title_size,color=home_color_plot)

        fig.suptitle(title, fontsize=16)
        fig.subplots_adjust(top=0.8)

        if sum([draw_home_b,lose_home_b,win_home_b]) != 0:
            for p in ax.patches:
                width = p.get_width()
                height = p.get_height()
                x, y = p.get_xy() 
                ax.annotate(f'{height/(draw_home_b+lose_home_b+win_home_b):.0%}', (x + width/2, y + height*1.15), ha='center',fontsize=13)
        else:
            for p in ax.patches:
                width = p.get_width()
                height = p.get_height()
                x, y = p.get_xy() 
                ax.annotate(str(0), (x + width/2, y + height*1.15), ha='center',fontsize=13)

        ax.set_ylim([0,max([draw_home_b,lose_home_b,win_home_b])+3])
        ax.set_yticks([], minor=True)
        ax.set_yticks([])
        
        plt.xticks(ind , labels,fontsize=ticks_size-2)
        plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_pressure+'/'+'pressure_reports_chart_section1'+graph_format, dpi = dpi_1)
        plt.show()
        
        page_number, pos_left, pos_top = matrix_positions_pressure['pressure_reports_chart_section1']['page_number'], matrix_positions_pressure['pressure_reports_chart_section1']['pos_left'],matrix_positions_pressure['pressure_reports_chart_section1']['pos_top'] 
        slide = prs.slides[page_number]
        pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_pressure+'/'+'pressure_reports_chart_section1'+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_6,prs))
        
    # Graph number 2
    if sum([draw_home_a,lose_home_a,win_home_a]) != 0 and sum([draw_home_b,lose_home_b,win_home_b]):
        if english == 1 :
            title = str("Last 20 games results when future games against similar opponents")
            labels = ["Draw","Lose","Win"] # Keep the same order
        elif english == 0 :
            title = str("Résultats 20 derniers matchs quand futurs matchs contre meme adversaires")
            labels = ["Nul","Défaite","Victoire"] # Keep the same order

        fig, ax = plt.subplots(1,1)
        ind = np.arange(len(labels))
        width = 0.30
        ax.bar(ind, [draw_home_b,lose_home_b,win_home_b], width, label=labels, color=[draw_color_plot,lose_color_plot,win_color_plot],edgecolor='black')

        ax.set_title(homeTeam_name,weight='bold',fontsize=title_size,color=home_color_plot)

        fig.suptitle(title, fontsize=16)
        fig.subplots_adjust(top=0.8)

        for p in ax.patches:
            width = p.get_width()
            height = p.get_height()
            x, y = p.get_xy() 
            ax.annotate(f'{height/(draw_home_b+lose_home_b+win_home_b):.0%}', (x + width/2, y + height*1.15), ha='center',fontsize=13)

        ax.set_ylim([0,max([draw_home_b,lose_home_b,win_home_b])+3])
        ax.set_yticks([], minor=True)
        ax.set_yticks([])
        
        plt.xticks(ind , labels,fontsize=ticks_size-2)
        plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_pressure+'/'+'pressure_reports_chart_section2'+graph_format, dpi = dpi_1)
        plt.show()
        
        page_number, pos_left, pos_top = matrix_positions_pressure['pressure_reports_chart_section2']['page_number'], matrix_positions_pressure['pressure_reports_chart_section2']['pos_left'],matrix_positions_pressure['pressure_reports_chart_section2']['pos_top'] 
        slide = prs.slides[page_number]
        pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_pressure+'/'+'pressure_reports_chart_section2'+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_6,prs))
        
    else:
        if english == 1 :
            title = str("Last 20 games results when having future games vs. better opponents")
            labels = ["Draw","Lose","Win"] # Keep the same order
            
        elif english == 0 :
            title = str("Résultats 20 derniers matchs quand futures match vs. meilleurs adversaires")
            labels = ["Nul","Défaite","Victoire"] # Keep the same order

        fig, ax = plt.subplots(1,1)
        ind = np.arange(len(labels))
        width = 0.30
        ax.bar(ind, [draw_home_c,lose_home_c,win_home_c], width, label=labels, color=[draw_color_plot,lose_color_plot,win_color_plot],edgecolor='black')

        ax.set_title(homeTeam_name,weight='bold',fontsize=title_size,color=home_color_plot)

        fig.suptitle(title, fontsize=16)
        fig.subplots_adjust(top=0.8)

        if sum([draw_home_c,lose_home_c,win_home_c]) != 0:
            for p in ax.patches:
                width = p.get_width()
                height = p.get_height()
                x, y = p.get_xy() 
                ax.annotate(f'{height/(draw_home_c+lose_home_c+win_home_c):.0%}', (x + width/2, y + height*1.15), ha='center',fontsize=13)

        else:
            for p in ax.patches:
                width = p.get_width()
                height = p.get_height()
                x, y = p.get_xy() 
                ax.annotate(str(0), (x + width/2, y + height*1.15), ha='center',fontsize=13)
                
        ax.set_ylim([0,max([draw_home_c,lose_home_c,win_home_c])+3])
        ax.set_yticks([], minor=True)
        ax.set_yticks([])
        
        plt.xticks(ind , labels,fontsize=ticks_size-2)
        plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_pressure+'/'+'pressure_reports_chart_section2'+graph_format, dpi = dpi_1)
        plt.show()
        
        page_number, pos_left, pos_top = matrix_positions_pressure['pressure_reports_chart_section2']['page_number'], matrix_positions_pressure['pressure_reports_chart_section2']['pos_left'],matrix_positions_pressure['pressure_reports_chart_section2']['pos_top'] 
        slide = prs.slides[page_number]
        pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_pressure+'/'+'pressure_reports_chart_section2'+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_6,prs))



    # Graph number 3
    if sum([draw_away_a,lose_away_a,win_away_a]) != 0:
        if english == 1 :
            title = str("Last 20 games results with similar stake")
            labels = ["Draw","Lose","Win"] # Keep the same order
        elif english == 0 :
            title = str("Résultats 20 derniers matchs lors du même enjeu")
            labels = ["Nul","Défaite","Victoire"] # Keep the same order

        fig, ax = plt.subplots(1,1)
        ind = np.arange(len(labels))
        width = 0.30
        ax.bar(ind, [draw_away_a,lose_away_a,win_away_a], width, label=labels, color=[draw_color_plot,lose_color_plot,win_color_plot],edgecolor='black')

        ax.set_title(awayTeam_name,weight='bold',fontsize=title_size,color=away_color_plot)

        fig.suptitle(title, fontsize=16)
        fig.subplots_adjust(top=0.8)

        for p in ax.patches:
            width = p.get_width()
            height = p.get_height()
            x, y = p.get_xy() 
            ax.annotate(f'{height/(draw_away_a+lose_away_a+win_away_a):.0%}', (x + width/2, y + height*1.15), ha='center',fontsize=13)

        ax.set_ylim([0,max([draw_away_a,lose_away_a,win_away_a])+3])
        ax.set_yticks([], minor=True)
        ax.set_yticks([])
        
        plt.xticks(ind , labels,fontsize=ticks_size-2)
        plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_pressure+'/'+'pressure_reports_chart_section3'+graph_format, dpi = dpi_1)
        plt.show()
        
        page_number, pos_left, pos_top = matrix_positions_pressure['pressure_reports_chart_section3']['page_number'], matrix_positions_pressure['pressure_reports_chart_section3']['pos_left'],matrix_positions_pressure['pressure_reports_chart_section3']['pos_top'] 
        slide = prs.slides[page_number]
        pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_pressure+'/'+'pressure_reports_chart_section3'+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_6,prs))
    
    else:
        if english == 1 :
            title = str("Last 20 games results when future games against similar opponents")
            labels = ["Draw","Lose","Win"] # Keep the same order
        elif english == 0 :
            title = str("Résultats 20 derniers matchs quand futurs matchs contre meme adversaires")
            labels = ["Nul","Défaite","Victoire"] # Keep the same order

        fig, ax = plt.subplots(1,1)
        ind = np.arange(len(labels))
        width = 0.30
        ax.bar(ind, [draw_away_b,lose_away_b,win_away_b], width, label=labels, color=[draw_color_plot,lose_color_plot,win_color_plot],edgecolor='black')

        ax.set_title(awayTeam_name,weight='bold',fontsize=title_size,color=away_color_plot)

        fig.suptitle(title, fontsize=16)
        fig.subplots_adjust(top=0.8)

        if sum([draw_away_b,lose_away_b,win_away_b]) != 0:
            for p in ax.patches:
                width = p.get_width()
                height = p.get_height()
                x, y = p.get_xy() 
                ax.annotate(f'{height/(draw_away_b+lose_away_b+win_away_b):.0%}', (x + width/2, y + height*1.15), ha='center',fontsize=13)
        else:
            for p in ax.patches:
                width = p.get_width()
                height = p.get_height()
                x, y = p.get_xy() 
                ax.annotate(str(0), (x + width/2, y + height*1.15), ha='center',fontsize=13)

        ax.set_ylim([0,max([draw_away_b,lose_away_b,win_away_b])+3])
        ax.set_yticks([], minor=True)
        ax.set_yticks([])
        
        plt.xticks(ind , labels,fontsize=ticks_size-2)
        plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_pressure+'/'+'pressure_reports_chart_section3'+graph_format, dpi = dpi_1)
        plt.show()
        
        page_number, pos_left, pos_top = matrix_positions_pressure['pressure_reports_chart_section3']['page_number'], matrix_positions_pressure['pressure_reports_chart_section3']['pos_left'],matrix_positions_pressure['pressure_reports_chart_section3']['pos_top'] 
        slide = prs.slides[page_number]
        pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_pressure+'/'+'pressure_reports_chart_section3'+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_6,prs))
        
    # Graph number 4
    if sum([draw_away_a,lose_away_a,win_away_a]) != 0 and sum([draw_away_b,lose_away_b,win_away_b]):
        if english == 1 :
            title = str("Last 20 games results when future games against similar opponents")
            labels = ["Draw","Lose","Win"] # Keep the same order
        elif english == 0 :
            title = str("Résultats 20 derniers matchs quand futurs matchs contre meme adversaires")
            labels = ["Nul","Défaite","Victoire"] # Keep the same order

        fig, ax = plt.subplots(1,1)
        ind = np.arange(len(labels))
        width = 0.30
        ax.bar(ind, [draw_away_b,lose_away_b,win_away_b], width, label=labels, color=[draw_color_plot,lose_color_plot,win_color_plot],edgecolor='black')

        ax.set_title(awayTeam_name,weight='bold',fontsize=title_size,color=away_color_plot)

        fig.suptitle(title, fontsize=16)
        fig.subplots_adjust(top=0.8)

        for p in ax.patches:
            width = p.get_width()
            height = p.get_height()
            x, y = p.get_xy() 
            ax.annotate(f'{height/(draw_away_b+lose_away_b+win_away_b):.0%}', (x + width/2, y + height*1.15), ha='center',fontsize=13)

        ax.set_ylim([0,max([draw_away_b,lose_away_b,win_away_b])+3])
        ax.set_yticks([], minor=True)
        ax.set_yticks([])
        
        plt.xticks(ind , labels,fontsize=ticks_size-2)
        plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_pressure+'/'+'pressure_reports_chart_section4'+graph_format, dpi = dpi_1)
        plt.show()
        
        page_number, pos_left, pos_top = matrix_positions_pressure['pressure_reports_chart_section4']['page_number'], matrix_positions_pressure['pressure_reports_chart_section4']['pos_left'],matrix_positions_pressure['pressure_reports_chart_section4']['pos_top'] 
        slide = prs.slides[page_number]
        pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_pressure+'/'+'pressure_reports_chart_section4'+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_6,prs))
        
    else:
        if english == 1 :
            title = str("Last 20 games results when having future games vs. better opponents")
            labels = ["Draw","Lose","Win"] # Keep the same order
            
        elif english == 0 :
            title = str("Résultats 20 derniers matchs quand futures match vs. meilleurs adversaires")
            labels = ["Nul","Défaite","Victoire"] # Keep the same order

        fig, ax = plt.subplots(1,1)
        ind = np.arange(len(labels))
        width = 0.30
        ax.bar(ind, [draw_away_c,lose_away_c,win_away_c], width, label=labels, color=[draw_color_plot,lose_color_plot,win_color_plot],edgecolor='black')

        ax.set_title(awayTeam_name,weight='bold',fontsize=title_size,color=away_color_plot)

        fig.suptitle(title, fontsize=16)
        fig.subplots_adjust(top=0.8)

        if sum([draw_away_c,lose_away_c,win_away_c]) != 0:
            for p in ax.patches:
                width = p.get_width()
                height = p.get_height()
                x, y = p.get_xy() 
                ax.annotate(f'{height/(draw_away_c+lose_away_c+win_away_c):.0%}', (x + width/2, y + height*1.15), ha='center',fontsize=13)

        else:
            for p in ax.patches:
                width = p.get_width()
                height = p.get_height()
                x, y = p.get_xy() 
                ax.annotate(str(0), (x + width/2, y + height*1.15), ha='center',fontsize=13)
                
        ax.set_ylim([0,max([draw_away_c,lose_away_c,win_away_c])+3])
        ax.set_yticks([], minor=True)
        ax.set_yticks([])
        
        plt.xticks(ind , labels,fontsize=ticks_size-2)
        plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_pressure+'/'+'pressure_reports_chart_section4'+graph_format, dpi = dpi_1)
        plt.show()
        
        page_number, pos_left, pos_top = matrix_positions_pressure['pressure_reports_chart_section4']['page_number'], matrix_positions_pressure['pressure_reports_chart_section4']['pos_left'],matrix_positions_pressure['pressure_reports_chart_section4']['pos_top'] 
        slide = prs.slides[page_number]
        pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_pressure+'/'+'pressure_reports_chart_section4'+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_6,prs))


    """ 
    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)
      """
    prs.save(pres_path)
    
    return draw_home_a,lose_home_a,win_home_a,draw_away_a,lose_away_a,win_away_a,draw_home_b,lose_home_b,win_home_b,draw_away_b,lose_away_b,win_away_b,draw_home_c,lose_home_c,win_home_c,draw_away_c,lose_away_c,win_away_c
    

# Imports
import requests
import json
import pandas as pd  
import numpy as np
import matplotlib as mpl
import matplotlib.pyplot as plt
from matplotlib import rcParams
from matplotlib.ticker import StrMethodFormatter
from pathlib import Path
import time
from matplotlib.lines import Line2D
import matplotlib.patches as mpatches
import plotly.graph_objects as go
from pptx.util import Inches
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
import shutil  
from shutil import copyfile
#from pptx.enum.shapes import MSO_SHAPE
#from pptx.enum.dml import MSO_FILL
from colour import Color
import webcolors
import math

from PIL import ImageColor

from pptx.util import Inches, Pt,Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import ColorFormat, RGBColor
import codecs

from parameters_package import *
from core_functions import * # For nan_to_zero function
from output_v2 import *
from summarize import *

import seaborn as sns
import warnings; warnings.filterwarnings(action='once')


'''
Package parameters
------------------------------------------------------------------------------------
'''

# The parameters for the font
mpl.rcParams["font.family"] = font_1

'''
Transfer functions
------------------------------------------------------------------------------------
'''

global team_id 
global opp_id  
global team_name 
global opp_name  
global game_fixture_id  

# OK - 
def h2h_master(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_h2h,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3):

    print("start of h2h_master")
    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    
    prs = Presentation(pres_path)

    # Charts
    
    # Other photos and docs to be added

    # Texts
    h2h_adv = transfer_text_to_ppt_h2h(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_h2h,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3)
    prs = Presentation(pres_path)
    
    # Advantage
    move_box_advantage(pres_path,"L'historique des matchs entre les 2 équipes est à l'avantage de",str(h2h_adv),matrix_positions_h2h['h2h_advantage']['pos_left']+2*buffer_adv,matrix_positions_h2h['h2h_advantage']['pos_top']+2*buffer_adv)
    
    print("end of h2h_master")

def transfer_text_to_ppt_h2h(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_h2h,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3):
    
    # Advantage parameters
    advantage_h2h_home = 0
    advantage_h2h_away = 0

    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    

    # Useful data part 1
    fixtures_home_away = pd.read_pickle(link_data+'{}'.format(team_id)+'/'+str(game_fixture_id) + '/' +'fixtures_package'+'/'+'fixtures_team_id_next_s'+ '_'+ str(team_id) +'_'+ str(nex) +'.pkl')
    rounds = fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('round')]
    
    league_data = pd.read_pickle(link_data+ central_folder +'/'+link_data_package_leagues_central+'/'+'leagues_league_id_master_s'+ '_'+ str(fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')]) +'.pkl')
    teams_league_data = pd.read_pickle(link_data+ central_folder +'/'+link_data_package_teams_central+'/'+'teams_league_id_s'+ '_'+ str(fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')]) +'.pkl')
    homeTeam_data = teams_league_data.loc[teams_league_data['team_id'] == int(homeTeam_id)]
    
    # Header
    text_modification_paragraph_center(pres_path,"h2h OF PLAYING","TYPOLOGIE ET STYLES DE JEU",font_size_titleh1,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"Observing the past performances between the two teams","Analyse der dernières rencontres entre les 2 équipes",14,font_1,black_color_report,None,True)

    text_modification_paragraph_center(pres_path,"Past games results","Résultats des derniers face-à-face",14,font_1,black_color_report,None,True)
    text_modification_paragraph_center(pres_path,"Pas games analysis","Analyse des derniers face-à-face≤",14,font_1,black_color_report,None,True)


    # Titles
    text_modification_paragraph_center(pres_path,"H2H RESULTS","RESULTATS DES FACE-A-FACE",font_size_titleh2,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"H2H ANALYSIS","ANALYSE DES FACE-A-FACE",font_size_titleh2,font_1,black_color_report,True,None)

    # Beginning of populating of the section
    
    # Advantage
    text_modification_paragraph_left(pres_path,"h2h results gives advantage to","L'historique des matchs entre les 2 équipes est à l'avantage de",14,font_1,white_color_report,True,None)

    prs = Presentation(pres_path)
    prs.save(pres_path)

    h2h_adv = 0
    
    if advantage_h2h_home>advantage_h2h_away:
        text_modification_paragraph_center(pres_path,"+ ADVANTAGE h2h","+ " + str(homeTeam_name),24,font_1,white_color_report,True,None)  
        h2h_adv = math.ceil(10*advantage_h2h_home/(advantage_h2h_home+advantage_h2h_away))
        
    elif advantage_h2h_home<advantage_h2h_away:
        text_modification_paragraph_center(pres_path,"+ ADVANTAGE h2h","+ " + str(awayTeam_name),24,font_1,white_color_report,True,None)  
        h2h_adv = math.ceil(10*advantage_h2h_away/(advantage_h2h_home+advantage_h2h_away))
        
    else:
        text_modification_paragraph_center(pres_path,"+ ADVANTAGE h2h","Perfect equality",24,font_1,white_color_report,True,None)  

    return h2h_adv


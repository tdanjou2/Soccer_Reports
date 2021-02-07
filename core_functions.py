# Imports 
import requests
import json
import pandas as pd  
import numpy as np
from pathlib import Path
import time

import openpyxl 
from openpyxl import Workbook
from openpyxl.styles import PatternFill
from openpyxl.workbook import Workbook
from openpyxl.styles import Font, Fill
from openpyxl.styles import colors

from parameters_package import *

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

'''
Package parameters
------------------------------------------------------------------------------------
'''

'''
Functions
------------------------------------------------------------------------------------
'''

# OK -
def nan_to_zero (df,x):
    
    for i in df.index:
        try:
            df.iloc[i,df.columns.get_loc(x)] = int(df.iloc[i,df.columns.get_loc(x)])
        except:
            df.iloc[i,df.columns.get_loc(x)] = 0
    return df

# OK -
def nan_to_zero_all (df):
    
    for i in range (df.shape[0]):
        for j in range (df.shape[1]):
            try:
                df.iloc[i,j] = int(df.iloc[i,j])
            except:
                df.iloc[i,j] = 0
    return df

# OK -
def columns_to_excel_mod (m):

    ranks = []
    for (columnName, columnData) in m.iteritems():
        columnName = str(columnName)
        if columnName[-3:] == "mod" :
            ranks.append(m.columns.get_loc(columnName))
    return ranks

# OK -
def customize_excel(link):
    link = str(link)
    wb = openpyxl.load_workbook(link)
    sheet = wb['Sheet1']
    greyfill = PatternFill(fill_type='solid',start_color='9A99A9',end_color='9A99A9')
    bluefill = PatternFill(fill_type='solid',start_color='DCDBF2',end_color='DCDBF2')
    a = 1
    for i in range (1,last_fixtures_lineups_duration+2):
        a,b = 0,0
        for j in range (1,len(cols)+2):
            if sheet.cell(i,j).value == None:
                sheet.cell(i,j).fill = greyfill
            if j<len(cols)+1:
                if cols[j-1] == "^GAME^":
                    a=1
                if cols[j-1] == "^NOT_GAME^":
                    b=1
                if a == 1 and b == 1:
                    sheet.cell(i,j+1).fill = bluefill

    sheet.sheet_view.showGridLines = False
    wb.save(link)

# OK -
def round_formating(x):
    x = str(x)
    y = ""
    for i in range (len(x)):
        if x[i] == " ":
            y += "_"
        else:
            y += x[i]
    return y

# OK -
def match_item_list(item,list_target):
    i = 0
    while i<len(list_target) and int(item) != int(list_target[i]) :
        i += 1
    if i == len(list_target) :
        return False
    else:
         return True

# OK -
def match_item_list_general(item,list_target):
    i = 0
    while i<len(list_target) and item != list_target[i] :
        i += 1
    if i == len(list_target):
        return False
    else:
         return True

# OK -
def position(item,list_target):
    i = 0
    while i<len(list_target) and item != list_target[i] :
        i += 1
    return i

# OK - 
def test_team_league(team_id_var,league_id_var):
    try:
        test = 0
        team_id_var, league_id_var = int(team_id_var), int(league_id_var)
        var = pd.read_pickle(link_data+central_folder+'/'+link_data_package_teams_central+'/'+'teams_league_id_s'+ '_'+ str(league_id_var) + '.pkl')
        for i in range (len(var)):
            if int(var.iloc[i,var.columns.get_loc('team_id')]) == team_id_var:
                test = 1
        if test == 0 :
            return False
        else:
            return True
    except:
        print("no matching of team within league_id for league_id{}".format(league_id_var))            


"""
For the reports
"""

 # OK -
def make_autopct(values):
    def my_autopct(pct):
        total = sum(values)
        val = int(round(pct*total/100))
        if pct>0:
            return '{p:.0f}%'.format(p=pct)

        #return '{p:.0f}%  ({v:d})'.format(p=pct,v=val)
    return my_autopct

# OK -
def cm_to_ppx(x,prs):
    return (x*prs.slide_width)/prs.slide_width.cm


# OK -
def text_modification_paragraph_center(pres_path,old_text,new_text,font_size,font_name,font_color,bold_bol,italic_bol):
    pres_path, old_text, new_text = str(pres_path), str(old_text), str(new_text)
    # Normal shapes
    prs = Presentation(pres_path)    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                tf = shape.text_frame
                for p in tf.paragraphs:
                    #print(p.text)
                    if p.text == old_text:
                        try:    
                            p.text = new_text
                            p.font.name = font_name
                            p.font.size = Pt(int(font_size))
                            p.font.color.rgb = font_color
                            p.alignment = PP_ALIGN.CENTER
                            p.font.bold = bold_bol
                            p.font.italic = italic_bol  
                        except Exception as e:
                            print(e) 
    # Group shapes
    for slide in prs.slides:
        group_shapes = [shp for shp in slide.shapes if shp.shape_type == MSO_SHAPE_TYPE.GROUP]
        for group_shape in group_shapes:
            for shape in group_shape.shapes:
                if shape.has_text_frame:
                    tf = shape.text_frame
                    for p in tf.paragraphs:
                        #print(p.text)
                        if p.text == old_text:
                            try:    
                                p.text = new_text
                                p.font.name = font_name
                                p.font.size = Pt(int(font_size))
                                p.font.color.rgb = font_color
                                p.alignment = PP_ALIGN.CENTER
                                p.font.bold = bold_bol
                                p.font.italic = italic_bol  
                            except Exception as e:
                                print(e) 
    prs.save(pres_path)

# OK -
def text_modification_paragraph_left(pres_path,old_text,new_text,font_size,font_name,font_color,bold_bol,italic_bol):
    pres_path, old_text, new_text = str(pres_path), str(old_text), str(new_text)
    # Normal shapes
    prs = Presentation(pres_path)    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                tf = shape.text_frame
                for p in tf.paragraphs:
                    #print(p.text)
                    if p.text == old_text:
                        try:    
                            p.text = new_text
                            p.font.name = font_name
                            p.font.size = Pt(int(font_size))
                            p.font.color.rgb = font_color
                            p.alignment = PP_ALIGN.LEFT
                            p.font.bold = bold_bol
                            p.font.italic = italic_bol  
                        except Exception as e:
                            print(e) 
    # Group shapes
    for slide in prs.slides:
        group_shapes = [shp for shp in slide.shapes if shp.shape_type == MSO_SHAPE_TYPE.GROUP]
        for group_shape in group_shapes:
            for shape in group_shape.shapes:
                if shape.has_text_frame:
                    tf = shape.text_frame
                    for p in tf.paragraphs:
                        #print(p.text)
                        if p.text == old_text:
                            try:    
                                p.alignment = PP_ALIGN.LEFT
                                p.text = new_text
                                p.font.name = font_name
                                p.font.size = Pt(int(font_size))
                                p.font.color.rgb = font_color
                                p.font.bold = bold_bol
                                p.font.italic = italic_bol  
                            except Exception as e:
                                print(e) 
    prs.save(pres_path)


# OK -
def move_box_advantage(pres_path,old_text_1,old_text_2,pos_left,pos_top):
    
    pres_path, old_text_1, old_text_2 = str(pres_path), str(old_text_1), str(old_text_2)
    prs = Presentation(pres_path)

    cpt = 0
    for slide in prs.slides:
        group_shapes = [shp for shp in slide.shapes if shp.shape_type == MSO_SHAPE_TYPE.GROUP]
        for group_shape in group_shapes:
            cpt = 0
            for shape in group_shape.shapes:
                if shape.has_text_frame:
                    tf = shape.text_frame
                    for p in tf.paragraphs:
                        if p.text == old_text_1 or (p.text == old_text_2 and p.runs[0].font.bold == True):
                            cpt += 1
                            if cpt == 2 :                    
                                try:    
                                    group_shape.left = int(cm_to_ppx(int(pos_left),prs))
                                    group_shape.top = int(cm_to_ppx(int(pos_top),prs))
                                except Exception as e:
                                    print(e) 
    prs.save(pres_path) 


def hex_to_RGB(hex):
    ''' "#FFFFFF" -> [255,255,255] '''
    # Pass 16 to the integer function for change of base
    return [int(hex[i:i+2], 16) for i in range(1,6,2)]


def RGB_to_hex(RGB):
    ''' [255,255,255] -> "#FFFFFF" '''
    # Components need to be integers for hex to make sense
    RGB = [int(x) for x in RGB]
    return "#"+"".join(["0{0:x}".format(v) if v < 16 else
                        "{0:x}".format(v) for v in RGB])

def color_dict(gradient):
    ''' Takes in a list of RGB sub-lists and returns dictionary of
      colors in RGB and hex form for use in a graphing function
      defined later on '''
    return {"hex":[RGB_to_hex(RGB) for RGB in gradient],
            "r":[RGB[0] for RGB in gradient],
            "g":[RGB[1] for RGB in gradient],
            "b":[RGB[2] for RGB in gradient]}

def linear_gradient(s, f="#FFFFFF", n=10):
    ''' returns a gradient list of (n) colors between
      two hex colors. start_hex and finish_hex
      should be the full six-digit color string,
      inlcuding the number sign ("#FFFFFF") '''
    # Starting and ending colors in RGB form
    #s = hex_to_RGB(start_hex)
    #f = hex_to_RGB(finish_hex)
    # Initilize a list of the output colors with the starting color
    RGB_list = [s]
    # Calcuate a color at each evenly spaced value of t from 1 to n
    for t in range(1, n):
        # Interpolate RGB vector for color at the current value of t
        curr_vector = [int(s[j] + (float(t)/(n-1))*(f[j]-s[j])) for j in range(3)]
        # Add it to our list of output colors
        RGB_list.append(curr_vector)

    return color_dict(RGB_list)



# OK -
def test_2(pres_path):
    pres_path = str(pres_path)
    # Normal shapes
    prs = Presentation(pres_path)    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                shape = shape.table
                cell = shape.cell(0, 0)
                cell.text = '3'
                print("yes")

    prs.save(pres_path) 


# OK -
def text_modification_table_center(pres_path,rows,cols,old_text,new_text,font_size,font_name,font_color,bold_bol,italic_bol):
    pres_path, new_text, old_text = str(pres_path), str(new_text), str(old_text)
    # Normal shapes
    prs = Presentation(pres_path)    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                shape = shape.table
                cell = shape.cell(rows, cols)
                if cell.text == old_text:
                    cell.text = new_text
                    run = cell.text_frame.paragraphs[0].runs[0]
                    run.font.size = Pt(int(font_size))
                    run.font.name = font_name
                    run.font.color.rgb = font_color
                    run.font.bold = bold_bol
                    run.font.italic = italic_bol
                    run.alignment = PP_ALIGN.CENTER
    prs.save(pres_path)

# OK -
def text_modification_table_left(pres_path,rows,cols,old_text,new_text,font_size,font_name,font_color,bold_bol,italic_bol):
    pres_path, new_text, old_text = str(pres_path), str(new_text), str(old_text)
    # Normal shapes
    prs = Presentation(pres_path)    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                shape = shape.table
                cell = shape.cell(rows, cols)
                if cell.text == old_text:
                    cell.text = new_text
                    run = cell.text_frame.paragraphs[0].runs[0]
                    run.font.size = Pt(int(font_size))
                    run.font.name = font_name
                    run.font.color.rgb = font_color
                    run.font.bold = bold_bol
                    run.font.italic = italic_bol
                    run.alignment = PP_ALIGN.CENTER
    prs.save(pres_path)



def text_modification_paragraph_left_line(pres_path,old_text,new_text,font_size,font_name,font_color,bold_bol,italic_bol,line_data):
    pres_path, old_text, new_text = str(pres_path), str(old_text), str(new_text)
    # Normal shapes
    prs = Presentation(pres_path)    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                tf = shape.text_frame
                for p in tf.paragraphs:
                    #print(p.text)
                    if p.text == old_text:
                        try:    
                            line = shape.line
                            line.color.rgb = black_color_report
                            line.width = int(line_data)
                            p.text = new_text
                            p.font.name = font_name
                            p.font.size = Pt(int(font_size))
                            p.font.color.rgb = font_color
                            p.alignment = PP_ALIGN.LEFT
                            p.font.bold = bold_bol
                            p.font.italic = italic_bol  
                        except Exception as e:
                            print(e) 
    # Group shapes
    for slide in prs.slides:
        group_shapes = [shp for shp in slide.shapes if shp.shape_type == MSO_SHAPE_TYPE.GROUP]
        for group_shape in group_shapes:
            for shape in group_shape.shapes:
                if shape.has_text_frame:
                    tf = shape.text_frame
                    for p in tf.paragraphs:
                        #print(p.text)
                        if p.text == old_text:
                            try:    
                                line = shape.line
                                line.color.rgb = black_color_report
                                line.width = int(line_data)
                                p.alignment = PP_ALIGN.LEFT
                                p.text = new_text
                                p.font.name = font_name
                                p.font.size = Pt(int(font_size))
                                p.font.color.rgb = font_color
                                p.font.bold = bold_bol
                                p.font.italic = italic_bol  
                            except Exception as e:
                                print(e) 
    prs.save(pres_path)

def text_modification_paragraph_center_line(pres_path,old_text,new_text,font_size,font_name,font_color,bold_bol,italic_bol,line_data):
    pres_path, old_text, new_text = str(pres_path), str(old_text), str(new_text)
    # Normal shapes
    prs = Presentation(pres_path)    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                tf = shape.text_frame
                for p in tf.paragraphs:
                    #print(p.text)
                    if p.text == old_text:
                        try:    
                            line = shape.line
                            line.color.rgb = black_color_report
                            line.width = int(line_data)
                            p.text = new_text
                            p.font.name = font_name
                            p.font.size = Pt(int(font_size))
                            p.font.color.rgb = font_color
                            p.alignment = PP_ALIGN.CENTER
                            p.font.bold = bold_bol
                            p.font.italic = italic_bol  
                        except Exception as e:
                            print(e) 
    # Group shapes
    for slide in prs.slides:
        group_shapes = [shp for shp in slide.shapes if shp.shape_type == MSO_SHAPE_TYPE.GROUP]
        for group_shape in group_shapes:
            for shape in group_shape.shapes:
                if shape.has_text_frame:
                    tf = shape.text_frame
                    for p in tf.paragraphs:
                        #print(p.text)
                        if p.text == old_text:
                            try:    
                                line = shape.line
                                line.color.rgb = black_color_report
                                line.width = int(line_data)
                                p.alignment = PP_ALIGN.CENTER
                                p.text = new_text
                                p.font.name = font_name
                                p.font.size = Pt(int(font_size))
                                p.font.color.rgb = font_color
                                p.font.bold = bold_bol
                                p.font.italic = italic_bol  
                            except Exception as e:
                                print(e) 
    prs.save(pres_path)
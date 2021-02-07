# Imports
import requests
import json
import pandas as pd  
import numpy as np
import matplotlib as mpl
import matplotlib.pyplot as plt
from matplotlib import rcParams
from matplotlib.ticker import StrMethodFormatter
from pathlib import Path
import time
from matplotlib.lines import Line2D
import matplotlib.patches as mpatches
import plotly.graph_objects as go
from pptx.util import Inches
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
import shutil  
from shutil import copyfile
#from pptx.enum.shapes import MSO_SHAPE
#from pptx.enum.dml import MSO_FILL
from colour import Color
import webcolors
import math

from PIL import ImageColor

from pptx.util import Inches, Pt,Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import ColorFormat, RGBColor
import codecs

from parameters_package import *
from core_functions import * # For nan_to_zero function
from output_v2 import *
from summarize import *

import seaborn as sns
import warnings; warnings.filterwarnings(action='once')

'''
Package parameters
------------------------------------------------------------------------------------
'''

# The parameters for the font
mpl.rcParams["font.family"] = font_1

'''
Transfer functions
------------------------------------------------------------------------------------
'''

global team_id 
global opp_id  
global team_name 
global opp_name  
global game_fixture_id  


# OK - 
def composition_master(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_composition,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3):

    print('start of composition_master')

    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 

    # Save the current version of the report
    prs = Presentation(pres_path)

    # Charts
    formation_home, formation_home_data = composition_reports_section1(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_composition,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3)
    formation_opp, formation_opp_data = composition_reports_section2(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_composition,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3)
    fav_home_formation, fav_away_formation, home_results, away_results = composition_reports_section3(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_composition,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3)
    perf_home, perf_away = composition_reports_section4(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_composition,fav_home_formation,fav_away_formation,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3)
    home_mean_rotation = composition_reports_section5(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_composition,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3)
    away_mean_rotation = composition_reports_section6(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_composition,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3)
    
    # Other photos and media
    transfer_image_to_ppt_composition(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_composition,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3)

    # Texts
    composition_adv = transfer_text_to_ppt_composition(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_composition,fav_home_formation, fav_away_formation,home_results, away_results, formation_home, formation_home_data,formation_opp, formation_opp_data,home_mean_rotation,away_mean_rotation,perf_home, perf_away,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3)
    prs = Presentation(pres_path)

    # Advantage
    move_box_advantage(pres_path,"L'analyse des compositions et joueurs donnent l'avantage à",str(composition_adv),matrix_positions_composition['composition_advantage']['pos_left']+2*buffer_adv,matrix_positions_composition['composition_advantage']['pos_top']+2*buffer_adv)

    print('end of composition_master')


# OK - 
def transfer_image_to_ppt_composition(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_composition,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3):
    
    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    
    # Charge ppt presentation
    prs = Presentation(pres_path)
    
    # Logo home
    page_number, pos_left, pos_top = matrix_positions_composition['home_logo_1']['page_number'], matrix_positions_composition['home_logo_1']['pos_left'],matrix_positions_composition['home_logo_1']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(homeTeam_id) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_7,prs))
    
    page_number, pos_left, pos_top = matrix_positions_composition['home_logo_2']['page_number'], matrix_positions_composition['home_logo_2']['pos_left'],matrix_positions_composition['home_logo_2']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(homeTeam_id) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_7,prs))
    
    # Logo away
    page_number, pos_left, pos_top = matrix_positions_composition['away_logo_1']['page_number'], matrix_positions_composition['away_logo_1']['pos_left'],matrix_positions_composition['away_logo_1']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(awayTeam_id) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_7,prs))
    
    page_number, pos_left, pos_top = matrix_positions_composition['away_logo_2']['page_number'], matrix_positions_composition['away_logo_2']['pos_left'],matrix_positions_composition['away_logo_2']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + "logo_" + str(awayTeam_id) + pic_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_7,prs))

    # Save ppt presentation 
    prs.save(pres_path)


def transfer_text_to_ppt_composition(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_composition, fav_home_formation, fav_away_formation,home_results, away_results, formation_home, formation_home_data,formation_opp, formation_opp_data,home_mean_rotation,away_mean_rotation,perf_home, perf_away,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3):
    
    # Advantage parameters
    advantage_composition_home = 0
    advantage_composition_away = 0

    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    #prs = Presentation(pres_path)

    # Useful data part 1
    fixtures_home_away = pd.read_pickle(link_data+'{}'.format(team_id)+'/'+str(game_fixture_id) + '/' +'fixtures_package'+'/'+'fixtures_team_id_next_s'+ '_'+ str(team_id) +'_'+ str(nex) +'.pkl')
    rounds = fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('round')]
    
    league_data = pd.read_pickle(link_data+ central_folder +'/'+link_data_package_leagues_central+'/'+'leagues_league_id_master_s'+ '_'+ str(fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')]) +'.pkl')
    teams_league_data = pd.read_pickle(link_data+ central_folder +'/'+link_data_package_teams_central+'/'+'teams_league_id_s'+ '_'+ str(fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')]) +'.pkl')
    homeTeam_data = teams_league_data.loc[teams_league_data['team_id'] == int(homeTeam_id)]
    
    # Header
    text_modification_paragraph_center(pres_path,"COMPOSITION & PLAYERS","COMPOSITIONS & JOUEURS",font_size_titleh1,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"Observing for the past x games parameters including:","Analyse des paramètres suivants dans les {} derniers matchs".format(last_fixtures_lineups_duration),14,font_1,black_color_report,None,True)
    text_modification_paragraph_center(pres_path,"Team formations","Formations des équipes",14,font_1,black_color_report,None,True)
    text_modification_paragraph_center(pres_path,"Rotation rates & whole team dynamics","Taux de rotation des effectifs",14,font_1,black_color_report,None,True)
    text_modification_paragraph_center(pres_path,"Players performances","Performances individuelles",14,font_1,black_color_report,None,True)

    # Titles
    text_modification_paragraph_center(pres_path,"ANALYSIS OF LINEUPS","ANALYSE DES FORMATIONS",font_size_titleh2,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"ROTATION RATES","ROTATION D'EFFECTIFS",font_size_titleh2,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"WHOLE TEAM DYNAMICS","DYNAMIQUE D'ÉQUIPE",font_size_titleh2,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"PLAYERS PERFORMANCES","PERFORMANCES INDIVIDUELLES",font_size_titleh2,font_1,black_color_report,True,None)

    text_modification_paragraph_center(pres_path,"Fav home",fav_home_formation,24,font_1,white_color_report,True,None)
    text_modification_paragraph_center(pres_path,"Fav away",fav_away_formation,24,font_1,white_color_report,True,None)
    text_modification_paragraph_center(pres_path,"Formation home",homeTeam_name,12,font_1,home_color_report,True,None)
    text_modification_paragraph_center(pres_path,"Formation away",awayTeam_name,12,font_1,away_color_report,True,None)
    text_modification_paragraph_center(pres_path,"Number of goals score by subs in last x games","Nombre de buts par remplaçants ({} derniers matchs)".format(last_fixtures_lineups_duration),14,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"Number of assists by subs in last x games","Nombre de passes décisives par remplaçants ({} derniers matchs)".format(last_fixtures_lineups_duration),14,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"Average per game","Moyenne par match",14,font_1,draw_color_report,True,True)
    text_modification_paragraph_center(pres_path,"Hometeam dyn",homeTeam_name,12,font_1,white_color_report,True,None)
    text_modification_paragraph_center(pres_path,"Awayteam dyn",awayTeam_name,12,font_1,white_color_report,True,None)
    
    # Beginning of populating of the section
    number_id = 1
    text_modification_paragraph_center(pres_path,"ucom{}".format(number_id),str(int(round(100*formation_home_data/last_fixtures_lineups_duration,0))) +  "%",20,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"mcom{}".format(number_id),"matchs en " + str(formation_home),font_size_bubble,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"lcom{}".format(number_id),homeTeam_name,font_size_bubble,font_1,home_color_report,True,None)

    number_id = 2
    text_modification_paragraph_center(pres_path,"ucom{}".format(number_id),str(int(round(100*formation_opp_data/last_fixtures_lineups_duration,0))) +  "%",20,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"mcom{}".format(number_id),"matchs vs. " + str(formation_home),font_size_bubble,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"lcom{}".format(number_id),homeTeam_name,font_size_bubble,font_1,home_color_report,True,None)

    number_id = 3
    text_modification_paragraph_center(pres_path,"ucom{}".format(number_id),str(int(round(100*home_results[1]/last_fixtures_lineups_duration,0))) +  "%",20,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"mcom{}".format(number_id),"victoires en {}".format(fav_home_formation),font_size_bubble,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"lcom{}".format(number_id),homeTeam_name,font_size_bubble,font_1,home_color_report,True,None)

    text_modification_paragraph_left(pres_path,"Comment on home","{} a gagné {}".format(homeTeam_name,int(round(100*perf_home,0))) +"%"+ " de ses matchs contre la formation de {} ({})".format(awayTeam_name,fav_away_formation),14,font_1,black_color_report,True,None)
    text_modification_paragraph_left(pres_path,"Comment on away","{} a gagné {}".format(awayTeam_name,int(round(100*perf_away,0))) +"%"+ " de ses matchs contre la formation de {} ({})".format(homeTeam_name,fav_home_formation),14,font_1,black_color_report,True,None)

    number_id = 4
    text_modification_paragraph_center(pres_path,"ucom{}".format(number_id),str(int(round(100*home_mean_rotation[0]/11,0))) +  "%",20,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"mcom{}".format(number_id),"rotation moyenne",font_size_bubble,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"lcom{}".format(number_id),homeTeam_name,font_size_bubble,font_1,home_color_report,True,None)

    number_id = 5
    text_modification_paragraph_center(pres_path,"ucom{}".format(number_id),str(int(round(100*away_mean_rotation[0]/11,0))) +  "%",20,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"mcom{}".format(number_id),"rotation moyenne",font_size_bubble,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"lcom{}".format(number_id),awayTeam_name,font_size_bubble,font_1,away_color_report,True,None)
    
    sgh, sga, asgh, asga, sah, saa, asah, asaa = composition_reports_section7_data(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_composition,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3)

    text_modification_paragraph_center(pres_path,"sgh",sgh,14,font_1,white_color_report,True,None)
    text_modification_paragraph_center(pres_path,"sga",sga,14,font_1,white_color_report,True,None)
    text_modification_paragraph_center(pres_path,"asgh",asgh,14,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"asga",asga,14,font_1,black_color_report,True,None)
    
    text_modification_paragraph_center(pres_path,"sah",sah,14,font_1,white_color_report,True,None)
    text_modification_paragraph_center(pres_path,"saa",saa,14,font_1,white_color_report,True,None)
    text_modification_paragraph_center(pres_path,"asah",asah,14,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"asaa",asaa,14,font_1,black_color_report,True,None)

    if ((asaa+asaa)/2)>((asah+asgh)/2):
        text_modification_paragraph_center(pres_path,"Team x seem to be using more the whole team","{} fait une meilleure utilisation de son effectif que {}".format(awayTeam_name,homeTeam_name),16,font_1,black_color_report,None,None)
    else:
        text_modification_paragraph_center(pres_path,"Team x seem to be using more the whole team","{} fait une meilleure utilisation de son effectif que {}".format(homeTeam_name,awayTeam_name),16,font_1,black_color_report,None,None)
    
    home_players, away_players = composition_reports_section8(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_composition,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3)
    
    if home_players>away_players:
        number_id = 6
        text_modification_paragraph_center(pres_path,"ucom{}".format(number_id),"+ " + str(int(100*round((home_players-away_players),2))) +"%",20,font_1,black_color_report,True,None)
        text_modification_paragraph_center(pres_path,"mcom{}".format(number_id),"joueurs utilisés par match",font_size_bubble,font_1,black_color_report,True,None)
        text_modification_paragraph_center(pres_path,"lcom{}".format(number_id),"par " + homeTeam_name,font_size_bubble,font_1,away_color_report,True,None)

    else:
        number_id = 6
        text_modification_paragraph_center(pres_path,"ucom{}".format(number_id),"+ " + str(int(100*round((-home_players+away_players),2))) +"%",20,font_1,black_color_report,True,None)
        text_modification_paragraph_center(pres_path,"mcom{}".format(number_id),"joueurs utilisés par match",font_size_bubble,font_1,black_color_report,True,None)
        text_modification_paragraph_center(pres_path,"lcom{}".format(number_id),"par " + awayTeam_name,font_size_bubble,font_1,away_color_report,True,None)
    
    # 1 criteria
    if perf_home>perf_away:
        advantage_composition_home += 1
    else:
        advantage_composition_away += 1
    
    # 2 criteria
    if home_mean_rotation[0]>away_mean_rotation[0]:
        advantage_composition_home += 1
    else:
        advantage_composition_away += 1

    # 3 criteria

    # Advantage
    text_modification_paragraph_left(pres_path,"Composition & players gives advantage to","L'analyse des compositions et joueurs donnent l'avantage à",14,font_1,white_color_report,True,None)

    prs = Presentation(pres_path)
    prs.save(pres_path)

    composition_adv = 0

    if advantage_composition_home>advantage_composition_away:
        text_modification_paragraph_center(pres_path,"+ ADVANTAGE COMPOSITION","+ " + str(homeTeam_name),24,font_1,white_color_report,True,None)  
        composition_adv = math.ceil(10*advantage_composition_home/(advantage_composition_home+advantage_composition_away))
        
    elif advantage_composition_home<advantage_composition_away:
        text_modification_paragraph_center(pres_path,"+ ADVANTAGE COMPOSITION","+ " + str(awayTeam_name),24,font_1,white_color_report,True,None)  
        composition_adv = math.ceil(10*advantage_composition_away/(advantage_composition_home+advantage_composition_away))
        
    else:
        text_modification_paragraph_center(pres_path,"+ ADVANTAGE COMPOSITION","Perfect equality",24,font_1,white_color_report,True,None)  

    return composition_adv


# OK -
# Here to show the formations animations

def composition_reports_section1(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_composition,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3):
    
    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)
    
    fixtures_home_away = pd.read_pickle(link_data+'{}'.format(team_id)+'/'+str(game_fixture_id) + '/' +'fixtures_package'+'/'+'fixtures_team_id_next_s'+ '_'+ str(team_id) +'_'+ str(nex) +'.pkl')

    df_home = df_raw_home_1.groupby('formation').size() 
    df_home = df_home.sort_values(ascending=False)

    df_away = df_raw_away_1.groupby('formation').size() 
    df_away = df_away.sort_values(ascending=False)

    list_formations = list(df_home.index)
    for f in list(df_away.index):
        change = 1
        for p in list_formations:
            if f == p :
                change = 0
        if change == 1:
            list_formations.append(f)
    ticks_size = 13
    label_size = ticks_size*1.15
    legend_size = ticks_size*1.15
    title_size = ticks_size*1.4

    if english == 1 :
        title = str("Lineups used during last 20 games")
        
    elif english == 0 :
        title = str("Formation utilisée 20 derniers matchs")
        
    fig, ax = plt.subplots(1,2)

    # B. Plot Graph 1    
    patches, texts, autotexts = ax[0].pie(df_home, autopct=make_autopct(df_home),shadow=False, labels = df_home.index,colors=[colors_data[df_home.index[i]] for i in range(len(df_home))], wedgeprops = {'linewidth': 6,'edgecolor' : white_color_plot}, rotatelabels=False,pctdistance=0.6,labeldistance=None, textprops = {'fontsize': label_size,'color':black_color_plot,'weight':"bold"})
    #plt.setp(autotexts, size=8, weight="bold") # This is for the labels in the wedges 
    plt.setp(texts, color='grey',weight='bold',size=label_size) # This is for the main labels
    
    center, r = patches[0].center, patches[0].r
    centre_circle = plt.Circle(center,0.25,fc=white_color_plot)
    edge_circle = plt.Circle(center,r,fill=False,ec=home_color_plot)
    ax[0].add_artist(centre_circle)
    ax[0].add_artist(edge_circle)

    #fig = plt.gcf()
    #ax = fig.gca()

    ax[0].axis('equal')
    ax[0].set_title(homeTeam_name,weight='bold',fontsize=title_size,color=home_color_plot)
    
    # C. Plot Graph 2
    patches, texts, autotexts = ax[1].pie(df_away, autopct=make_autopct(df_away),shadow=False, labels=df_away.index,colors=[colors_data[df_away.index[i]] for i in range(len(df_away))], wedgeprops = {'linewidth': 6,'edgecolor' : white_color_plot}, rotatelabels=False,pctdistance=0.6,labeldistance=None, textprops = {'fontsize': label_size,'color':black_color_plot,'weight':"bold"})
    plt.setp(texts, color='grey',weight='bold',size=label_size) # This is for the main labels
    
    center, r = patches[0].center, patches[0].r
    centre_circle = plt.Circle(center,0.25,fc=white_color_plot)
    edge_circle = plt.Circle(center,r,fill=False,ec=away_color_plot)
    ax[1].add_artist(centre_circle)
    ax[1].add_artist(edge_circle)

    ax[1].axis('equal')
    ax[1].set_title(awayTeam_name,weight='bold',fontsize=title_size,color=away_color_plot)
    
    fig.suptitle(title, fontsize=16)
    fig.subplots_adjust(top=0.8)

    #ncol= len(list_formations)
    # D. Plot Figure
    hand = [mpatches.Patch(color = colors_data[k], label=k) for k in list_formations]
    ax[1].legend(handles=hand,ncol= 3,borderpad=1,labelspacing=1,frameon=False, prop={'size': legend_size},loc = 'lower center', bbox_to_anchor = (0,-0.05,1,1),bbox_transform = plt.gcf().transFigure)
    
    #plt.tight_layout()
    plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_composition+'/'+'composition_reports_chart_section1'+graph_format, dpi = dpi_1)
    plt.show()
    
    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)

    page_number, pos_left, pos_top = matrix_positions_composition['composition_reports_chart_section1']['page_number'], matrix_positions_composition['composition_reports_chart_section1']['pos_left'],matrix_positions_composition['composition_reports_chart_section1']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_composition+'/'+'composition_reports_chart_section1'+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_1,prs))

    prs.save(pres_path)
    
    return df_home.index[0], df_home[0]
    

def composition_reports_section2(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_composition,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3):
    
    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)
    
    fixtures_home_away = pd.read_pickle(link_data+'{}'.format(team_id)+'/'+str(game_fixture_id) + '/' +'fixtures_package'+'/'+'fixtures_team_id_next_s'+ '_'+ str(team_id) +'_'+ str(nex) +'.pkl')

    df_home = df_raw_home_1.groupby('formation_opp').size() 
    df_home = df_home.sort_values(ascending=False)
    df_home = df_home[:3].append(pd.Series([last_fixtures_lineups_duration-(df_home[0]+df_home[1]+df_home[2])],index=['other']))

    df_away = df_raw_away_1.groupby('formation_opp').size() 
    df_away = df_away.sort_values(ascending=False)
    df_away = df_away[:3].append(pd.Series([last_fixtures_lineups_duration-(df_away[0]+df_away[1]+df_away[2])],index=['other']))

    list_formations = list(df_home.index)

    for f in list(df_away.index):
        change = 1
        for p in list_formations:
            if f == p :
                change = 0
        if change == 1:
            list_formations.append(f)
    ticks_size = 13
    label_size = ticks_size*1.15
    legend_size = ticks_size*1.15
    title_size = ticks_size*1.4
    
    list_formations.sort()

    if english == 1 :
        title = str("Lineups of opponents during last 20 games")
        
    elif english == 0 :
        title = str("Formation adversaires 20 derniers matchs")
        
    fig, ax = plt.subplots(1,2)

    # B. Plot Graph 1    
    patches, texts, autotexts = ax[0].pie(df_home, autopct=make_autopct(df_home),shadow=False, labels = df_home.index,colors=[colors_data[df_home.index[i]] for i in range(len(df_home))], wedgeprops = {'linewidth': 6,'edgecolor' : white_color_plot}, rotatelabels=False,pctdistance=0.6,labeldistance=None, textprops = {'fontsize': label_size,'color':black_color_plot,'weight':"bold"})
    #plt.setp(autotexts, size=8, weight="bold") # This is for the labels in the wedges 
    plt.setp(texts, color='grey',weight='bold',size=label_size) # This is for the main labels
    
    center, r = patches[0].center, patches[0].r
    centre_circle = plt.Circle(center,0.25,fc=white_color_plot)
    edge_circle = plt.Circle(center,r,fill=False,ec=home_color_plot)
    ax[0].add_artist(centre_circle)
    ax[0].add_artist(edge_circle)

    #fig = plt.gcf()
    #ax = fig.gca()

    ax[0].axis('equal')
    ax[0].set_title(homeTeam_name,weight='bold',fontsize=title_size,color=home_color_plot)
    
    # C. Plot Graph 2
    patches, texts, autotexts = ax[1].pie(df_away, autopct=make_autopct(df_away),shadow=False, labels=df_away.index,colors=[colors_data[df_away.index[i]] for i in range(len(df_away))], wedgeprops = {'linewidth': 6,'edgecolor' : white_color_plot}, rotatelabels=False,pctdistance=0.6,labeldistance=None, textprops = {'fontsize': label_size,'color':black_color_plot,'weight':"bold"})
    plt.setp(texts, color='grey',weight='bold',size=label_size) # This is for the main labels
    
    center, r = patches[0].center, patches[0].r
    centre_circle = plt.Circle(center,0.25,fc=white_color_plot)
    edge_circle = plt.Circle(center,r,fill=False,ec=away_color_plot)
    ax[1].add_artist(centre_circle)
    ax[1].add_artist(edge_circle)

    ax[1].axis('equal')
    ax[1].set_title(awayTeam_name,weight='bold',fontsize=title_size,color=away_color_plot)
    
    fig.suptitle(title, fontsize=16)
    fig.subplots_adjust(top=0.8)

    # D. Plot Figure
    hand = [mpatches.Patch(color = colors_data[k], label=k) for k in list_formations]
    ax[1].legend(handles=hand,ncol= 3,borderpad=1,labelspacing=1,frameon=False, prop={'size': legend_size},loc = 'lower center', bbox_to_anchor = (0,-0.125,1,1),bbox_transform = plt.gcf().transFigure)
    
    #plt.tight_layout()
    plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_composition+'/'+'composition_reports_chart_section2'+graph_format, dpi = dpi_1)
    plt.show()
    
    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)

    page_number, pos_left, pos_top = matrix_positions_composition['composition_reports_chart_section2']['page_number'], matrix_positions_composition['composition_reports_chart_section2']['pos_left'],matrix_positions_composition['composition_reports_chart_section2']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_composition+'/'+'composition_reports_chart_section2'+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_1,prs))

    prs.save(pres_path)

    return df_home.index[0], df_home[0]


def composition_reports_section3(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_composition,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3):
    
    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)
    
    fixtures_home_away = pd.read_pickle(link_data+'{}'.format(team_id)+'/'+str(game_fixture_id) + '/' +'fixtures_package'+'/'+'fixtures_team_id_next_s'+ '_'+ str(team_id) +'_'+ str(nex) +'.pkl')

    df_home = df_raw_home_1.groupby('formation').size() 
    df_home = df_home.sort_values(ascending=False)

    df_away = df_raw_away_1.groupby('formation').size() 
    df_away = df_away.sort_values(ascending=False)

    W_home, D_home, L_home, W_away, D_away, L_away = 0,0,0,0,0,0
    for i in range (len(df_raw_home_1)):
        if df_raw_home_1.iloc[i,df_raw_home_1.columns.get_loc('formation')] == df_home.index[0]:
            if df_raw_home_1.iloc[i,df_raw_home_1.columns.get_loc('W-D-L')] == 'W':
                W_home += 1
            elif df_raw_home_1.iloc[i,df_raw_home_1.columns.get_loc('W-D-L')] == 'D':
                D_home += 1
            else:
                L_home += 1
    for i in range (len(df_raw_away_1)):
        if df_raw_away_1.iloc[i,df_raw_away_1.columns.get_loc('formation')] == df_away.index[0]:
            if df_raw_away_1.iloc[i,df_raw_away_1.columns.get_loc('W-D-L')] == 'W':
                W_away += 1
            elif df_raw_away_1.iloc[i,df_raw_away_1.columns.get_loc('W-D-L')] == 'D':
                D_away += 1
            else:
                L_away += 1

    home_results, away_results = [D_home, L_home, W_home],[D_away, L_away, W_away]

    ticks_size = 13
    label_size = ticks_size*1.15
    legend_size = ticks_size*1.15
    title_size = ticks_size*1.4

    # A. Language parameters
    if english == 1 :
        labels = ["Draw","Lose","Win"] # Keep the same order
    elif english == 0 :
        labels = ["Nul","Défaite","Victoire"] # Keep the same order
    
    fig, ax = plt.subplots(1,2)

    # B. Plot Graph 1    
    patches, texts, autotexts = ax[0].pie(home_results, autopct=make_autopct(home_results),shadow=False, labels = labels, colors = [draw_color_plot,lose_color_plot,win_color_plot], wedgeprops = {'linewidth': 6,'edgecolor' : white_color_plot}, rotatelabels=False,pctdistance=0.6,labeldistance=None, textprops = {'fontsize': label_size,'color':black_color_plot,'weight':"bold"})
    #plt.setp(autotexts, size=8, weight="bold") # This is for the labels in the wedges 
    plt.setp(texts, color='grey',weight='bold',size=label_size) # This is for the main labels
    
    center, r = patches[0].center, patches[0].r
    centre_circle = plt.Circle(center,0.25,fc=white_color_plot)
    edge_circle = plt.Circle(center,r,fill=False,ec=home_color_plot)
    ax[0].add_artist(centre_circle)
    ax[0].add_artist(edge_circle)

    #fig = plt.gcf()
    #ax = fig.gca()

    ax[0].axis('equal')
    ax[0].set_title(homeTeam_name,weight='bold',fontsize=title_size,color=home_color_plot)
    
    # C. Plot Graph 2
    patches, texts, autotexts = ax[1].pie(away_results, autopct=make_autopct(away_results),shadow=False, labels=labels, colors=[draw_color_plot,lose_color_plot,win_color_plot], wedgeprops = {'linewidth': 6,'edgecolor' : white_color_plot}, rotatelabels=False,pctdistance=0.6,labeldistance=None, textprops = {'fontsize': label_size,'color':black_color_plot,'weight':"bold"})
    plt.setp(texts, color='grey',weight='bold',size=label_size) # This is for the main labels
    
    center, r = patches[0].center, patches[0].r
    centre_circle = plt.Circle(center,0.25,fc=white_color_plot)
    edge_circle = plt.Circle(center,r,fill=False,ec=away_color_plot)
    ax[1].add_artist(centre_circle)
    ax[1].add_artist(edge_circle)

    ax[1].axis('equal')
    ax[1].set_title(awayTeam_name,weight='bold',fontsize=title_size,color=away_color_plot)
    
    # D. Plot Figure
    w_patch, d_patch, l_patch = mpatches.Patch(color=win_color_plot, label=labels[2]), mpatches.Patch(color=draw_color_plot, label=labels[0]),mpatches.Patch(color=lose_color_plot, label=labels[1])
    ax[1].legend(handles=[w_patch,d_patch,l_patch] ,ncol= len(labels),borderpad=1,labelspacing=1,frameon=False, prop={'size': legend_size},loc = 'lower center', bbox_to_anchor = (0,0,1,1),bbox_transform = plt.gcf().transFigure)
    
    #plt.tight_layout()
    plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_composition+'/'+'composition_reports_chart_section3'+graph_format, dpi = dpi_1)
    plt.show()

    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)

    # Add the 2 formation pictures for home and away teams
    page_number, pos_left, pos_top = matrix_positions_composition['formation_home']['page_number'], matrix_positions_composition['formation_home']['pos_left'],matrix_positions_composition['formation_home']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' + str(df_home.index[0]) + png_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_5,prs))

    page_number, pos_left, pos_top = matrix_positions_composition['formation_away']['page_number'], matrix_positions_composition['formation_away']['pos_left'],matrix_positions_composition['formation_away']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports + central_folder + '/' + "other" + '/'+ str(df_away.index[0]) + png_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_5,prs))

    page_number, pos_left, pos_top = matrix_positions_composition['composition_reports_chart_section3']['page_number'], matrix_positions_composition['composition_reports_chart_section3']['pos_left'],matrix_positions_composition['composition_reports_chart_section3']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_composition+'/'+'composition_reports_chart_section3'+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_1,prs))

    prs.save(pres_path)

    fav_home_formation, fav_away_formation = df_home.index[0], df_away.index[0]
    
    return fav_home_formation, fav_away_formation, home_results, away_results


def composition_reports_section4(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_composition,fav_home_formation,fav_away_formation,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3):
    
    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)
    
    fixtures_home_away = pd.read_pickle(link_data+'{}'.format(team_id)+'/'+str(game_fixture_id) + '/' +'fixtures_package'+'/'+'fixtures_team_id_next_s'+ '_'+ str(team_id) +'_'+ str(nex) +'.pkl')
    
    W_home, D_home, L_home, W_away, D_away, L_away = 0,0,0,0,0,0
    for i in range (len(df_raw_home_1)):
        if df_raw_home_1.iloc[i,df_raw_home_1.columns.get_loc('formation_opp')] == fav_away_formation:
            if df_raw_home_1.iloc[i,df_raw_home_1.columns.get_loc('W-D-L')] == 'W':
                W_home += 1
            elif df_raw_home_1.iloc[i,df_raw_home_1.columns.get_loc('W-D-L')] == 'D':
                D_home += 1
            else:
                L_home += 1
    for i in range (len(df_raw_away_1)):
        if df_raw_away_1.iloc[i,df_raw_away_1.columns.get_loc('formation_opp')] == fav_home_formation:
            if df_raw_away_1.iloc[i,df_raw_away_1.columns.get_loc('W-D-L')] == 'W':
                W_away += 1
            elif df_raw_away_1.iloc[i,df_raw_away_1.columns.get_loc('W-D-L')] == 'D':
                D_away += 1
            else:
                L_away += 1

    home_results, away_results = [D_home, L_home, W_home],[D_away, L_away, W_away]
    perc_home, perc_away = W_home/(D_home+L_home+W_home), W_away/(D_away+L_away+W_away)

    ticks_size = 13
    label_size = ticks_size*1.15
    legend_size = ticks_size*1.15
    title_size = ticks_size*1.4

    # A. Language parameters
    if english == 1 :
        labels = ["Draw","Lose","Win"] # Keep the same order
        title = "Historic performance again opponent's formation"
    elif english == 0 :
        labels = ["Nul","Défaite","Victoire"] # Keep the same order
        title = "Performance vs. formation de l'adversaire"
    
    fig, ax = plt.subplots(1,2)

    # B. Plot Graph 1    
    patches, texts, autotexts = ax[0].pie(home_results, autopct=make_autopct(home_results),shadow=False, labels = labels, colors = [draw_color_plot,lose_color_plot,win_color_plot], wedgeprops = {'linewidth': 6,'edgecolor' : white_color_plot}, rotatelabels=False,pctdistance=0.6,labeldistance=None, textprops = {'fontsize': label_size,'color':black_color_plot,'weight':"bold"})
    #plt.setp(autotexts, size=8, weight="bold") # This is for the labels in the wedges 
    plt.setp(texts, color='grey',weight='bold',size=label_size) # This is for the main labels
    
    center, r = patches[0].center, patches[0].r
    centre_circle = plt.Circle(center,0.25,fc=white_color_plot)
    edge_circle = plt.Circle(center,r,fill=False,ec=home_color_plot)
    ax[0].add_artist(centre_circle)
    ax[0].add_artist(edge_circle)

    #fig = plt.gcf()
    #ax = fig.gca()

    ax[0].axis('equal')
    ax[0].set_title(homeTeam_name + " vs. " + fav_away_formation,weight='bold',fontsize=title_size-3,color=home_color_plot)
    
    # C. Plot Graph 2
    patches, texts, autotexts = ax[1].pie(away_results, autopct=make_autopct(away_results),shadow=False, labels=labels, colors=[draw_color_plot,lose_color_plot,win_color_plot], wedgeprops = {'linewidth': 6,'edgecolor' : white_color_plot}, rotatelabels=False,pctdistance=0.6,labeldistance=None, textprops = {'fontsize': label_size,'color':black_color_plot,'weight':"bold"})
    plt.setp(texts, color='grey',weight='bold',size=label_size) # This is for the main labels
    
    center, r = patches[0].center, patches[0].r
    centre_circle = plt.Circle(center,0.25,fc=white_color_plot)
    edge_circle = plt.Circle(center,r,fill=False,ec=away_color_plot)
    ax[1].add_artist(centre_circle)
    ax[1].add_artist(edge_circle)

    ax[1].axis('equal')
    ax[1].set_title(awayTeam_name + " vs. " + fav_home_formation,weight='bold',fontsize=title_size-3,color=away_color_plot)
    
    fig.suptitle(title, fontsize=16)
    fig.subplots_adjust(top=0.8)

    # D. Plot Figure
    w_patch, d_patch, l_patch = mpatches.Patch(color=win_color_plot, label=labels[2]), mpatches.Patch(color=draw_color_plot, label=labels[0]),mpatches.Patch(color=lose_color_plot, label=labels[1])
    ax[1].legend(handles=[w_patch,d_patch,l_patch] ,ncol= len(labels),borderpad=1,labelspacing=1,frameon=False, prop={'size': legend_size},loc = 'lower center', bbox_to_anchor = (0,0,1,1),bbox_transform = plt.gcf().transFigure)
    
    #plt.tight_layout()
    plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_composition+'/'+'composition_reports_chart_section4'+graph_format, dpi = dpi_1)
    plt.show()

    page_number, pos_left, pos_top = matrix_positions_composition['composition_reports_chart_section4']['page_number'], matrix_positions_composition['composition_reports_chart_section4']['pos_left'],matrix_positions_composition['composition_reports_chart_section4']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_composition+'/'+'composition_reports_chart_section4'+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_1,prs))

    prs.save(pres_path)
    
    return perc_home, perc_away


def composition_reports_section5(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_composition,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3):
    
    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)
    
    fixtures_home_away = pd.read_pickle(link_data+'{}'.format(team_id)+'/'+str(game_fixture_id) + '/' +'fixtures_package'+'/'+'fixtures_team_id_next_s'+ '_'+ str(team_id) +'_'+ str(nex) +'.pkl')

    players_home = [0]*(last_fixtures_lineups_duration-1)
    
    if english == 1 :
        average = "Average"
    elif english == 0 :
        average = "Moyenne"

    for i in range (len(df_raw_home_1)-1):
        for j in range (1,17):
            change = 0 
            x = df_raw_home_1.iloc[i,df_raw_home_1.columns.get_loc("P{}_stats_mod".format(j))]
            if type(x) == type([]) and x[0]['substitute'] == 'False':
                id = x[0]['player_id']
                for k in range(1,17):
                    if type(df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("P{}_stats_mod".format(k))]) == type([]) and df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("P{}_stats_mod".format(k))][0]['substitute'] == "False" and df_raw_home_1.iloc[i+1,df_raw_home_1.columns.get_loc("P{}_stats_mod".format(k))][0]['player_id'] == id:
                        change = 1
                if change == 0:
                    players_home[i] += 1
    ticks_size = 13
    label_size = ticks_size*1.15
    legend_size = ticks_size*1.15
    title_size = ticks_size*1.4
    
    if english == 1 :
        average = "Average"
    elif english == 0 :
        average = "Moyenne"

    fig, ax = plt.subplots(1,1)

    N = len(players_home)
    players_home.reverse()
    mean_home_rotation = [np.mean(players_home)]*N
    x = ["N-{}".format(i) for i in range (N,0,-1)]
    
    ax.plot(x, players_home,color=home_color_plot)
    ax.plot(x,mean_home_rotation, color=black_color_plot, linestyle='--')

    #fig = plt.gcf()
    #ax = fig.gca()

    ax.axis('equal')
    ax.set_title(homeTeam_name,weight='bold',fontsize=title_size,color=home_color_plot)
    
    every_nth = 3
    for n, label in enumerate(ax.xaxis.get_ticklabels()):
        if n % every_nth != 0:
            label.set_visible(False) 
    
    patch = [matplotlib.lines.Line2D(xdata=[1],ydata=[1],linestyle='--', color='black',label=average)]
    ax.legend(handles=patch,loc='upper right',fontsize=legend_size)

    #plt.tight_layout()
    plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_composition+'/'+'composition_reports_chart_section5'+graph_format, dpi = dpi_1)
    plt.show()

    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)

    page_number, pos_left, pos_top = matrix_positions_composition['composition_reports_chart_section5']['page_number'], matrix_positions_composition['composition_reports_chart_section5']['pos_left'],matrix_positions_composition['composition_reports_chart_section5']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_composition+'/'+'composition_reports_chart_section5'+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_6,prs))

    prs.save(pres_path)

    return mean_home_rotation


def composition_reports_section6(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_composition,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3):
    
    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)
    
    fixtures_home_away = pd.read_pickle(link_data+'{}'.format(team_id)+'/'+str(game_fixture_id) + '/' +'fixtures_package'+'/'+'fixtures_team_id_next_s'+ '_'+ str(team_id) +'_'+ str(nex) +'.pkl')

    players_away = [0]*(last_fixtures_lineups_duration-1)
    
    for i in range (len(df_raw_away_1)-1):
        for j in range (1,17):
            change = 0 
            x = df_raw_away_1.iloc[i,df_raw_away_1.columns.get_loc("P{}_stats_mod".format(j))]
            if type(x) == type([]) and x[0]['substitute'] == 'False':
                id = x[0]['player_id']
                for k in range(1,17):
                    if type(df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("P{}_stats_mod".format(k))]) == type([]) and df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("P{}_stats_mod".format(k))][0]['substitute'] == "False" and df_raw_away_1.iloc[i+1,df_raw_away_1.columns.get_loc("P{}_stats_mod".format(k))][0]['player_id'] == id:
                        change = 1
                if change == 0:
                    players_away[i] += 1
    ticks_size = 13
    label_size = ticks_size*1.15
    legend_size = ticks_size*1.15
    title_size = ticks_size*1.4
    
    if english == 1 :
        average = "Average"
    elif english == 0 :
        average = "Moyenne"

    fig, ax = plt.subplots(1,1)

    N = len(players_away)
    players_away.reverse()
    mean_away_rotation = [np.mean(players_away)]*N
    x = ["N-{}".format(i) for i in range (N,0,-1)]
    
    ax.plot(x, players_away,color=away_color_plot)
    ax.plot(x,mean_away_rotation, color=black_color_plot, linestyle='--')

    #fig = plt.gcf()
    #ax = fig.gca()

    ax.axis('equal')
    ax.set_title(awayTeam_name,weight='bold',fontsize=title_size,color=away_color_plot)
    
    every_nth = 3
    for n, label in enumerate(ax.xaxis.get_ticklabels()):
        if n % every_nth != 0:
            label.set_visible(False) 
    
    patch = [matplotlib.lines.Line2D(xdata=[1],ydata=[1],linestyle='--', color='black',label=average)]
    ax.legend(handles=patch,loc='upper right',fontsize=legend_size)

    #plt.tight_layout()
    plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_composition+'/'+'composition_reports_chart_section6'+graph_format, dpi = dpi_1)
    plt.show()

    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)

    page_number, pos_left, pos_top = matrix_positions_composition['composition_reports_chart_section6']['page_number'], matrix_positions_composition['composition_reports_chart_section6']['pos_left'],matrix_positions_composition['composition_reports_chart_section6']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_composition+'/'+'composition_reports_chart_section6'+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_6,prs))

    prs.save(pres_path)

    return mean_away_rotation


def composition_reports_section7_data(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_composition,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3):
    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)

    sgh, sga, asgh, asga, sah, saa, asah, asaa = 0,0,0,0,0,0,0,0
    cpt_1, cpt_2 =0,0
    for i in range (len(df_raw_home_1)):
        try:
            sgh += df_raw_home_1.iloc[i,df_raw_home_1.columns.get_loc("BP_subs")]
            sah += df_raw_home_1.iloc[i,df_raw_home_1.columns.get_loc("ABP_subs")]
            cpt_1 +=1
        except:
            pass
        try:
            sga += df_raw_away_1.iloc[i,df_raw_away_1.columns.get_loc("BP_subs")]
            saa += df_raw_away_1.iloc[i,df_raw_away_1.columns.get_loc("ABP_subs")]
            cpt_2 +=1
        except: 
            pass
    
    asgh, asga = round(sgh/cpt_1,1), round(sga/cpt_2,1)
    asah, asaa = round(sah/cpt_1,1), round(saa/cpt_2,1)

    # Add the logos
    page_number, pos_left, pos_top = matrix_positions_composition['home_logo_1']['page_number'], matrix_positions_composition['home_logo_1']['pos_left'],matrix_positions_composition['home_logo_1']['pos_top'] 
    pic = prs.slides[page_number].shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' +'logo_'+ str(homeTeam_id) + png_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(1.68,prs))

    page_number, pos_left, pos_top = matrix_positions_composition['home_logo_2']['page_number'], matrix_positions_composition['home_logo_2']['pos_left'],matrix_positions_composition['home_logo_2']['pos_top'] 
    pic = prs.slides[page_number].shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' +'logo_'+ str(homeTeam_id) + png_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(1.68,prs))

    page_number, pos_left, pos_top = matrix_positions_composition['away_logo_1']['page_number'], matrix_positions_composition['away_logo_1']['pos_left'],matrix_positions_composition['away_logo_1']['pos_top'] 
    pic = prs.slides[page_number].shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' +'logo_'+ str(awayTeam_id) + png_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(1.68,prs))

    page_number, pos_left, pos_top = matrix_positions_composition['away_logo_2']['page_number'], matrix_positions_composition['away_logo_2']['pos_left'],matrix_positions_composition['away_logo_2']['pos_top'] 
    pic = prs.slides[page_number].shapes.add_picture(link_reports + central_folder + '/' + "other" + '/' +'logo_'+ str(awayTeam_id) + png_format , left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(1.68,prs))

    prs.save(pres_path)

    return sgh, sga, asgh, asga, sah, saa, asah, asaa


def composition_reports_section8(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_composition,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3):
    
    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)
    
    ticks_size = 13
    label_size = ticks_size*1.15
    legend_size = ticks_size*1.15
    title_size = ticks_size*1.4

    if english == 1 :
        title = str("Average number of players used per game")
        
    elif english == 0 :
        title = str("Nombre moyen de joueurs utilisés par match")
        
    labels = [homeTeam_name,awayTeam_name] 
    
    home_players, away_players = 0,0
    home_players_n, away_players_n = 0,0
    a,b = 0,0
    for i in range (len(df_raw_home_1)):
        a,b = 0,0
        for j in range (1,17):
            if type(df_raw_home_1.iloc[i,df_raw_home_1.columns.get_loc("P{}_stats_mod".format(j))]) == type([]):
                home_players += 1
                a = 1
            if type(df_raw_away_1.iloc[i,df_raw_away_1.columns.get_loc("P{}_stats_mod".format(j))]) == type([]):
                away_players += 1
                b =1
        if a == 1 :
            home_players_n += 1
        if b == 1:
            away_players_n += 1

    fig, ax = plt.subplots(1,1)
    ind = np.arange(len(labels))
    width = 0.30
    ax.bar(ind, [round(home_players/home_players_n,1),round(away_players/away_players_n,1)], width, label=labels, color=[home_color_plot,away_color_plot],edgecolor='black')

    ax.set_title(title,weight='bold',fontsize=title_size,color=black_color_plot)    
    
    for p in ax.patches:
        width = p.get_width()
        height = p.get_height()
        x, y = p.get_xy() 
        ax.annotate(str(height), (x+width/2 , y + height*1.01), ha='center',fontsize=15)  

    ax.set_ylim([min([home_players/home_players_n,away_players/away_players_n])-1,max([home_players/home_players_n,away_players/away_players_n])+0.5])
    ax.set_yticks([], minor=True)
    ax.set_yticks([])
    
    plt.xticks(ind , labels,fontsize=title_size)
    
    plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_composition+'/'+'composition_reports_chart_section8'+graph_format, dpi = dpi_1)
    plt.show()
    
    page_number, pos_left, pos_top = matrix_positions_composition['composition_reports_chart_section8']['page_number'], matrix_positions_composition['composition_reports_chart_section8']['pos_left'],matrix_positions_composition['composition_reports_chart_section8']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_composition+'/'+'composition_reports_chart_section8'+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_1,prs))

    prs.save(pres_path)

    home_players, away_players = home_players/home_players_n, away_players/away_players_n

    return home_players, away_players

    


# OK - 
# Here to plot the players profiles
def composition_reports_section_defender(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_composition,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3):
        
    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)
    
    categories = ["%"+" duels gagnés","%"+" passes réussies","%"+" passes total équipe","%"+" fautes total équipe","%"+" cartons total équipe"]
    
    N = len(categories)
    
    pi = math.pi
    values_home = {'D1':[],'D2':[],'D3':[],'D4':[],'D5':[]}
    values_home_other = {'D1':[],'D2':[],'D3':[],'D4':[],'D5':[]}

    values_away = {'D1':[],'D2':[],'D3':[],'D4':[],'D5':[]}
    values_away_other = {'D1':[],'D2':[],'D3':[],'D4':[],'D5':[]}
    
    # What will be the angle of each axis in the plot? (we divide the plot / number of variable)
    angles = [n / float(N) * 2 * pi for n in range(N)]
    angles += angles[:1]

    # HOME PLAYERS
    for i in range (1,6): # Number of defender in the past formation
        x = df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('^T_BC_D{}'.format(i))]
        if type(x) == type([]):
            values_home['D{}'.format(i)].append(x[0]['duels_w_XI']) # 1
            values_home['D{}'.format(i)].append(x[0]['passes_acc_XI%']) # 2
            values_home['D{}'.format(i)].append(x[0]['passes_team_XI%']) # 3
            values_home['D{}'.format(i)].append(x[0]['yellow_team_XI%']) # 4
            values_home['D{}'.format(i)].append(x[0]['fouls_XI_team%']) # 5

            values_home_other['D{}'.format(i)].append(x[0]['#XI']) # 1
            values_home_other['D{}'.format(i)].append(x[0]['#subs']) # 2

        print(values_home['D{}'.format(i)])
        # Initialise the spider plot
        ax = plt.subplot(111, polar=True)
    
        # If you want the first axis to be on top:
        ax.set_theta_offset(pi / 2)
        ax.set_theta_direction(+1)
    
        # Draw one axe per variable + add labels labels yet
        plt.xticks(angles[:-1], categories,fontsize=15)
        
        # Draw ylabels
        ax.set_rlabel_position(0)
        plt.yticks([10,20,30,40,50,60,70,80,90,100], ["10","20","30","40","50","60","70","80","90","100"], color="grey", size=7)
        plt.ylim(0,100)
    
        values_home['D{}'.format(i)] += values_home['D{}'.format(i)][:1]
        ax.plot(angles, values_home['D{}'.format(i)], defender_color_plot,linewidth=1, linestyle='solid')
        ax.fill(angles, values_home['D{}'.format(i)], ec=defender_color_plot , fc=defender_color_plot, alpha=0.1)
        plt.tight_layout()
        plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_composition+'/'+'composition_reports_section_defender_home_'+str(i)+graph_format, dpi = dpi_1)
        plt.show()
        prs.save(pres_path)

        page_number, pos_left, pos_top = matrix_positions_composition['composition_reports_section_defender_home_'+str(i)]['page_number'], matrix_positions_composition['composition_reports_section_defender_home_'+str(i)]['pos_left'],matrix_positions_composition['composition_reports_section_defender_home_'+str(i)]['pos_top'] 
        slide = prs.slides[page_number]
        pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_composition+'/'+'composition_reports_section_defender_home_'+str(i)+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_8,prs))

    # AWAY PLAYERS
    for i in range (1,6): # Number of defender in the past formation
        x = df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('^T_BC_D{}'.format(i))]
        if type(x) == type([]):
            values_away['D{}'.format(i)].append(x[0]['duels_w_XI']) # 1
            values_away['D{}'.format(i)].append(x[0]['passes_acc_XI%']) # 2
            values_away['D{}'.format(i)].append(x[0]['passes_team_XI%']) # 3
            values_away['D{}'.format(i)].append(x[0]['yellow_team_XI%']) # 4
            values_away['D{}'.format(i)].append(x[0]['fouls_XI_team%']) # 5
            
            values_away_other['D{}'.format(i)].append(x[0]['#XI']) # 1
            values_away_other['D{}'.format(i)].append(x[0]['#subs']) # 2

        print(values_away['D{}'.format(i)])
        # Initialise the spider plot
        ax = plt.subplot(111, polar=True)
    
        # If you want the first axis to be on top:
        ax.set_theta_offset(pi / 2)
        ax.set_theta_direction(+1)
    
        # Draw one axe per variable + add labels labels yet
        plt.xticks(angles[:-1], categories,fontsize=15)
        
        # Draw ylabels
        ax.set_rlabel_position(0)
        plt.yticks([10,20,30,40,50,60,70,80,90,100], ["10","20","30","40","50","60","70","80","90","100"], color="grey", size=7)
        plt.ylim(0,100)
    
        values_away['D{}'.format(i)] += values_away['D{}'.format(i)][:1]
        ax.plot(angles, values_away['D{}'.format(i)], defender_color_plot,linewidth=1, linestyle='solid')
        ax.fill(angles, values_away['D{}'.format(i)], ec=defender_color_plot , fc=defender_color_plot, alpha=0.1)
        plt.tight_layout()
        plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_composition+'/'+'composition_reports_section_defender_away_'+str(i)+graph_format, dpi = dpi_1)
        plt.show()
        prs.save(pres_path)
        
        page_number, pos_left, pos_top = matrix_positions_composition['composition_reports_section_defender_away_'+str(i)]['page_number'], matrix_positions_composition['composition_reports_section_defender_away_'+str(i)]['pos_left'],matrix_positions_composition['composition_reports_section_defender_away_'+str(i)]['pos_top'] 
        slide = prs.slides[page_number]
        pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_composition+'/'+'composition_reports_section_defender_away_'+str(i)+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_8,prs))

    return values_home_other, values_away_other

# OK - 
# Here to plot the players profiles
def composition_reports_section_midfielder(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_composition,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3):
        
    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)
    
    categories = ["%"+" passes résussies","%"+" passes total équipe","%"+" passes décisives","%"+" tirs total équipe","%"+" fautes total équipe"]
    
    N = len(categories)
    
    pi = math.pi
    values_home = {'M1':[],'M2':[],'M3':[],'M4':[],'M5':[],'M6':[]}
    values_home_other = {'M1':[],'M2':[],'M3':[],'M4':[],'M5':[],'M6':[]}

    values_away = {'M1':[],'M2':[],'M3':[],'M4':[],'M5':[],'M6':[]}
    values_away_other = {'M1':[],'M2':[],'M3':[],'M4':[],'M5':[],'M6':[]}

    # What will be the angle of each axis in the plot? (we divide the plot / number of variable)
    angles = [n / float(N) * 2 * pi for n in range(N)]
    angles += angles[:1]

    # HOME PLAYERS
    for i in range (1,7): # Number of midfielder in the past formation
        x = df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('^T_Pass_M{}'.format(i))]
        if type(x) == type([]):
            values_home['M{}'.format(i)].append(x[0]['passes_acc_XI%']) # 1
            values_home['M{}'.format(i)].append(x[0]['passes_XI_team%']) # 2
            values_home['M{}'.format(i)].append(x[0]['assists_XI']) # 3
            values_home['M{}'.format(i)].append(x[0]['shots_XI']) # 4
            values_home['M{}'.format(i)].append(x[0]['fouls_XI']) # 5

            values_home_other['M{}'.format(i)].append(x[0]['#XI']) # 1
            values_home_other['M{}'.format(i)].append(x[0]['#subs']) # 2


        print(values_home['M{}'.format(i)])
        # Initialise the spider plot
        ax = plt.subplot(111, polar=True)
    
        # If you want the first axis to be on top:
        ax.set_theta_offset(pi / 2)
        ax.set_theta_direction(+1)
    
        # Draw one axe per variable + add labels labels yet
        plt.xticks(angles[:-1], categories,fontsize=15)
        
        # Draw ylabels
        ax.set_rlabel_position(0)
        plt.yticks([10,20,30,40,50,60,70,80,90,100], ["10","20","30","40","50","60","70","80","90","100"], color="grey", size=7)
        plt.ylim(0,100)
    
        values_home['M{}'.format(i)] += values_home['M{}'.format(i)][:1]
        ax.plot(angles, values_home['M{}'.format(i)], midfielder_color_plot,linewidth=1, linestyle='solid')
        ax.fill(angles, values_home['M{}'.format(i)], ec=midfielder_color_plot , fc=midfielder_color_plot, alpha=0.1)
        plt.tight_layout()
        plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_composition+'/'+'composition_reports_section_midfielder_home_'+str(i)+graph_format, dpi = dpi_1)
        plt.show()
        prs.save(pres_path)

        page_number, pos_left, pos_top = matrix_positions_composition['composition_reports_section_midfielder_home_'+str(i)]['page_number'], matrix_positions_composition['composition_reports_section_midfielder_home_'+str(i)]['pos_left'],matrix_positions_composition['composition_reports_section_midfielder_home_'+str(i)]['pos_top'] 
        slide = prs.slides[page_number]
        pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_composition+'/'+'composition_reports_section_midfielder_home_'+str(i)+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_8,prs))

    # AWAY PLAYERS
    for i in range (1,7): # Number of midfielder in the past formation
        x = df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('^T_Pass_M{}'.format(i))]
        if type(x) == type([]):
            values_away['M{}'.format(i)].append(x[0]['passes_acc_XI%']) # 1
            values_away['M{}'.format(i)].append(x[0]['passes_XI_team%']) # 2
            values_away['M{}'.format(i)].append(x[0]['assists_subs']) # 3
            values_away['M{}'.format(i)].append(x[0]['shots_XI']) # 4
            values_away['M{}'.format(i)].append(x[0]['fouls_XI']) # 5

            values_away_other['M{}'.format(i)].append(x[0]['#XI']) # 1
            values_away_other['M{}'.format(i)].append(x[0]['#subs']) # 2

        print(values_away['M{}'.format(i)])
        # Initialise the spider plot
        ax = plt.subplot(111, polar=True)
    
        # If you want the first axis to be on top:
        ax.set_theta_offset(pi / 2)
        ax.set_theta_direction(+1)
    
        # Draw one axe per variable + add labels labels yet
        plt.xticks(angles[:-1], categories,fontsize=15)
        
        # Draw ylabels
        ax.set_rlabel_position(0)
        plt.yticks([10,20,30,40,50,60,70,80,90,100], ["10","20","30","40","50","60","70","80","90","100"], color="grey", size=7)
        plt.ylim(0,100)
    
        values_away['M{}'.format(i)] += values_away['M{}'.format(i)][:1]
        ax.plot(angles, values_away['M{}'.format(i)], midfielder_color_plot,linewidth=1, linestyle='solid')
        ax.fill(angles, values_away['M{}'.format(i)], ec=midfielder_color_plot , fc=midfielder_color_plot, alpha=0.1)
        plt.tight_layout()
        plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_composition+'/'+'composition_reports_section_midfielder_away_'+str(i)+graph_format, dpi = dpi_1)
        plt.show()
        prs.save(pres_path)
        
        page_number, pos_left, pos_top = matrix_positions_composition['composition_reports_section_midfielder_away_'+str(i)]['page_number'], matrix_positions_composition['composition_reports_section_midfielder_away_'+str(i)]['pos_left'],matrix_positions_composition['composition_reports_section_midfielder_away_'+str(i)]['pos_top'] 
        slide = prs.slides[page_number]
        pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_composition+'/'+'composition_reports_section_midfielder_away_'+str(i)+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_8,prs))

    return values_home_other, values_away_other


# OK - 
# Here to plot the players profiles
def composition_reports_section_forward(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_composition,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3):
        
    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)
    
    categories = #["%"+" duels gagnés","%"+" passes réussies","%"+" passes total équipe","%"+" fautes total équipe","%"+" cartons total équipe"]
    
    N = len(categories)
    
    pi = math.pi
    values_home = {'F1':[],'F2':[],'F3':[],'F4':[],'F5':[],'F6':[]}
    values_home_other = {'F1':[],'F2':[],'F3':[],'F4':[],'F5':[],'F6':[]}

    values_away = {'F1':[],'F2':[],'F3':[],'F4':[],'F5':[],'F6':[]}
    values_away_other = {'F1':[],'F2':[],'F3':[],'F4':[],'F5':[],'F6':[]}
    
    # What will be the angle of each axis in the plot? (we divide the plot / number of variable)
    angles = [n / float(N) * 2 * pi for n in range(N)]
    angles += angles[:1]

    # HOME PLAYERS
    for i in range (1,7): # Number of forward in the past formation
        x = df_raw_home_1.iloc[0,df_raw_home_1.columns.get_loc('^T_BP_F{}'.format(i))]
        if type(x) == type([]):
            values_home['F{}'.format(i)].append(x[0]['goals_XI']) # 1
            values_home['F{}'.format(i)].append(x[0]['assists_XI']) # 2
            values_home['F{}'.format(i)].append(x[0]['dribbles_w_XI']) # 3
            values_home['F{}'.format(i)].append(x[0]['shots_on_XI']) # 4
            values_home['F{}'.format(i)].append(x[0]['shots_XI']) # 5
            
            values_home_other['F{}'.format(i)].append(x[0]['#XI']) # 1
            values_home_other['F{}'.format(i)].append(x[0]['#subs']) # 2

        print(values_home['F{}'.format(i)])
        # Initialise the spider plot
        ax = plt.subplot(111, polar=True)
    
        # If you want the first axis to be on top:
        ax.set_theta_offset(pi / 2)
        ax.set_theta_direction(+1)
    
        # Draw one axe per variable + add labels labels yet
        plt.xticks(angles[:-1], categories,fontsize=15)
        
        # Draw ylabels
        ax.set_rlabel_position(0)
        plt.yticks([10,20,30,40,50,60,70,80,90,100], ["10","20","30","40","50","60","70","80","90","100"], color="grey", size=7)
        plt.ylim(0,100)
    
        values_home['F{}'.format(i)] += values_home['F{}'.format(i)][:1]
        ax.plot(angles, values_home['F{}'.format(i)], forward_color_plot,linewidth=1, linestyle='solid')
        ax.fill(angles, values_home['F{}'.format(i)], ec=forward_color_plot , fc=forward_color_plot, alpha=0.1)
        plt.tight_layout()
        plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_composition+'/'+'composition_reports_section_forward_home_'+str(i)+graph_format, dpi = dpi_1)
        plt.show()
        prs.save(pres_path)

        page_number, pos_left, pos_top = matrix_positions_composition['composition_reports_section_forward_home_'+str(i)]['page_number'], matrix_positions_composition['composition_reports_section_forward_home_'+str(i)]['pos_left'],matrix_positions_composition['composition_reports_section_forward_home_'+str(i)]['pos_top'] 
        slide = prs.slides[page_number]
        pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_composition+'/'+'composition_reports_section_forward_home_'+str(i)+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_8,prs))

    # AWAY PLAYERS
    for i in range (1,7): # Number of forward in the past formation
        x = df_raw_away_1.iloc[0,df_raw_away_1.columns.get_loc('^T_BP_F{}'.format(i))]
        if type(x) == type([]):
            
            values_away['F{}'.format(i)].append(x[0]['goals_XI']) # 1
            values_away['F{}'.format(i)].append(x[0]['assists_XI']) # 2
            values_away['F{}'.format(i)].append(x[0]['dribbles_w_XI']) # 3
            values_away['F{}'.format(i)].append(x[0]['shots_on_XI']) # 4
            values_away['F{}'.format(i)].append(x[0]['shots_XI']) # 5
            
            values_away_other['F{}'.format(i)].append(x[0]['#XI']) # 1
            values_away_other['F{}'.format(i)].append(x[0]['#subs']) # 2
        
        print(values_away['F{}'.format(i)])
        # Initialise the spider plot
        ax = plt.subplot(111, polar=True)
    
        # If you want the first axis to be on top:
        ax.set_theta_offset(pi / 2)
        ax.set_theta_direction(+1)
    
        # Draw one axe per variable + add labels labels yet
        plt.xticks(angles[:-1], categories,fontsize=15)
        
        # Draw ylabels
        ax.set_rlabel_position(0)
        plt.yticks([10,20,30,40,50,60,70,80,90,100], ["10","20","30","40","50","60","70","80","90","100"], color="grey", size=7)
        plt.ylim(0,100)
    
        values_away['F{}'.format(i)] += values_away['F{}'.format(i)][:1]
        ax.plot(angles, values_away['F{}'.format(i)], forward_color_plot,linewidth=1, linestyle='solid')
        ax.fill(angles, values_away['F{}'.format(i)], ec=forward_color_plot , fc=forward_color_plot, alpha=0.1)
        plt.tight_layout()
        plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_composition+'/'+'composition_reports_section_forward_away_'+str(i)+graph_format, dpi = dpi_1)
        plt.show()
        prs.save(pres_path)
        
        page_number, pos_left, pos_top = matrix_positions_composition['composition_reports_section_forward_away_'+str(i)]['page_number'], matrix_positions_composition['composition_reports_section_forward_away_'+str(i)]['pos_left'],matrix_positions_composition['composition_reports_section_forward_away_'+str(i)]['pos_top'] 
        slide = prs.slides[page_number]
        pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_composition+'/'+'composition_reports_section_forward_away_'+str(i)+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_8,prs))

    return values_home_other, values_away_other


'''
Other functions
------------------------------------------------------------------------------------
'''


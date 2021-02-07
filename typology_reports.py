# Imports
import requests
import json
import pandas as pd  
import numpy as np
import matplotlib as mpl
import matplotlib.pyplot as plt
from matplotlib import rcParams
from matplotlib.ticker import StrMethodFormatter
from pathlib import Path
import time
from matplotlib.lines import Line2D
import matplotlib.patches as mpatches
import plotly.graph_objects as go
from pptx.util import Inches
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
import shutil  
from shutil import copyfile
#from pptx.enum.shapes import MSO_SHAPE
#from pptx.enum.dml import MSO_FILL
from colour import Color
import webcolors
import math

from PIL import ImageColor

from pptx.util import Inches, Pt,Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import ColorFormat, RGBColor
import codecs

from parameters_package import *
from core_functions import * # For nan_to_zero function
from output_v2 import *
from summarize import *

import seaborn as sns
import warnings; warnings.filterwarnings(action='once')


'''
Package parameters
------------------------------------------------------------------------------------
'''

# The parameters for the font
mpl.rcParams["font.family"] = font_1

'''
Transfer functions
------------------------------------------------------------------------------------
'''

global team_id 
global opp_id  
global team_name 
global opp_name  
global game_fixture_id  

# OK - 
def typology_master(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_typology,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3):

    print("start of typology_master")
    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    
    prs = Presentation(pres_path)

    # Charts
    BP_30_home,BP_60_home, BP_90_home, BC_30_home, BC_60_home, BC_90_home = typology_reports_section1(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_typology,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3)
    BP_30_away,BP_60_away, BP_90_away, BC_30_away, BC_60_away, BC_90_away = typology_reports_section2(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_typology,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3)
    fouls_home_mean, yellow_home_mean, fouls_away_mean, yellow_away_mean = typology_reports_section3(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_typology,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3)
    values_BP_home,  values_BC_home = typology_reports_section4(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_typology,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3)
    values_BP_away,  values_BC_away = typology_reports_section5(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_typology,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3)
    home_mean_poss, away_mean_poss = typology_reports_section6(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_typology,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3)
    BP_direct_home, BP_direct_away = typology_reports_section7(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_typology,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3)
    # Other photos and docs to be added

    # Texts
    typology_adv = transfer_text_to_ppt_typology(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_typology,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3,BP_30_home,BP_60_home, BP_90_home, BC_30_home, BC_60_home, BC_90_home,BP_30_away,BP_60_away, BP_90_away, BC_30_away, BC_60_away, BC_90_away,fouls_home_mean, yellow_home_mean, fouls_away_mean, yellow_away_mean, values_BP_home,  values_BC_home, values_BP_away,  values_BC_away,home_mean_poss, away_mean_poss, BP_direct_home, BP_direct_away)
    prs = Presentation(pres_path)
    
    # Advantage
    move_box_advantage(pres_path,"Le style de jeu est à l'avantage de",str(typology_adv),matrix_positions_typology['typology_advantage']['pos_left']+2*buffer_adv,matrix_positions_typology['typology_advantage']['pos_top']+2*buffer_adv)
    
    print("end of typology_master")

def transfer_text_to_ppt_typology(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_typology,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3,BP_30_home,BP_60_home, BP_90_home, BC_30_home, BC_60_home, BC_90_home,BP_30_away,BP_60_away, BP_90_away, BC_30_away, BC_60_away, BC_90_away,fouls_home_mean, yellow_home_mean, fouls_away_mean, yellow_away_mean,values_BP_home,values_BC_home,values_BP_away,values_BC_away,home_mean_poss,away_mean_poss, BP_direct_home, BP_direct_away):
    
    # Advantage parameters
    advantage_typology_home = 0
    advantage_typology_away = 0

    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    

    # Useful data part 1
    fixtures_home_away = pd.read_pickle(link_data+'{}'.format(team_id)+'/'+str(game_fixture_id) + '/' +'fixtures_package'+'/'+'fixtures_team_id_next_s'+ '_'+ str(team_id) +'_'+ str(nex) +'.pkl')
    rounds = fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('round')]
    
    league_data = pd.read_pickle(link_data+ central_folder +'/'+link_data_package_leagues_central+'/'+'leagues_league_id_master_s'+ '_'+ str(fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')]) +'.pkl')
    teams_league_data = pd.read_pickle(link_data+ central_folder +'/'+link_data_package_teams_central+'/'+'teams_league_id_s'+ '_'+ str(fixtures_home_away.iloc[0,fixtures_home_away.columns.get_loc('league_id')]) +'.pkl')
    homeTeam_data = teams_league_data.loc[teams_league_data['team_id'] == int(homeTeam_id)]
    
    # Header
    text_modification_paragraph_center(pres_path,"TYPOLOGY OF PLAYING","TYPOLOGIE ET STYLES DE JEU",font_size_titleh1,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"Observing the teams' dynamics and typology of playing","Analyse des dynamiques et styles de jeu",14,font_1,black_color_report,None,True)

    text_modification_paragraph_center(pres_path,"When do teams score and concede goals?","A quels moments du match les équipes sont le plus efficaces?",14,font_1,black_color_report,None,True)
    text_modification_paragraph_center(pres_path,"Roughness and fouls committed","Force et dureté dans le jeu",14,font_1,black_color_report,None,True)
    text_modification_paragraph_center(pres_path,"Team-oriented or individuals-oriented?","Style de jeu tourné vers le collectif ou les individualités?",14,font_1,black_color_report,None,True)
    text_modification_paragraph_center(pres_path,"XXX","XXX",14,font_1,black_color_report,None,True)

    # Titles
    text_modification_paragraph_center(pres_path,"TEAM'S TIMING DURING THE GAME","TIMING DE L'ÉQUIPE",font_size_titleh2,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"ROUGHNESS AND FOULS COMMITTED","Force et dureté dans le jeu",font_size_titleh2,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"TEAM-ORIENTED vs. INDIVIDUALS-ORIENTED","STYLE DE JEU COLLECTIF vs. INDIVIDUALISTISTE",font_size_titleh2,font_1,black_color_report,True,None)


    BP_home_x = {sum(BP_30_home):"BP 0'-30'",sum(BP_60_home):"BP 30'-60'", sum(BP_90_home):"BP 60'-90'"}
    BP_home = [max(BP_home_x.keys()),BP_home_x[max(BP_home_x.keys())]]

    BC_home_x = {sum(BC_30_home):"BC 0'-30'",sum(BC_60_home):"BC 30'-60'", sum(BC_90_home):"BC 60'-90'"}
    BC_home = [max(BC_home_x.keys()),BC_home_x[max(BC_home_x.keys())]]
    
    BP_away_x = {sum(BP_30_away):"BP 0'-30'",sum(BP_60_away):"BP 30'-60'", sum(BP_90_away):"BP 60'-90'"}
    BP_away = [max(BP_away_x.keys()),BP_away_x[max(BP_away_x.keys())]]

    BC_away_x = {sum(BC_30_away):"BC 0'-30'",sum(BC_60_away):"BC 30'-60'", sum(BC_90_away):"BC 60'-90'"}
    BC_away = [max(BC_away_x.keys()),BC_away_x[max(BC_away_x.keys())]]

    number_id = 1
    text_modification_paragraph_center(pres_path,"utyp{}".format(number_id),BP_home[0],20,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"mtyp{}".format(number_id),BP_home[1],font_size_bubble,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"ltyp{}".format(number_id),str(homeTeam_name),font_size_bubble,font_1,home_color_report,True,None)

    number_id = 2
    text_modification_paragraph_center(pres_path,"utyp{}".format(number_id),BC_home[0],20,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"mtyp{}".format(number_id),BC_home[1],font_size_bubble,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"ltyp{}".format(number_id),str(homeTeam_name),font_size_bubble,font_1,home_color_report,True,None)

    number_id = 3
    text_modification_paragraph_center(pres_path,"utyp{}".format(number_id),BP_away[0],20,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"mtyp{}".format(number_id),BP_away[1],font_size_bubble,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"ltyp{}".format(number_id),str(awayTeam_name),font_size_bubble,font_1,away_color_report,True,None)

    number_id = 4
    text_modification_paragraph_center(pres_path,"utyp{}".format(number_id),BC_away[0],20,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"mtyp{}".format(number_id),BC_away[1],font_size_bubble,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"ltyp{}".format(number_id),str(awayTeam_name),font_size_bubble,font_1,away_color_report,True,None)

    # Beginning of populating of the section
    a, b = 0,0
    BP_home, BP_home_min = 0,0
    
    for i in range (len(df_raw_home_1)):
        try:
            if df_raw_home_1.iloc[i,df_raw_home_1.columns.get_loc('First')][0] == homeTeam_id :
                BP_home += 1
                BP_home_min += df_raw_home_1.iloc[i,df_raw_home_1.columns.get_loc('First')][1]
                b += 1
            a += 1
        except:
            pass
    
    BP_home, BP_home_min = int(100*round(BP_home/a,2)), int(round(BP_home_min/b,0))
    
    a, b = 0,0
    BP_away, BP_away_min = 0,0

    for i in range (len(df_raw_away_1)):
        try:
            if df_raw_away_1.iloc[i,df_raw_away_1.columns.get_loc('First')][0] == awayTeam_id :
                BP_away += 1
                BP_away_min += df_raw_away_1.iloc[i,df_raw_away_1.columns.get_loc('First')][1]
                b += 1
            a += 1
        except:
            pass
    
    BP_away, BP_away_min = int(100*round(BP_away/a,2)), int(round(BP_away_min/b))


    number_id = 5
    text_modification_paragraph_center(pres_path,"utyp{}".format(number_id),BP_home+" %",20,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"mtyp{}".format(number_id),"matchs premier but marqué",font_size_bubble,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"ltyp{}".format(number_id),str(homeTeam_name),font_size_bubble,font_1,home_color_report,True,None)

    number_id = 6
    text_modification_paragraph_center(pres_path,"utyp{}".format(number_id),BP_away+" %",20,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"mtyp{}".format(number_id),"matchs premier but marqué",font_size_bubble,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"ltyp{}".format(number_id),str(awayTeam_name),font_size_bubble,font_1,away_color_report,True,None)

    number_id = 7
    text_modification_paragraph_center(pres_path,"utyp{}".format(number_id),BP_home_min,44,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"mtyp{}".format(number_id),"moyenne",font_size_bubble,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"ltyp{}".format(number_id),"min premier but marqué",font_size_bubble,font_1,black_color_report,True,None)
    
    number_id = 8
    text_modification_paragraph_center(pres_path,"utyp{}".format(number_id),BP_away_min,44,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"mtyp{}".format(number_id),"moyenne",font_size_bubble,font_1,black_color_report,True,None)  
    text_modification_paragraph_center(pres_path,"ltyp{}".format(number_id),"min premier but marqué",font_size_bubble,font_1,black_color_report,True,None)  
    
    number_id = 9
    text_modification_paragraph_center(pres_path,"utyp{}".format(number_id),int(100*round(fouls_home_mean,2)),44,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"mtyp{}".format(number_id),"moyenne",font_size_bubble,font_1,black_color_report,True,None)  
    text_modification_paragraph_center(pres_path,"ltyp{}".format(number_id),"fautes {}".format(homeTeam_name),font_size_bubble,font_1,black_color_report,True,None)  

    number_id = 10
    text_modification_paragraph_center(pres_path,"utyp{}".format(number_id),int(100*round(fouls_away_mean,2)),44,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"mtyp{}".format(number_id),"moyenne",font_size_bubble,font_1,black_color_report,True,None)  
    text_modification_paragraph_center(pres_path,"ltyp{}".format(number_id),"fautes {}".format(awayTeam_name),font_size_bubble,font_1,black_color_report,True,None)  

    text_modification_paragraph_center(pres_path,"Not all goals are alike... analysis of significance of goals for and against","Tous les buts ne se ressemblent pas... analyse des 'type' de buts marqués",14,font_1,black_color_report,True,True)  
    text_modification_paragraph_center(pres_path,"Abnormal goals are:","Des buts 'anormaux' sont:",14,font_1,black_color_report,None,True)  
    text_modification_paragraph_left(pres_path,"Goals after opponent got a red card","Des buts marqués/pris après un carton rouge pour l'adversaire/équipe",14,font_1,black_color_report,None,True)  
    text_modification_paragraph_left(pres_path,"Goals after opponent got a yellow card","Des buts marqués/pris après un carton jaune pour l'adversaire/équipe",14,font_1,black_color_report,None,True)  
    text_modification_paragraph_left(pres_path,"Goals after opponent conceded a penalty","Des buts marqués/pris après un pénalty pour l'équipe/adversaire",14,font_1,black_color_report,None,True)  
    
    number_id = 11
    text_modification_paragraph_center(pres_path,"utyp{}".format(number_id),int(100*round(values_BP_home[1]/(values_BP_home[0]+values_BP_home[1]),2)),44,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"mtyp{}".format(number_id),"buts 'normaux'",font_size_bubble,font_1,black_color_report,True,None)  
    text_modification_paragraph_center(pres_path,"ltyp{}".format(number_id),"pour " + str(homeTeam_name),font_size_bubble,font_1,black_color_report,True,None)  

    number_id = 12
    text_modification_paragraph_center(pres_path,"utyp{}".format(number_id),int(100*round(values_BP_away[1]/(values_BP_away[0]+values_BP_away[1]),2)),44,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"mtyp{}".format(number_id),"buts 'normaux'",font_size_bubble,font_1,black_color_report,True,None)  
    text_modification_paragraph_center(pres_path,"ltyp{}".format(number_id),"pour " + str(awayTeam_name),font_size_bubble,font_1,black_color_report,True,None)  

    number_id = 13
    text_modification_paragraph_center(pres_path,"utyp{}".format(number_id),int(100*round(home_mean_poss,2)),44,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"mtyp{}".format(number_id),"possesion moyenne",font_size_bubble,font_1,black_color_report,True,None)  
    text_modification_paragraph_center(pres_path,"ltyp{}".format(number_id),str(homeTeam_name),font_size_bubble,font_1,black_color_report,True,None)  

    number_id = 14
    text_modification_paragraph_center(pres_path,"utyp{}".format(number_id),int(100*round(away_mean_poss,2)),44,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"mtyp{}".format(number_id),"possesion moyenne",font_size_bubble,font_1,black_color_report,True,None)  
    text_modification_paragraph_center(pres_path,"ltyp{}".format(number_id),str(awayTeam_name),font_size_bubble,font_1,black_color_report,True,None)  

    number_id = 15
    text_modification_paragraph_center(pres_path,"utyp{}".format(number_id),BP_direct_home,44,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"mtyp{}".format(number_id),"buts directs",font_size_bubble,font_1,black_color_report,True,None)  
    text_modification_paragraph_center(pres_path,"ltyp{}".format(number_id),str(homeTeam_name),font_size_bubble,font_1,black_color_report,True,None)  

    number_id = 16
    text_modification_paragraph_center(pres_path,"utyp{}".format(number_id),BP_direct_away,44,font_1,black_color_report,True,None)
    text_modification_paragraph_center(pres_path,"mtyp{}".format(number_id),"buts directs",font_size_bubble,font_1,black_color_report,True,None)  
    text_modification_paragraph_center(pres_path,"ltyp{}".format(number_id),str(awayTeam_name),font_size_bubble,font_1,black_color_report,True,None)  


    # Advantage
    text_modification_paragraph_left(pres_path,"Typology of playing gives advantage to","Le style de jeu est à l'avantage de",14,font_1,white_color_report,True,None)

    prs = Presentation(pres_path)
    prs.save(pres_path)

    typology_adv = 0
    
    if advantage_typology_home>advantage_typology_away:
        text_modification_paragraph_center(pres_path,"+ ADVANTAGE typology","+ " + str(homeTeam_name),24,font_1,white_color_report,True,None)  
        typology_adv = math.ceil(10*advantage_typology_home/(advantage_typology_home+advantage_typology_away))
        
    elif advantage_typology_home<advantage_typology_away:
        text_modification_paragraph_center(pres_path,"+ ADVANTAGE typology","+ " + str(awayTeam_name),24,font_1,white_color_report,True,None)  
        typology_adv = math.ceil(10*advantage_typology_away/(advantage_typology_home+advantage_typology_away))
        
    else:
        text_modification_paragraph_center(pres_path,"+ ADVANTAGE typology","Perfect equality",24,font_1,white_color_report,True,None)  

    return typology_adv


def typology_reports_section1(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_typology,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3):
    
    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)

    ticks_size = 13
    label_size = ticks_size*1.15
    legend_size = ticks_size*1.15
    title_size = ticks_size*1.4
    
    a = 0
    BP_30,BP_60,BP_90 = [],[],[]
    BC_30,BC_60,BC_90 = [],[],[]
    for i in range (len(df_raw_home_1)):
        try:
            BP_30.append(df_raw_home_1.iloc[i,df_raw_home_1.columns.get_loc('BP_0-30')])
            a += 1
        except:
            pass
        try:
            BP_60.append(df_raw_home_1.iloc[i,df_raw_home_1.columns.get_loc('BP_30-60')])
        except:
            pass
        try:
            BP_90.append(df_raw_home_1.iloc[i,df_raw_home_1.columns.get_loc('BP_60-100')])
        except:
            pass     

        try:
            BC_30.append(df_raw_home_1.iloc[i,df_raw_home_1.columns.get_loc('BC_0-30')])
        except:
            pass
        try:
            BC_60.append(df_raw_home_1.iloc[i,df_raw_home_1.columns.get_loc('BC_30-60')])
        except:
            pass
        try:
            BC_90.append(df_raw_home_1.iloc[i,df_raw_home_1.columns.get_loc('BC_60-100')])
        except:
            pass  
    
    BP_30.reverse()
    BP_60.reverse()
    BP_90.reverse()
    BC_30.reverse()
    BC_60.reverse()
    BC_90.reverse() 

    BP_30=np.array(BP_30)
    BP_60=np.array(BP_60)
    BP_90=np.array(BP_90)
    
    BC_30=np.array(BC_30)
    BC_60=np.array(BC_60)
    BC_90=np.array(BC_90)
    
    width = 0.35       # the width of the bars: can also be len(x) sequence

    # BP for home team
    fig, ax = plt.subplots(1,1)

    ind = np.arange(a)
    print(a)
    plt.xticks(ind, ["N-{}".format(i) for i in range (a,0,-1)],fontsize=ticks_size-2)
    plt.yticks([0,1,2,3,4,5,6,7,8,9,10,11])

    every_nth = 2
    for n, label in enumerate(ax.xaxis.get_ticklabels()):
        if n % every_nth != 0:
            label.set_visible(False)


    if english == 1 :
        title = str("Timing of scored goals by ")+homeTeam_name
        ylabel = str("Goals scored")
        
    elif english == 0 :
        title = str("Timing des buts marqués par ")+homeTeam_name
        ylabel = str("Buts marqués")

    p1 = ax.bar(ind, BP_30, width)
    p2 = ax.bar(ind, BP_60,width , bottom=BP_30)
    p3 = ax.bar(ind, BP_90,width , bottom=BP_30+BP_60)

    plt.ylabel(ylabel)
    plt.title(title)
    
    plt.legend((p1[0], p2[0] ,p3[0]), ("[0'-30']", "[30'-60']","[60'-90']"))

    #plt.tight_layout()
    plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_typology+'/'+'typology_reports_chart_section1'+graph_format, dpi = dpi_1)
    plt.show()
    
    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)

    page_number, pos_left, pos_top = matrix_positions_typology['typology_reports_chart_section1']['page_number'], matrix_positions_typology['typology_reports_chart_section1']['pos_left'],matrix_positions_typology['typology_reports_chart_section1']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_typology+'/'+'typology_reports_chart_section1'+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_6,prs))

    # BC for home Team
    if english == 1 :
        title = str("Timing of against goals for ")+homeTeam_name
        ylabel = str("Goals against")
        
    elif english == 0 :
        title = str("Timing des buts concédés par ")+homeTeam_name
        ylabel = str("Buts concédés")

    fig, ax = plt.subplots(1,1)

    ind = np.arange(a)
    print(a)
    plt.xticks(ind, ["N-{}".format(i) for i in range (a,0,-1)],fontsize=ticks_size-2)
    plt.yticks([0,1,2,3,4,5,6,7,8,9,10,11])

    every_nth = 2
    for n, label in enumerate(ax.xaxis.get_ticklabels()):
        if n % every_nth != 0:
            label.set_visible(False)
    
    p1 = ax.bar(ind, BC_30, width)
    p2 = ax.bar(ind, BC_60,width , bottom=BC_30)
    p3 = ax.bar(ind, BC_90,width , bottom=BC_30+BC_60)

    plt.ylabel(ylabel)
    plt.title(title)
    
    plt.legend((p1[0], p2[0] ,p3[0]), ("[0'-30']", "[30'-60']","[60'-90']"))

    #plt.tight_layout()
    plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_typology+'/'+'typology_reports_chart_section2'+graph_format, dpi = dpi_1)
    plt.show()

    page_number, pos_left, pos_top = matrix_positions_typology['typology_reports_chart_section2']['page_number'], matrix_positions_typology['typology_reports_chart_section2']['pos_left'],matrix_positions_typology['typology_reports_chart_section2']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_typology+'/'+'typology_reports_chart_section2'+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_6,prs))

    prs.save(pres_path)

    return BP_30,BP_60,BP_90,BC_30,BC_60,BC_90


def typology_reports_section2(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_typology,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3):
    
    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)

    ticks_size = 13
    label_size = ticks_size*1.15
    legend_size = ticks_size*1.15
    title_size = ticks_size*1.4
    
    a = 0
    BP_30,BP_60,BP_90 = [],[],[]
    BC_30,BC_60,BC_90 = [],[],[]
    for i in range (len(df_raw_away_1)):
        try:
            BP_30.append(df_raw_away_1.iloc[i,df_raw_away_1.columns.get_loc('BP_0-30')])
            a += 1
        except:
            pass
        try:
            BP_60.append(df_raw_away_1.iloc[i,df_raw_away_1.columns.get_loc('BP_30-60')])
        except:
            pass
        try:
            BP_90.append(df_raw_away_1.iloc[i,df_raw_away_1.columns.get_loc('BP_60-100')])
        except:
            pass     

        try:
            BC_30.append(df_raw_away_1.iloc[i,df_raw_away_1.columns.get_loc('BC_0-30')])
        except:
            pass
        try:
            BC_60.append(df_raw_away_1.iloc[i,df_raw_away_1.columns.get_loc('BC_30-60')])
        except:
            pass
        try:
            BC_90.append(df_raw_away_1.iloc[i,df_raw_away_1.columns.get_loc('BC_60-100')])
        except:
            pass      
    BP_30.reverse()
    BP_60.reverse()
    BP_90.reverse()
    BC_30.reverse()
    BC_60.reverse()
    BC_90.reverse() 

    BP_30=np.array(BP_30)
    BP_60=np.array(BP_60)
    BP_90=np.array(BP_90)
    
    BC_30=np.array(BC_30)
    BC_60=np.array(BC_60)
    BC_90=np.array(BC_90)
    
    width = 0.35       # the width of the bars: can also be len(x) sequence

    # BP for away team
    fig, ax = plt.subplots(1,1)

    ind = np.arange(a)
    print(a)
    plt.xticks(ind, ["N-{}".format(i) for i in range (a,0,-1)],fontsize=ticks_size-2)
    plt.yticks([0,1,2,3,4,5,6,7,8,9,10,11])

    every_nth = 2
    for n, label in enumerate(ax.xaxis.get_ticklabels()):
        if n % every_nth != 0:
            label.set_visible(False)
    if english == 1 :
        title = str("Timing of scored goals by ")+awayTeam_name
        ylabel = str("Goals scored")
        
    elif english == 0 :
        title = str("Timing des buts marqués par ")+awayTeam_name
        ylabel = str("Buts marqués")

    p1 = ax.bar(ind, BP_30, width)
    p2 = ax.bar(ind, BP_60,width , bottom=BP_30)
    p3 = ax.bar(ind, BP_90,width , bottom=BP_30+BP_60)

    plt.ylabel(ylabel)
    plt.title(title)
    
    plt.legend((p1[0], p2[0] ,p3[0]), ("[0'-30']", "[30'-60']","[60'-90']"))

    #plt.tight_layout()
    plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_typology+'/'+'typology_reports_chart_section3'+graph_format, dpi = dpi_1)
    plt.show()
    
    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)

    page_number, pos_left, pos_top = matrix_positions_typology['typology_reports_chart_section3']['page_number'], matrix_positions_typology['typology_reports_chart_section3']['pos_left'],matrix_positions_typology['typology_reports_chart_section3']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_typology+'/'+'typology_reports_chart_section3'+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_6,prs))

    # BC for away Team
    if english == 1 :
        title = str("Timing of against goals for ")+awayTeam_name
        ylabel = str("Goals against")
        
    elif english == 0 :
        title = str("Timing des buts concédés par ")+awayTeam_name
        ylabel = str("Buts concédés")

    fig, ax = plt.subplots(1,1)

    ind = np.arange(a)
    print(a)
    plt.xticks(ind, ["N-{}".format(i) for i in range (a,0,-1)],fontsize=ticks_size-2)
    plt.yticks([0,1,2,3,4,5,6,7,8,9,10,11])

    every_nth = 2
    for n, label in enumerate(ax.xaxis.get_ticklabels()):
        if n % every_nth != 0:
            label.set_visible(False)
    
    p1 = ax.bar(ind, BC_30, width)
    p2 = ax.bar(ind, BC_60,width , bottom=BC_30)
    p3 = ax.bar(ind, BC_90,width , bottom=BC_30+BC_60)

    plt.ylabel(ylabel)
    plt.title(title)
    
    plt.legend((p1[0], p2[0] ,p3[0]), ("[0'-30']", "[30'-60']","[60'-90']"))

    #plt.tight_layout()
    plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_typology+'/'+'typology_reports_chart_section4'+graph_format, dpi = dpi_1)
    plt.show()

    page_number, pos_left, pos_top = matrix_positions_typology['typology_reports_chart_section4']['page_number'], matrix_positions_typology['typology_reports_chart_section4']['pos_left'],matrix_positions_typology['typology_reports_chart_section4']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_typology+'/'+'typology_reports_chart_section4'+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_6,prs))

    prs.save(pres_path)

    return BP_30,BP_60,BP_90,BC_30,BC_60,BC_90


def typology_reports_section3(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_typology,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3):
    
    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)

    fouls_home = []
    yellow_home = []
    fouls_away = []
    yellow_away = []

    a_yellow_home = 0
    a_fouls_home = 0
    a_yellow_away = 0
    a_fouls_away = 0

    for i in range (len(df_raw_home_1)):
        try:
            yellow_home.append(df_raw_home_1.iloc[i,df_raw_home_1.columns.get_loc('yellow_team')])
            a_yellow_home += 1
        except: 
            pass
        try:
            fouls_home.append(df_raw_home_1.iloc[i,df_raw_home_1.columns.get_loc('fouls_team')])
            a_fouls_home += 1
        except:
            pass

    for i in range (len(df_raw_away_1)):
        try:
            yellow_away.append(df_raw_away_1.iloc[i,df_raw_away_1.columns.get_loc('yellow_team')])
            a_yellow_away += 1
        except: 
            pass
        try:
            fouls_away.append(df_raw_away_1.iloc[i,df_raw_away_1.columns.get_loc('fouls_team')])
            a_fouls_away += 1
        except:
            pass

    N_home = len(df_raw_home_1)
    N_away = len(df_raw_away_1)

    fouls_mean_home = [sum(fouls_home)/a_fouls_home]*N_home
    yellow_mean_home = [sum(yellow_home)/a_yellow_home]*N_home

    fouls_mean_away = [sum(fouls_away)/a_fouls_away]*N_away
    yellow_mean_away = [sum(yellow_away)/a_yellow_away]*N_away

    fouls_home.reverse()
    yellow_home.reverse()
    fouls_away.reverse()
    yellow_away.reverse()

    # HOME TEAM
    fig, ax = plt.subplots(1,1)
    
    ticks_size = 13
    label_size = ticks_size*1.15
    legend_size = ticks_size*1.15
    title_size = ticks_size*1.4

    if english == 1 :
        title = str("Fouls and yellow cards ")+homeTeam_name
        yellow_legends = str("Yellow cards")
        fouls_legends = str('Fouls')
        average = str('Average')
        
    elif english == 0 :
        title = str("Fautes et cartons jaunes ")+homeTeam_name
        yellow_legends = str("Cartons jaunes")
        fouls_legends = str('Fautes')
        average = str('Moyenne')

    ind = np.arange(N_home) 
    width = 0.30       
    ax.bar(ind, fouls_home, width, label=fouls_legends, color=N_home*[black_color_plot],edgecolor='black')
    ax.bar(ind + width, yellow_home, width,label=yellow_legends, color=N_home*[yellow_color_plot],edgecolor='black')

    ax.plot(ind,fouls_mean_home, color=black_color_plot, linestyle='--')
    ax.plot(ind,yellow_mean_home, color=yellow_color_plot, linestyle='--')
    
    ax.set_title(title,fontsize=title_size,weight='bold')
    plt.xticks(ind + width / 2, ["N-{}".format(i) for i in range (N_home,0,-1)],fontsize=ticks_size-2)
    plt.yticks(fontsize=ticks_size)
    patch = [mpatches.Patch(color=black_color_plot, label=fouls_legends), mpatches.Patch(color=yellow_color_plot, label=yellow_legends) , matplotlib.lines.Line2D(xdata=[1],ydata=[1],linestyle='--', color='black',label=average)]
    ax.legend(handles=patch,loc='lower center',frameon=False, ncol=len(patch), bbox_to_anchor = (0,0,1,1),bbox_transform = plt.gcf().transFigure,fontsize=legend_size)
    
    every_nth = 2
    for n, label in enumerate(ax.xaxis.get_ticklabels()):
        if n % every_nth != 0:
            label.set_visible(False)

    plt.tight_layout()
    fig.autofmt_xdate()

    plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_typology+'/'+'typology_reports_chart_section5'+graph_format, dpi = dpi_1)
    plt.show()

    page_number, pos_left, pos_top = matrix_positions_typology['typology_reports_chart_section5']['page_number'], matrix_positions_typology['typology_reports_chart_section5']['pos_left'],matrix_positions_typology['typology_reports_chart_section5']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_typology+'/'+'typology_reports_chart_section5'+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_6,prs))
    
    prs.save(pres_path)

    # AWAY TEAM
    fig, ax = plt.subplots(1,1)
    ticks_size = 13
    label_size = ticks_size*1.15
    legend_size = ticks_size*1.15
    title_size = ticks_size*1.4

    if english == 1 :
        title = str("Fouls and yellow cards ")+awayTeam_name
        yellow_legends = str("Yellow cards")
        fouls_legends = str('Fouls')
        average = str('Average')
        
    elif english == 0 :
        title = str("Fautes et cartons jaunes ")+awayTeam_name
        yellow_legends = str("Cartons jaunes")
        fouls_legends = str('Fautes')
        average = str('Moyenne')

    ind = np.arange(N_away) 
    width = 0.30       
    ax.bar(ind, fouls_away, width, label=fouls_legends, color=N_away*[black_color_plot],edgecolor='black')
    ax.bar(ind + width, yellow_away, width,label=yellow_legends, color=N_away*[yellow_color_plot],edgecolor='black')

    ax.plot(ind,fouls_mean_away, color=black_color_plot, linestyle='--')
    ax.plot(ind,yellow_mean_away, color=yellow_color_plot, linestyle='--')
    
    ax.set_title(title,fontsize=title_size,weight='bold')
    plt.xticks(ind + width / 2, ["N-{}".format(i) for i in range (N_away,0,-1)],fontsize=ticks_size-2)
    plt.yticks(fontsize=ticks_size)
    patch = [mpatches.Patch(color=black_color_plot, label=fouls_legends), mpatches.Patch(color=yellow_color_plot, label=yellow_legends) , matplotlib.lines.Line2D(xdata=[1],ydata=[1],linestyle='--', color='black',label=average)]
    ax.legend(handles=patch,loc='lower center',frameon=False, ncol=len(patch), bbox_to_anchor = (0,0,1,1),bbox_transform = plt.gcf().transFigure,fontsize=legend_size)
    
    every_nth = 2
    for n, label in enumerate(ax.xaxis.get_ticklabels()):
        if n % every_nth != 0:
            label.set_visible(False)

    plt.tight_layout()
    fig.autofmt_xdate()

    plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_typology+'/'+'typology_reports_chart_section6'+graph_format, dpi = dpi_1)
    plt.show()

    page_number, pos_left, pos_top = matrix_positions_typology['typology_reports_chart_section6']['page_number'], matrix_positions_typology['typology_reports_chart_section6']['pos_left'],matrix_positions_typology['typology_reports_chart_section6']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_typology+'/'+'typology_reports_chart_section6'+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_6,prs))
    
    prs.save(pres_path)

    return sum(fouls_home)/a_fouls_home, sum(yellow_home)/a_yellow_home, sum(fouls_away)/a_fouls_away, sum(yellow_away)/a_yellow_away


# For home team
def typology_reports_section4(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_typology,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3):

    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)
    
    ticks_size = 13
    label_size = ticks_size*1.15
    legend_size = ticks_size*1.15
    title_size = ticks_size*1.4
        
    color_data = ['silver','whitesmoke']
    # A. Language parameters
    if english == 1 :
        labels = ["'Abnormal' goals","'Normal' goals"] # Keep the same order
    elif english == 0 :
        labels = ["Buts 'anormaux'", "Buts 'normaux'"] # Keep the same order
    
    fig, ax = plt.subplots(1,2)
    
    BP_tot = 0
    BC_tot = 0
    BP_normal = 0
    BC_normal = 0
    BP_anormal = 0
    BC_anormal = 0

    # B. Plot Graph 1
    for i in range (len(df_raw_home_1)):
        try:
            BP_tot += df_raw_home_1.iloc[i,df_raw_home_1.columns.get_loc('BP')]
            BC_tot += df_raw_home_1.iloc[i,df_raw_home_1.columns.get_loc('BC')]
            
            BP_normal += df_raw_home_1.iloc[i,df_raw_home_1.columns.get_loc('BP_normal')]
            BC_normal += df_raw_home_1.iloc[i,df_raw_home_1.columns.get_loc('BC_normal')]

        except: 
            pass
    
    BP_anormal = BP_tot - BP_normal
    BC_anormal = BC_tot - BC_normal
    values_BP = [BP_anormal,BP_normal]
    values_BC = [BC_anormal,BC_normal]

    patches, texts, autotexts = ax[0].pie(values_BP, autopct=make_autopct(values_BP),shadow=False, labels = labels, colors = color_data, wedgeprops = {'linewidth': 6,'edgecolor' : 'white'}, rotatelabels=False,pctdistance=0.6,labeldistance=None, textprops = {'fontsize': label_size,'color':"black",'weight':"bold"})
    patch = [mpatches.Patch(color=color_data[0], label=labels[0]), mpatches.Patch(color=color_data[1], label=labels[1])]
    ax[0].legend(handles=patch,loc='lower center',frameon=False, ncol=len(patch), bbox_to_anchor = (0,0,1,1),bbox_transform = plt.gcf().transFigure,fontsize=legend_size)

    center, r = patches[0].center, patches[0].r
    centre_circle = plt.Circle(center,0.25,fc='white')
    edge_circle = plt.Circle(center,r,fill=False,ec=black_color_plot)
    ax[0].add_artist(centre_circle)
    ax[0].add_artist(edge_circle)

    ax[0].axis('equal')
    ax[0].set_title("Buts pour",weight='bold',fontsize=title_size,color= black_color_plot)
    
    # C. Plot Graph 2
    patches, texts, autotexts = ax[1].pie(values_BC, autopct=make_autopct(values_BC),shadow=False, labels = labels, colors = color_data, wedgeprops = {'linewidth': 6,'edgecolor' : 'white'}, rotatelabels=False,pctdistance=0.6,labeldistance=None, textprops = {'fontsize': label_size,'color':"black",'weight':"bold"})
    patch = [mpatches.Patch(color=color_data[0], label=labels[0]), mpatches.Patch(color=color_data[1], label=labels[1])]
    ax[1].legend(handles=patch,loc='lower center',frameon=False, ncol=len(patch), bbox_to_anchor = (0,0,1,1),bbox_transform = plt.gcf().transFigure,fontsize=legend_size)

    center, r = patches[1].center, patches[1].r
    centre_circle = plt.Circle(center,0.25,fc='white')
    edge_circle = plt.Circle(center,r,fill=False,ec=black_color_plot)
    ax[1].add_artist(centre_circle)
    ax[1].add_artist(edge_circle)

    ax[1].axis('equal')
    ax[1].set_title("Buts contre",weight='bold',fontsize=title_size,color= black_color_plot)
    
    fig.suptitle("Analyse des buts " + awayTeam_name, fontsize=16)
    fig.subplots_adjust(top=0.8)

    plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_typology+'/'+'typology_reports_chart_section7'+graph_format, dpi = dpi_1)
    plt.show()
    
    page_number, pos_left, pos_top = matrix_positions_typology['typology_reports_chart_section7']['page_number'], matrix_positions_typology['typology_reports_chart_section7']['pos_left'],matrix_positions_typology['typology_reports_chart_section7']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_typology+'/'+'typology_reports_chart_section7'+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_1,prs))
    
    prs.save(pres_path)

    return values_BP,  values_BC

# For away team
def typology_reports_section5(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_typology,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3):
    
    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)

    ticks_size = 13
    label_size = ticks_size*1.15
    legend_size = ticks_size*1.15
    title_size = ticks_size*1.4
    
    color_data = ['silver','whitesmoke']
    # A. Language parameters
    if english == 1 :
        labels = ["'Abnormal' goals","'Normal' goals"] # Keep the same order
    elif english == 0 :
        labels = ["Buts 'anormaux'", "Buts 'normaux'"] # Keep the same order
    
    fig, ax = plt.subplots(1,2)
    
    BP_tot = 0
    BC_tot = 0
    BP_normal = 0
    BC_normal = 0
    BP_anormal = 0
    BC_anormal = 0

    # B. Plot Graph 1
    for i in range (len(df_raw_away_1)):
        try:
            BP_tot += df_raw_away_1.iloc[i,df_raw_away_1.columns.get_loc('BP')]
            BC_tot += df_raw_away_1.iloc[i,df_raw_away_1.columns.get_loc('BC')]
            
            BP_normal += df_raw_away_1.iloc[i,df_raw_away_1.columns.get_loc('BP_normal')]
            BC_normal += df_raw_away_1.iloc[i,df_raw_away_1.columns.get_loc('BC_normal')]

        except: 
            pass
    
    BP_anormal = BP_tot - BP_normal
    BC_anormal = BC_tot - BC_normal
    values_BP = [BP_anormal,BP_normal]
    values_BC = [BC_anormal,BC_normal]

    patches, texts, autotexts = ax[0].pie(values_BP, autopct=make_autopct(values_BP),shadow=False, labels = labels, colors = color_data, wedgeprops = {'linewidth': 6,'edgecolor' : 'white'}, rotatelabels=False,pctdistance=0.6,labeldistance=None, textprops = {'fontsize': label_size,'color':"black",'weight':"bold"})
    patch = [mpatches.Patch(color=color_data[0], label=labels[0]), mpatches.Patch(color=color_data[1], label=labels[1])]
    ax[0].legend(handles=patch,loc='lower center',frameon=False, ncol=len(patch), bbox_to_anchor = (0,0,1,1),bbox_transform = plt.gcf().transFigure,fontsize=legend_size)

    center, r = patches[0].center, patches[0].r
    centre_circle = plt.Circle(center,0.25,fc='white')
    edge_circle = plt.Circle(center,r,fill=False,ec=black_color_plot)
    ax[0].add_artist(centre_circle)
    ax[0].add_artist(edge_circle)

    ax[0].axis('equal')
    ax[0].set_title("Buts pour",weight='bold',fontsize=title_size,color= black_color_plot)
    
    # C. Plot Graph 2
    patches, texts, autotexts = ax[1].pie(values_BC, autopct=make_autopct(values_BC),shadow=False, labels = labels, colors = color_data, wedgeprops = {'linewidth': 6,'edgecolor' : 'white'}, rotatelabels=False,pctdistance=0.6,labeldistance=None, textprops = {'fontsize': label_size,'color':"black",'weight':"bold"})
    patch = [mpatches.Patch(color=color_data[0], label=labels[0]), mpatches.Patch(color=color_data[1], label=labels[1])]
    ax[1].legend(handles=patch,loc='lower center',frameon=False, ncol=len(patch), bbox_to_anchor = (0,0,1,1),bbox_transform = plt.gcf().transFigure,fontsize=legend_size)

    center, r = patches[1].center, patches[1].r
    centre_circle = plt.Circle(center,0.25,fc='white')
    edge_circle = plt.Circle(center,r,fill=False,ec=black_color_plot)
    ax[1].add_artist(centre_circle)
    ax[1].add_artist(edge_circle)

    ax[1].axis('equal')
    ax[1].set_title("Buts contre",weight='bold',fontsize=title_size,color= black_color_plot)

    
    fig.suptitle("Analyse des buts " + awayTeam_name, fontsize=16)
    fig.subplots_adjust(top=0.8)

    plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_typology+'/'+'typology_reports_chart_section8'+graph_format, dpi = dpi_1)
    plt.show()

    page_number, pos_left, pos_top = matrix_positions_typology['typology_reports_chart_section8']['page_number'], matrix_positions_typology['typology_reports_chart_section8']['pos_left'],matrix_positions_typology['typology_reports_chart_section8']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_typology+'/'+'typology_reports_chart_section8'+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_1,prs))
    
    prs.save(pres_path)

    return values_BP,  values_BC


# Team possession
def typology_reports_section6(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_typology,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3):
    
    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)
    
    if english == 1 :
        average = "Average"
    elif english == 0 :
        average = "Moyenne"
    
    fig, ax = plt.subplots(1,1)

    ticks_size = 13
    label_size = ticks_size*1.15
    legend_size = ticks_size*1.15
    title_size = ticks_size*1.4

    Poss_home = []
    Poss_home_n = 0
    Poss_away = []
    Poss_away_n = 0
    
    N_home = len(df_raw_home_1)
    N_away = len(df_raw_away_1)
    
    # B. Plot Graph 1
    for i in range (len(df_raw_away_1)):
        try:
            Poss_away.append(df_raw_away_1.iloc[i,df_raw_away_1.columns.get_loc('Poss team')])
            Poss_away_n += 1
        except: 
            pass
    
    for i in range (len(df_raw_home_1)):
        try:
            Poss_home.append(df_raw_home_1.iloc[i,df_raw_home_1.columns.get_loc('Poss team')])
            Poss_home_n += 1
        except: 
            pass
    
    Poss_home.reverse()
    Poss_away.reverse()

    Poss_home_mean = [sum(Poss_home)/Poss_home_n]*len(df_raw_away_1)
    Poss_away_mean = [sum(Poss_away)/Poss_away_n]*len(df_raw_home_1)

    x_home = ["N-{}".format(i) for i in range (N_home,0,-1)]
    x_away = ["N-{}".format(i) for i in range (N_away,0,-1)]
    
    # First plot, for HOME TEAM
    ax.plot(x_home, Poss_home,color=home_color_plot)
    ax.plot(x_home,Poss_home_mean, color=black_color_plot, linestyle='--')

    ax.set_title(homeTeam_name,weight='bold',fontsize=title_size,color=home_color_plot)
    
    every_nth = 3
    for n, label in enumerate(ax.xaxis.get_ticklabels()):
        if n % every_nth != 0:
            label.set_visible(False) 
    
    patch = [matplotlib.lines.Line2D(xdata=[1],ydata=[1],linestyle='--', color='black',label=average)]
    ax.legend(handles=patch,loc='upper right',fontsize=legend_size)

    plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_typology+'/'+'typology_reports_chart_section9'+graph_format, dpi = dpi_1)
    plt.show()

    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)

    page_number, pos_left, pos_top = matrix_positions_typology['typology_reports_chart_section9']['page_number'], matrix_positions_typology['typology_reports_chart_section9']['pos_left'],matrix_positions_typology['typology_reports_chart_section9']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_typology+'/'+'typology_reports_chart_section9'+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_6,prs))

    prs.save(pres_path)

    # Second plot, for AWAY TEAM
    fig, ax = plt.subplots(1,1)
    ax.plot(x_away, Poss_away,color=away_color_plot)
    ax.plot(x_away,Poss_away_mean, color=black_color_plot, linestyle='--')

    ax.set_title(awayTeam_name,weight='bold',fontsize=title_size,color=away_color_plot)
    
    every_nth = 3
    for n, label in enumerate(ax.xaxis.get_ticklabels()):
        if n % every_nth != 0:
            label.set_visible(False) 
    
    patch = [matplotlib.lines.Line2D(xdata=[1],ydata=[1],linestyle='--', color='black',label=average)]
    ax.legend(handles=patch,loc='upper right',fontsize=legend_size)

    plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_typology+'/'+'typology_reports_chart_section10'+graph_format, dpi = dpi_1)
    plt.show()

    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)

    page_number, pos_left, pos_top = matrix_positions_typology['typology_reports_chart_section10']['page_number'], matrix_positions_typology['typology_reports_chart_section10']['pos_left'],matrix_positions_typology['typology_reports_chart_section10']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_typology+'/'+'typology_reports_chart_section10'+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_6,prs))

    prs.save(pres_path)

    return sum(Poss_home)/Poss_home_n , sum(Poss_away)/Poss_away_n


# Team possession
def typology_reports_section7(team_id,team_name,opp_id,opp_name,game_fixture_id,name_folder_typology,homeTeam_id,awayTeam_id,homeTeam_name,awayTeam_name,df_raw_home_1,df_raw_home_2,df_raw_home_3,df_raw_away_1,df_raw_away_2,df_raw_away_3):
    
    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)
    
    if english == 1 :
        average = "Average"
        title = "Direct goals"
    elif english == 0 :
        average = "Moyenne"
        title = "Buts directs"
    
    fig, ax = plt.subplots(1,1)

    ticks_size = 13
    label_size = ticks_size*1.15
    legend_size = ticks_size*1.15
    title_size = ticks_size*1.4

    BP_direct_home = []
    BC_direct_home = []
    home_n = 0
    
    BP_direct_away = []
    BC_direct_away = []
    away_n = 0
    
    N_home = len(df_raw_home_1)
    N_away = len(df_raw_away_1)
    
    # B. Plot Graph 1
    for i in range (len(df_raw_away_1)):
        try:
            BP_direct_away.append(df_raw_away_1.iloc[i,df_raw_away_1.columns.get_loc('BP_direct')])
            BC_direct_away.append(df_raw_away_1.iloc[i,df_raw_away_1.columns.get_loc('BC_direct')])
            away_n += 1
        except: 
            pass
    
    for i in range (len(df_raw_home_1)):
        try:
            BP_direct_home.append(df_raw_home_1.iloc[i,df_raw_home_1.columns.get_loc('BP_direct')])
            BC_direct_home.append(df_raw_home_1.iloc[i,df_raw_home_1.columns.get_loc('BC_direct')])
            home_n += 1
        except: 
            pass
    
    BP_direct_away.reverse()
    BC_direct_away.reverse()
    BP_direct_home.reverse()
    BC_direct_home.reverse()

    BP_direct_away_mean = [sum(BP_direct_away)/away_n]*len(df_raw_away_1)
    BC_direct_away_mean = [sum(BC_direct_away)/away_n]*len(df_raw_away_1)

    BP_direct_home_mean = [sum(BP_direct_home)/home_n]*len(df_raw_home_1)
    BC_direct_home_mean = [sum(BC_direct_home)/home_n]*len(df_raw_home_1)

    x_home = ["N-{}".format(i) for i in range (N_home,0,-1)]
    x_away = ["N-{}".format(i) for i in range (N_away,0,-1)]
    
    # First plot, for HOME TEAM
    ind = np.arange(N_home) 

    width = 0.30       
    ax.bar(ind, BP_direct_home, width, label="BP directs " + homeTeam_name, color=N_home*[home_color_plot],edgecolor='black')
    ax.bar(ind + width, BC_direct_home, width,label="BC directs " + homeTeam_name, color=N_home*[away_color_plot],edgecolor='black')

    ax.plot(ind,BP_direct_home_mean, color=home_color_plot, linestyle='--')
    ax.plot(ind,BC_direct_home_mean, color=away_color_plot, linestyle='--')

    ax.set_title(title,fontsize=title_size,weight='bold')
    plt.xticks(ind + width / 2, x_home,fontsize=ticks_size-2)
    plt.yticks(fontsize=ticks_size)
    patch = [mpatches.Patch(color=home_color_plot, label="BP directs " + homeTeam_name), mpatches.Patch(color=away_color_plot, label="BC directs " + homeTeam_name) , matplotlib.lines.Line2D(xdata=[1],ydata=[1],linestyle='--', color='black',label=average)]
    ax.legend(handles=patch,loc='lower center',frameon=False, ncol=len(patch), bbox_to_anchor = (0,0,1,1),bbox_transform = plt.gcf().transFigure,fontsize=legend_size)
    
    every_nth = 3
    for n, label in enumerate(ax.xaxis.get_ticklabels()):
        if n % every_nth != 0:
            label.set_visible(False)

    plt.tight_layout()
    fig.autofmt_xdate()

    plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_typology+'/'+'typology_reports_chart_section11'+graph_format, dpi = dpi_1)
    plt.show()

    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)

    page_number, pos_left, pos_top = matrix_positions_typology['typology_reports_chart_section11']['page_number'], matrix_positions_typology['typology_reports_chart_section11']['pos_left'],matrix_positions_typology['typology_reports_chart_section11']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_typology+'/'+'typology_reports_chart_section11'+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_6,prs))

    prs.save(pres_path)

    # Second plot, for AWAY TEAM
    fig, ax = plt.subplots(1,1)
    ind = np.arange(N_away) 

    width = 0.30       
    ax.bar(ind, BP_direct_away, width, label="BP directs " + awayTeam_name, color=N_away*[away_color_plot],edgecolor='black')
    ax.bar(ind + width, BC_direct_away, width,label="BC directs " + awayTeam_name, color=N_away*[away_color_plot],edgecolor='black')

    ax.plot(ind,BP_direct_away_mean, color=away_color_plot, linestyle='--')
    ax.plot(ind,BC_direct_away_mean, color=away_color_plot, linestyle='--')

    ax.set_title(title,fontsize=title_size,weight='bold')
    plt.xticks(ind + width / 2, x_away,fontsize=ticks_size-2)
    plt.yticks(fontsize=ticks_size)
    patch = [mpatches.Patch(color=away_color_plot, label="BP directs " + awayTeam_name), mpatches.Patch(color=away_color_plot, label="BC directs " + awayTeam_name) , matplotlib.lines.Line2D(xdata=[1],ydata=[1],linestyle='--', color='black',label=average)]
    ax.legend(handles=patch,loc='lower center',frameon=False, ncol=len(patch), bbox_to_anchor = (0,0,1,1),bbox_transform = plt.gcf().transFigure,fontsize=legend_size)
    
    every_nth = 3
    for n, label in enumerate(ax.xaxis.get_ticklabels()):
        if n % every_nth != 0:
            label.set_visible(False)

    plt.tight_layout()
    fig.autofmt_xdate()

    plt.savefig(fname=link_reports +str(game_fixture_id) +'/'+name_folder_typology+'/'+'typology_reports_chart_section12'+graph_format, dpi = dpi_1)
    plt.show()

    pres_path = link_reports + game_fixture_id + '/' + 'report' +'/'+'report' + '.pptx' 
    prs = Presentation(pres_path)

    page_number, pos_left, pos_top = matrix_positions_typology['typology_reports_chart_section12']['page_number'], matrix_positions_typology['typology_reports_chart_section12']['pos_left'],matrix_positions_typology['typology_reports_chart_section12']['pos_top'] 
    slide = prs.slides[page_number]
    pic = slide.shapes.add_picture(link_reports +str(game_fixture_id) +'/'+name_folder_typology+'/'+'typology_reports_chart_section12'+graph_format,left=cm_to_ppx(pos_left+buffer,prs), top=cm_to_ppx(pos_top+buffer,prs), height= cm_to_ppx(height_6,prs))

    prs.save(pres_path)

    return BP_direct_home , BP_direct_away